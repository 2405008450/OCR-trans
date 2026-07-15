import io
import json
import sys
from pathlib import Path
from types import SimpleNamespace
from zipfile import ZIP_DEFLATED, ZipFile

import anyio
import pytest
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from sqlalchemy import create_engine, text
from sqlalchemy.orm import sessionmaker
from starlette.datastructures import UploadFile

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app.controller import task as task_controller
from app.db.database import Base
from app.model.entity import Task
from app.service import task_queue_service as queue_module
from app.service import word_count_service
from app.service.ocr_text_service import ocr_markup_to_plain_text
import pdf2docx as pdf2docx_module


def _make_png(path: Path) -> None:
    from PIL import Image

    image = Image.new("RGB", (24, 24), color=(48, 128, 220))
    image.save(path)


def _set_docx_app_props(path: Path, updates: dict[str, str]) -> None:
    temp_path = path.with_suffix(".tmp.docx")
    with ZipFile(path, "r") as source, ZipFile(temp_path, "w", ZIP_DEFLATED) as target:
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename == "docProps/app.xml":
                from lxml import etree

                root = etree.fromstring(data)
                for child in root:
                    tag = etree.QName(child).localname
                    if tag in updates:
                        child.text = str(updates[tag])
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            target.writestr(item, data)
    path.unlink()
    temp_path.rename(path)


def _allow_root(monkeypatch, root: Path):
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps([str(root)]))
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_UNC_MOUNT_MAP_JSON", "")
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")
    monkeypatch.setattr(queue_module.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps([str(root)]))
    monkeypatch.setattr(queue_module.settings, "WORD_COUNT_UNC_MOUNT_MAP_JSON", "")
    monkeypatch.setattr(queue_module.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")


def test_discover_and_prepare_selected_word_count_files(monkeypatch, tmp_path):
    root = tmp_path / "共享目录"
    root.mkdir()
    (root / "a.txt").write_text("甲乙", encoding="utf-8")
    (root / "b.txt").write_text("hello world", encoding="utf-8")
    (root / "ignore.exe").write_bytes(b"MZ")
    _allow_root(monkeypatch, root)

    discovered = word_count_service.discover_word_count_files(
        directory_path=str(root),
        recursive=True,
        extensions=[".txt"],
    )
    assert [item["relative_path"] for item in discovered["files"]] == ["a.txt", "b.txt"]

    prepared = word_count_service.prepare_word_count_request(
        directory_path=str(root),
        extensions=[".txt"],
        relative_paths=["b.txt"],
    )
    assert prepared["input_files"]["relative_paths"] == ["b.txt"]
    assert prepared["params"]["relative_paths"] == ["b.txt"]


def test_word_count_task_only_processes_selected_files(monkeypatch, tmp_path):
    root = tmp_path / "共享目录"
    root.mkdir()
    (root / "a.txt").write_text("甲乙", encoding="utf-8")
    (root / "b.txt").write_text("hello world", encoding="utf-8")
    _allow_root(monkeypatch, root)
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))

    result = word_count_service.run_word_count_task_sync(
        task_id="selected-task",
        display_no="WC-SELECTED",
        directory_path=str(root),
        recursive=True,
        include_hidden=False,
        extensions=[".txt"],
        relative_paths=["b.txt"],
    )

    assert [item["relative_path"] for item in result["files"]] == ["b.txt"]
    assert result["summary"]["counted_files"] == 1


def test_selected_word_count_file_respects_recursive_option(monkeypatch, tmp_path):
    root = tmp_path / "共享目录"
    nested = root / "子目录"
    nested.mkdir(parents=True)
    (nested / "a.txt").write_text("甲乙", encoding="utf-8")
    _allow_root(monkeypatch, root)

    with pytest.raises(ValueError, match="未启用子目录扫描"):
        word_count_service.prepare_word_count_request(
            directory_path=str(root),
            recursive=False,
            extensions=[".txt"],
            relative_paths=["子目录/a.txt"],
        )


def test_word_like_counter_counts_mixed_languages():
    text = "你好 world 123，テスト 한국어"

    assert word_count_service.count_words_word_like(text) == 11


def test_word_like_counter_exposes_script_buckets_for_quote():
    text = (
        "广药白云山 爱心满人间\n"
        "GUANGZHOU PHAR.'S BAIYUNSHAN\n"
        "GIVES FULL HEARTINESS TO THE WORLD\t广药集团 世界500强\n"
        "\t广州医药集团有限公司\n"
        "GUANGZHOU PHARMACEUTICAL HOLDINGS LIMITED\n"
    )

    metrics = word_count_service._count_text(text)

    assert metrics.word_count == 41
    assert metrics.script_count_total == metrics.word_count
    assert metrics.han_count == 27
    assert metrics.latin_word_count == 13
    assert metrics.number_token_count == 1
    assert metrics.billable_chinese_count == 27
    assert metrics.billable_latin_count == 13


def test_word_like_counter_splits_common_language_pairs():
    japanese = word_count_service._count_text("日本語テスト English 42")
    korean = word_count_service._count_text("한국어 테스트 English")
    russian_arabic = word_count_service._count_text("Привет world مرحبا 123")

    assert japanese.word_count == 8
    assert japanese.han_count == 3
    assert japanese.kana_count == 3
    assert japanese.billable_japanese_count == 6
    assert japanese.latin_word_count == 1
    assert japanese.number_token_count == 1

    assert korean.word_count == 7
    assert korean.hangul_count == 6
    assert korean.billable_korean_count == 6
    assert korean.latin_word_count == 1

    assert russian_arabic.word_count == 4
    assert russian_arabic.cyrillic_word_count == 1
    assert russian_arabic.latin_word_count == 1
    assert russian_arabic.arabic_word_count == 1
    assert russian_arabic.number_token_count == 1


def test_word_like_counter_matches_word_stats_sample():
    sample = Path("data/word/2. Blog Post 42 - ENG - Red Notice removal from INTERPOL, and when NOT to apply - word_translated.docx")
    if not sample.exists():
        pytest.skip("sample DOCX is not available")

    items = word_count_service._extract_docx_text_items(sample)
    main_text = "\n".join(item.text for item in items if not item.is_extra)
    metrics = word_count_service._count_text(main_text)

    assert metrics.word_count == 693
    assert metrics.non_space_chars == 724


def test_word_like_counter_handles_edge_punctuation():
    assert word_count_service.count_words_word_like("“中文”") == 4
    assert word_count_service.count_words_word_like("中文…") == 3


def test_ocr_markup_to_plain_text_removes_html_and_markdown():
    raw = """# 标题

**Hello** world
<table><tr><td>表格</td><td>123</td></tr></table>
"""

    plain = ocr_markup_to_plain_text(raw)
    metrics = word_count_service._count_text(plain)

    assert "<table>" not in plain
    assert "**" not in plain
    assert "标题" in plain
    assert "Hello" in plain
    assert metrics.word_count == 7
    assert metrics.script_count_total == metrics.word_count


def test_docx_nested_table_text_is_not_counted_twice(tmp_path):
    doc = Document()
    outer = doc.add_table(rows=1, cols=1)
    outer.cell(0, 0).text = "外层"
    nested = outer.cell(0, 0).add_table(rows=1, cols=1)
    nested.cell(0, 0).text = "内层"
    path = tmp_path / "nested.docx"
    doc.save(path)

    items = word_count_service._extract_docx_text_items(path)
    text = "\n".join(item.text for item in items if not item.is_extra)
    metrics = word_count_service._count_text(text)

    assert metrics.word_count == 4
    assert metrics.han_count == 4


def test_docx_main_and_extra_counts_are_separated(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))

    image_path = tmp_path / "doc-image.png"
    _make_png(image_path)
    doc = Document()
    doc.add_paragraph("你好 world")
    table = doc.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "金额 123"
    section = doc.sections[0]
    section.header.paragraphs[0].text = "页眉 text"
    doc.add_picture(str(image_path))
    doc.save(tmp_path / "sample.docx")
    _set_docx_app_props(tmp_path / "sample.docx", {"Pages": "2", "Lines": "9"})

    result = word_count_service.run_word_count_task_sync(
        task_id="task-1",
        display_no="000001",
        directory_path=str(tmp_path),
        recursive=True,
        include_hidden=False,
        extensions=[".docx"],
    )

    summary = result["summary"]
    assert summary["counted_files"] == 1
    assert summary["total_main_word_count"] == 6
    assert summary["total_extra_word_count"] == 3
    assert summary["script_count_total"] == 6
    assert summary["total_han_count"] == 4
    assert summary["total_latin_word_count"] == 1
    assert summary["total_number_token_count"] == 1
    assert summary["total_billable_chinese_count"] == 4
    assert summary["total_billable_latin_count"] == 1
    assert summary["total_extra_han_count"] == 2
    assert summary["total_page_count"] == 2
    assert summary["total_line_count"] == 9
    assert summary["total_image_count"] == 1
    file_result = result["files"][0]
    assert file_result["word_count"] == 6
    assert file_result["script_count_total"] == 6
    assert file_result["script_counts"]["han_count"] == 4
    assert file_result["quote_counts"]["billable_chinese_count"] == 4
    assert file_result["han_count"] == 4
    assert file_result["latin_word_count"] == 1
    assert file_result["number_token_count"] == 1
    assert file_result["extra_han_count"] == 2
    assert file_result["char_count_no_spaces"] > 0
    assert file_result["char_count_with_spaces"] >= file_result["char_count_no_spaces"]
    assert file_result["paragraph_count"] >= 3
    assert file_result["image_count"] == 1
    assert file_result["file_type"] == "Word"
    assert "DOCX" in file_result["stat_method"]
    assert file_result["counted_at"]
    assert Path(tmp_path / result["report_excel"]).exists()
    assert Path(tmp_path / result["report_json"]).exists()

    report = load_workbook(tmp_path / result["report_excel"], read_only=True)
    try:
        headers = [cell.value for cell in next(report["文件明细"].iter_rows(max_row=1))]
        assert "页数" in headers
        assert "字符数(不计空格)" in headers
        assert "图片数量" in headers
        assert "统计方法" in headers
        assert "中文候选" in headers
        assert "汉字" in headers
    finally:
        report.close()


def test_word_count_accepts_single_file_path(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))

    doc = Document()
    doc.add_paragraph("单文件 world 123")
    file_path = tmp_path / "single.docx"
    doc.save(file_path)

    prepared = word_count_service.prepare_word_count_request(directory_path=str(file_path))
    assert prepared["params"]["input_kind"] == "file"
    assert Path(prepared["params"]["directory_path"]) == file_path.resolve(strict=False)

    result = word_count_service.run_word_count_task_sync(
        task_id="task-single",
        display_no="000001-single",
        directory_path=str(file_path),
        recursive=True,
        include_hidden=False,
        extensions=[".docx"],
    )

    assert result["input_kind"] == "file"
    assert result["summary"]["counted_files"] == 1
    assert result["summary"]["total_main_word_count"] == 5
    assert result["files"][0]["relative_path"] == "single.docx"


def test_office_pdf_txt_and_blank_pdf_statuses(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))

    image_path = tmp_path / "embedded.png"
    _make_png(image_path)
    txt_path = tmp_path / "plain.txt"
    txt_path.write_text("中文 text 42\nsecond line", encoding="utf-8")

    workbook = Workbook()
    workbook.active["A1"] = "表格 abc"
    workbook.active["A2"] = "第二行"
    from openpyxl.drawing.image import Image as XlsxImage

    workbook.active.add_image(XlsxImage(str(image_path)), "B2")
    workbook.create_sheet("空表")
    workbook.save(tmp_path / "sheet.xlsx")

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "标题 one"
    tx_box = slide.shapes.add_textbox(0, 0, 914400, 914400)
    tx_box.text_frame.text = "文本框 two"
    table = slide.shapes.add_table(1, 1, 0, 914400, 914400, 914400).table
    table.cell(0, 0).text = "表格 three"
    slide.shapes.add_picture(str(image_path), 0, 1828800)
    slide.notes_slide.notes_text_frame.text = "备注 extra"
    presentation.save(tmp_path / "deck.pptx")

    import fitz

    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "PDF text")
    pdf.save(tmp_path / "text.pdf")
    pdf.close()

    blank_pdf = fitz.open()
    blank_pdf.new_page()
    blank_pdf.save(tmp_path / "blank.pdf")
    blank_pdf.close()

    result = word_count_service.run_word_count_task_sync(
        task_id="task-2",
        display_no="000002",
        directory_path=str(tmp_path),
        recursive=True,
        include_hidden=False,
        extensions=[".txt", ".xlsx", ".pptx", ".pdf"],
    )

    statuses = {item["filename"]: item["status"] for item in result["files"]}
    assert statuses["plain.txt"] == "counted"
    assert statuses["sheet.xlsx"] == "counted"
    assert statuses["deck.pptx"] == "counted"
    assert statuses["text.pdf"] == "counted"
    assert statuses["blank.pdf"] == "needs_ocr"
    assert result["summary"]["needs_ocr_files"] == 1
    files = {item["filename"]: item for item in result["files"]}
    assert files["plain.txt"]["page_count"] == 1
    assert files["plain.txt"]["line_count"] == 2
    assert files["sheet.xlsx"]["page_count"] == 2
    assert files["sheet.xlsx"]["line_count"] == 2
    assert files["sheet.xlsx"]["image_count"] == 1
    assert files["sheet.xlsx"]["file_type"] == "Excel"
    assert files["deck.pptx"]["page_count"] == 1
    assert files["deck.pptx"]["image_count"] == 1
    assert files["deck.pptx"]["extra_word_count"] > 0
    assert files["deck.pptx"]["file_type"] == "PPT"
    assert files["text.pdf"]["page_count"] == 1
    assert files["blank.pdf"]["page_count"] == 1
    assert result["summary"]["total_image_count"] >= 2


def test_ocr_auto_mode_depends_on_input_kind(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    image_path = tmp_path / "scan.png"
    _make_png(image_path)

    single = word_count_service.prepare_word_count_request(directory_path=str(image_path))
    directory = word_count_service.prepare_word_count_request(directory_path=str(tmp_path))
    directory_on = word_count_service.prepare_word_count_request(
        directory_path=str(tmp_path),
        ocr_mode="on",
    )

    assert single["params"]["ocr_enabled"] is True
    assert directory["params"]["ocr_enabled"] is False
    assert directory_on["params"]["ocr_enabled"] is True


def test_single_image_ocr_is_counted_and_archived(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))
    image_path = tmp_path / "scan.png"
    _make_png(image_path)

    def fake_extract_ocr_plain_text(**_kwargs):
        return {
            "text": "中文 Hello 123",
            "total_pages": 1,
            "processed_pages": [1],
            "failed_pages": [],
            "page_results": [
                {"page_number": 1, "text": "中文 Hello 123", "raw_text": "", "blank": False, "error": ""}
            ],
        }

    monkeypatch.setattr(word_count_service, "extract_ocr_plain_text", fake_extract_ocr_plain_text)
    result = word_count_service.run_word_count_task_sync(
        task_id="task-image-ocr",
        display_no="000-image-ocr",
        directory_path=str(image_path),
        recursive=True,
        include_hidden=False,
        extensions=[".png"],
    )

    file_result = result["files"][0]
    assert file_result["status"] == "counted"
    assert file_result["ocr_used"] is True
    assert file_result["ocr_page_count"] == 1
    assert file_result["main_word_count"] == 4
    assert result["summary"]["ocr_files"] == 1
    assert result["summary"]["ocr_pages"] == 1
    assert result["summary"]["total_main_word_count"] == 4

    archive_path = tmp_path / result["ocr_text_archive"]
    assert archive_path.exists()
    with ZipFile(archive_path) as archive:
        assert archive.namelist() == ["scan.png.txt"]
        assert archive.read("scan.png.txt").decode("utf-8") == "中文 Hello 123"

    report = load_workbook(tmp_path / result["report_excel"], read_only=True)
    try:
        headers = [cell.value for cell in next(report["文件明细"].iter_rows(max_row=1))]
        assert "是否使用 OCR" in headers
        assert "OCR 失败页" in headers
        assert "OCR 文本路径" in headers
    finally:
        report.close()


def test_directory_auto_mode_keeps_image_as_needs_ocr(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))
    _make_png(tmp_path / "scan.png")

    def unexpected_ocr(**_kwargs):
        raise AssertionError("目录自动模式不应调用 OCR")

    monkeypatch.setattr(word_count_service, "extract_ocr_plain_text", unexpected_ocr)
    result = word_count_service.run_word_count_task_sync(
        task_id="task-dir-auto",
        display_no="000-dir-auto",
        directory_path=str(tmp_path),
        recursive=True,
        include_hidden=False,
        extensions=[".png"],
        ocr_mode="auto",
    )

    assert result["files"][0]["status"] == "needs_ocr"
    assert result["summary"]["needs_ocr_files"] == 1
    assert result["ocr_text_archive"] == ""


def test_mixed_pdf_only_ocrs_scanned_page(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))
    image_path = tmp_path / "page.png"
    from PIL import Image, ImageDraw

    image = Image.new("RGB", (900, 1200), "white")
    ImageDraw.Draw(image).text((80, 120), "scanned page", fill="black")
    image.save(image_path)

    import fitz

    pdf_path = tmp_path / "mixed.pdf"
    pdf = fitz.open()
    text_page = pdf.new_page()
    text_page.insert_text((72, 72), "PDF text layer content")
    scan_page = pdf.new_page(width=595, height=842)
    scan_page.insert_image(scan_page.rect, filename=str(image_path))
    pdf.save(pdf_path)
    pdf.close()

    captured = {}

    def fake_extract_ocr_plain_text(**kwargs):
        captured.update(kwargs)
        return {
            "text": "扫描页 Scan",
            "total_pages": 2,
            "processed_pages": [2],
            "failed_pages": [],
            "page_results": [
                {"page_number": 2, "text": "扫描页 Scan", "raw_text": "", "blank": False, "error": ""}
            ],
        }

    monkeypatch.setattr(word_count_service, "extract_ocr_plain_text", fake_extract_ocr_plain_text)
    result = word_count_service.run_word_count_task_sync(
        task_id="task-mixed-pdf",
        display_no="000-mixed-pdf",
        directory_path=str(pdf_path),
        recursive=True,
        include_hidden=False,
        extensions=[".pdf"],
    )

    file_result = result["files"][0]
    assert list(captured["page_numbers"]) == [2]
    assert file_result["status"] == "counted"
    assert file_result["ocr_page_count"] == 1
    assert file_result["main_word_count"] == 8
    assert file_result["source_counts"] == {"pdf_page": 1, "pdf_ocr_page": 1}


def test_partial_ocr_failure_is_excluded_but_text_is_kept(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))
    image_path = tmp_path / "partial.png"
    _make_png(image_path)

    def fake_extract_ocr_plain_text(**_kwargs):
        return {
            "text": "成功页面",
            "total_pages": 2,
            "processed_pages": [1, 2],
            "failed_pages": [2],
            "page_results": [
                {"page_number": 1, "text": "成功页面", "raw_text": "", "blank": False, "error": ""},
                {"page_number": 2, "text": "", "raw_text": "", "blank": False, "error": "timeout"},
            ],
        }

    monkeypatch.setattr(word_count_service, "extract_ocr_plain_text", fake_extract_ocr_plain_text)
    result = word_count_service.run_word_count_task_sync(
        task_id="task-partial-ocr",
        display_no="000-partial-ocr",
        directory_path=str(image_path),
        recursive=True,
        include_hidden=False,
        extensions=[".png"],
    )

    file_result = result["files"][0]
    assert file_result["status"] == "failed"
    assert file_result["main_word_count"] == 0
    assert file_result["ocr_failed_pages"] == [2]
    assert result["summary"]["total_main_word_count"] == 0
    assert result["summary"]["ocr_failed_files"] == 1
    with ZipFile(tmp_path / result["ocr_text_archive"]) as archive:
        assert archive.read("partial.png.txt").decode("utf-8") == "成功页面"


def test_pdf2docx_ocr_supports_selected_pages_and_image_frames(tmp_path, monkeypatch):
    import fitz
    from PIL import Image

    pdf_path = tmp_path / "pages.pdf"
    pdf = fitz.open()
    for _ in range(3):
        pdf.new_page()
    pdf.save(pdf_path)
    pdf.close()

    monkeypatch.setattr(pdf2docx_module, "_ocr_single_image", lambda *_args, **_kwargs: "<p>Page two</p>")
    payload = pdf2docx_module.ocr_file(
        str(pdf_path),
        page_numbers=[2],
        return_metadata=True,
    )
    assert payload["processed_pages"] == [2]
    assert payload["failed_pages"] == []
    assert payload["page_results"][0]["page_number"] == 2

    first = Image.new("RGB", (20, 20), "white")
    second = Image.new("RGB", (20, 20), "black")
    tiff_path = tmp_path / "multi.tiff"
    first.save(tiff_path, save_all=True, append_images=[second])
    gif_path = tmp_path / "animated.gif"
    first.save(gif_path, save_all=True, append_images=[second], duration=100, loop=0)

    assert len(pdf2docx_module._image_frames_for_ocr(str(tiff_path))) == 2
    assert len(pdf2docx_module._image_frames_for_ocr(str(gif_path))) == 1


def test_directory_must_be_inside_allowed_roots(tmp_path, monkeypatch):
    allowed = tmp_path / "allowed"
    outside = tmp_path / "outside"
    allowed.mkdir()
    outside.mkdir()
    _allow_root(monkeypatch, allowed)

    with pytest.raises(PermissionError):
        word_count_service.prepare_word_count_request(directory_path=str(outside))


def test_local_paths_can_be_allowed_for_testing(tmp_path, monkeypatch):
    allowed = tmp_path / "allowed"
    outside = tmp_path / "outside"
    allowed.mkdir()
    outside.mkdir()
    _allow_root(monkeypatch, allowed)
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "True")

    result = word_count_service.prepare_word_count_request(directory_path=str(outside))

    assert Path(result["params"]["directory_path"]) == outside.resolve(strict=False)


def test_unc_server_root_allows_shares_on_same_server():
    root = Path("\\\\win-server\\")

    assert word_count_service._is_relative_to_path(Path("\\\\win-server\\服务器资料3\\项目"), root)
    assert not word_count_service._is_relative_to_path(Path("\\\\other-server\\服务器资料3\\项目"), root)


def test_unc_path_maps_to_container_mount_before_resolving(tmp_path, monkeypatch):
    mount_root = tmp_path / "mnt" / "win-server" / "服务器资料7"
    target_dir = mount_root / "客户" / "其他客户翻译任务" / "2026年" / "7月" / "美国玲翻译" / "0707" / "1. 原文"
    target_dir.mkdir(parents=True)
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps(["\\\\win-server\\服务器资料7\\"]))
    monkeypatch.setattr(
        word_count_service.settings,
        "WORD_COUNT_UNC_MOUNT_MAP_JSON",
        json.dumps({"\\\\win-server\\服务器资料7\\": str(mount_root)}, ensure_ascii=False),
    )
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")

    result = word_count_service.prepare_word_count_request(
        directory_path="\\\\win-server\\服务器资料7\\客户\\其他客户翻译任务\\2026年\\7月\\美国玲翻译\\0707\\1. 原文"
    )

    assert Path(result["params"]["directory_path"]) == target_dir.resolve(strict=False)
    assert Path(result["input_files"]["allowed_root"]) == mount_root.resolve(strict=False)


def test_unc_file_path_maps_to_single_file(tmp_path, monkeypatch):
    mount_root = tmp_path / "mnt" / "win-server" / "服务器资料7"
    target_dir = mount_root / "客户" / "其他客户翻译任务" / "2026年" / "7月" / "方舟移民" / "0701"
    target_dir.mkdir(parents=True)
    target_file = target_dir / "1.2.3.4.1 获奖通知-2020中国电子学会科技进步一等奖.docx"
    target_file.write_bytes(b"placeholder")
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps(["\\\\win-server\\服务器资料7\\"]))
    monkeypatch.setattr(
        word_count_service.settings,
        "WORD_COUNT_UNC_MOUNT_MAP_JSON",
        json.dumps({"\\\\win-server\\服务器资料7\\": str(mount_root)}, ensure_ascii=False),
    )
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")

    result = word_count_service.prepare_word_count_request(
        directory_path=(
            "\\\\win-server\\服务器资料7\\客户\\其他客户翻译任务\\2026年\\7月\\方舟移民\\0701\\"
            "1.2.3.4.1 获奖通知-2020中国电子学会科技进步一等奖.docx"
        )
    )

    assert result["params"]["input_kind"] == "file"
    assert Path(result["params"]["directory_path"]) == target_file.resolve(strict=False)
    assert Path(result["input_files"]["allowed_root"]) == mount_root.resolve(strict=False)


def test_unc_path_auto_maps_common_container_mount(tmp_path, monkeypatch):
    auto_root = tmp_path / "mnt"
    mount_root = auto_root / "win-server" / "服务器资料7"
    target_dir = mount_root / "客户" / "项目"
    target_dir.mkdir(parents=True)
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps(["\\\\win-server\\服务器资料7\\"]))
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_UNC_MOUNT_MAP_JSON", "")
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_UNC_AUTO_MOUNT_ROOTS_JSON", json.dumps([str(auto_root)], ensure_ascii=False))
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")

    result = word_count_service.prepare_word_count_request(directory_path="\\\\win-server\\服务器资料7\\客户\\项目")

    assert Path(result["params"]["directory_path"]) == target_dir.resolve(strict=False)
    assert Path(result["input_files"]["allowed_root"]) == mount_root.resolve(strict=False)


def test_unc_mapped_missing_path_reports_visible_mount_contents(tmp_path, monkeypatch):
    mount_root = tmp_path / "mnt" / "win-server" / "服务器资料7"
    (mount_root / "实际可见目录").mkdir(parents=True)
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps(["\\\\win-server\\服务器资料7\\"]))
    monkeypatch.setattr(
        word_count_service.settings,
        "WORD_COUNT_UNC_MOUNT_MAP_JSON",
        json.dumps({"\\\\win-server\\服务器资料7\\": str(mount_root)}, ensure_ascii=False),
    )
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")

    with pytest.raises(FileNotFoundError) as exc_info:
        word_count_service.prepare_word_count_request(directory_path="\\\\win-server\\服务器资料7\\客户\\项目")

    message = str(exc_info.value)
    assert str(mount_root.resolve(strict=False)) in message
    assert "实际可见目录" in message
    assert "Docker 当前挂载到了本地空目录" in message


def test_unc_path_error_lists_auto_mount_candidates(monkeypatch):
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_UNC_AUTO_MOUNT_ROOTS_JSON", json.dumps(["/mnt"], ensure_ascii=False))

    message = word_count_service._format_unc_auto_candidates("\\\\win-server\\服务器资料7\\客户\\项目")

    assert "/mnt/win-server/服务器资料7/客户/项目" in message.replace("\\", "/")
    assert "/mnt/服务器资料7/客户/项目" in message.replace("\\", "/")


def test_unc_mapping_still_requires_allowed_root(tmp_path, monkeypatch):
    mount_root = tmp_path / "mnt" / "win-server" / "服务器资料7"
    target_dir = mount_root / "客户"
    target_dir.mkdir(parents=True)
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps(["\\\\other-server\\share\\"]))
    monkeypatch.setattr(
        word_count_service.settings,
        "WORD_COUNT_UNC_MOUNT_MAP_JSON",
        json.dumps({"\\\\win-server\\服务器资料7\\": str(mount_root)}, ensure_ascii=False),
    )
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")

    with pytest.raises(PermissionError):
        word_count_service.prepare_word_count_request(directory_path="\\\\win-server\\服务器资料7\\客户")


def test_word_count_config_and_submit_endpoint(monkeypatch, tmp_path):
    _allow_root(monkeypatch, tmp_path)

    config = anyio.run(task_controller.get_word_count_page_config)
    assert config["allowed_roots"][0]["path"] == str(tmp_path)
    assert ".docx" in config["countable_extensions"]
    assert config["cad_support"]["direct_dxf"] is True
    assert ".dxf" in config["cad_support"]["supported_extensions"]
    assert "ODA_FILE_CONVERTER_PATH" not in json.dumps(config["cad_support"], ensure_ascii=False)
    assert config["default_ocr_mode"] == "auto"
    assert config["default_ocr_model"] in config["ocr_models"]

    captured = {}

    async def fake_submit_word_count_task(**kwargs):
        captured.update(kwargs)
        return SimpleNamespace(task_id="task-1", deduped=False)

    monkeypatch.setattr(task_controller.task_queue_service, "submit_word_count_task", fake_submit_word_count_task)

    body = task_controller.WordCountSubmitBody(
        directory_path=str(tmp_path),
        recursive=False,
        include_hidden=True,
        extensions=["docx"],
        ocr_mode="on",
        ocr_model="google/gemini-3.1-flash-lite",
    )
    result = anyio.run(task_controller.submit_word_count, body)

    assert result["status"] == "ACCEPTED"
    assert captured["directory_path"] == str(tmp_path)
    assert captured["recursive"] is False
    assert captured["include_hidden"] is True
    assert captured["extensions"] == ["docx"]
    assert captured["ocr_mode"] == "on"
    assert captured["ocr_model"] == "google/gemini-3.1-flash-lite"


def test_word_count_config_handles_inaccessible_unc_root(monkeypatch):
    unc_root = "\\\\win-server\\服务器资料7\\"
    expected_path = str(Path(unc_root).expanduser())
    original_exists = Path.exists

    def fake_exists(path):
        if str(path) == expected_path:
            raise OSError(1326, "用户名或密码不正确")
        return original_exists(path)

    monkeypatch.setattr(
        word_count_service.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps([unc_root], ensure_ascii=False)
    )
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_UNC_MOUNT_MAP_JSON", "")
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")
    monkeypatch.setattr(Path, "exists", fake_exists)

    config = word_count_service.get_word_count_config()

    assert config["allowed_roots"] == [
        {
            "path": expected_path,
            "exists": False,
            "scope_only": False,
            "mount_path": "",
        }
    ]


def test_word_count_page_keeps_ocr_model_fallback_options():
    project_root = Path(__file__).resolve().parents[1]
    html = (project_root / "static" / "word_count.html").read_text(encoding="utf-8")
    javascript = (project_root / "static" / "word_count.js").read_text(encoding="utf-8")

    assert '<option value="google/gemini-3-flash-preview" selected>' in html
    assert "FALLBACK_OCR_MODELS" in javascript
    assert "Object.keys(configuredModels).length ? configuredModels : FALLBACK_OCR_MODELS" in javascript


def test_word_count_page_explains_each_ocr_mode_on_hover():
    project_root = Path(__file__).resolve().parents[1]
    html = (project_root / "static" / "word_count.html").read_text(encoding="utf-8")
    javascript = (project_root / "static" / "word_count.js").read_text(encoding="utf-8")

    assert 'for="ocrModeAuto" tabindex="0" data-tooltip=' in html
    assert 'for="ocrModeOn" tabindex="0" data-tooltip=' in html
    assert 'for="ocrModeOff" tabindex="0" data-tooltip=' in html
    assert "event.target.closest?.('[data-tooltip]')" in javascript


def test_word_count_task_submission_uses_queue_dedupe(tmp_path, monkeypatch):
    _allow_root(monkeypatch, tmp_path)
    engine = create_engine(
        f"sqlite:///{tmp_path / 'tasks.db'}",
        connect_args={"check_same_thread": False},
    )
    Base.metadata.create_all(bind=engine)
    with engine.begin() as connection:
        connection.execute(
            text(
                "CREATE UNIQUE INDEX ux_task_active_request_fingerprint "
                "ON task (request_fingerprint) "
                "WHERE request_fingerprint IS NOT NULL AND status IN ('queued', 'running')"
            )
        )

    testing_session = sessionmaker(autocommit=False, autoflush=False, bind=engine)
    monkeypatch.setattr(queue_module, "SessionLocal", testing_session)

    service = queue_module.TaskQueueService()

    async def scenario():
        first = await service.submit_word_count_task(directory_path=str(tmp_path))
        second = await service.submit_word_count_task(directory_path=str(tmp_path))
        return first, second

    first, second = anyio.run(scenario)

    assert first.task_id == second.task_id
    assert first.deduped is False
    assert second.deduped is True
    with testing_session() as db:
        task = db.query(Task).one()
        assert task.task_type == "word_count"
        assert json.loads(task.input_files_json)["directory_path"] == str(tmp_path.resolve())


def test_prepare_word_count_upload_request_validates_file_and_config(monkeypatch):
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_UPLOAD_MAX_MB", 12)

    prepared = word_count_service.prepare_word_count_upload_request(
        filename="客户报价.docx",
        ocr_mode="auto",
    )

    assert prepared["filename"] == "客户报价.docx"
    assert prepared["extension"] == ".docx"
    assert prepared["params"]["input_source"] == "upload"
    assert prepared["params"]["input_kind"] == "file"
    assert prepared["params"]["ocr_enabled"] is True
    assert prepared["params"]["upload_max_file_mb"] == 12

    with pytest.raises(ValueError, match="不支持的文件格式"):
        word_count_service.prepare_word_count_upload_request(filename="报价.exe")


def test_uploaded_word_count_file_uses_trusted_upload_root_and_original_name(tmp_path, monkeypatch):
    upload_root = tmp_path / "uploads"
    stored_file = upload_root / "word_count" / "task-1" / "000001_input_abc123.txt"
    stored_file.parent.mkdir(parents=True)
    stored_file.write_text("中文报价 Hello 123", encoding="utf-8")
    monkeypatch.setattr(word_count_service.settings, "UPLOAD_DIR", str(upload_root))
    monkeypatch.setattr(word_count_service.settings, "OUTPUT_DIR", str(tmp_path / "outputs"))
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOW_LOCAL_PATHS", "False")
    monkeypatch.setattr(word_count_service.settings, "WORD_COUNT_ALLOWED_ROOTS_JSON", json.dumps([str(tmp_path / "other")]))

    result = word_count_service.run_word_count_task_sync(
        task_id="task-1",
        display_no="000001",
        directory_path=str(stored_file),
        recursive=False,
        include_hidden=False,
        extensions=[".txt"],
        ocr_mode="off",
        input_source="upload",
        original_filename="客户报价.txt",
    )

    assert result["input_source"] == "upload"
    assert result["files"][0]["relative_path"] == "客户报价.txt"
    assert result["files"][0]["filename"] == "客户报价.txt"
    assert result["summary"]["counted_files"] == 1

    outside_file = tmp_path / "outside.txt"
    outside_file.write_text("不能读取", encoding="utf-8")
    with pytest.raises(PermissionError, match="不在字数统计任务目录"):
        word_count_service.run_word_count_task_sync(
            task_id="task-2",
            display_no="000002",
            directory_path=str(outside_file),
            recursive=False,
            include_hidden=False,
            extensions=[".txt"],
            input_source="upload",
        )


def test_word_count_upload_submission_stages_file_and_dedupes(tmp_path, monkeypatch):
    engine = create_engine(
        f"sqlite:///{tmp_path / 'upload-tasks.db'}",
        connect_args={"check_same_thread": False},
    )
    Base.metadata.create_all(bind=engine)
    with engine.begin() as connection:
        connection.execute(
            text(
                "CREATE UNIQUE INDEX ux_task_active_request_fingerprint "
                "ON task (request_fingerprint) "
                "WHERE request_fingerprint IS NOT NULL AND status IN ('queued', 'running')"
            )
        )
    testing_session = sessionmaker(autocommit=False, autoflush=False, bind=engine)
    monkeypatch.setattr(queue_module, "SessionLocal", testing_session)
    monkeypatch.setattr(queue_module.settings, "UPLOAD_DIR", str(tmp_path / "uploads"))
    monkeypatch.setattr(queue_module.settings, "WORD_COUNT_UPLOAD_MAX_MB", 2)
    service = queue_module.TaskQueueService()

    async def scenario():
        first = await service.submit_word_count_upload_task(
            file=UploadFile(filename="客户报价.txt", file=io.BytesIO(b"hello 123")),
            ocr_mode="off",
        )
        second = await service.submit_word_count_upload_task(
            file=UploadFile(filename="客户报价.txt", file=io.BytesIO(b"hello 123")),
            ocr_mode="off",
        )
        return first, second

    first, second = anyio.run(scenario)

    assert first.task_id == second.task_id
    assert first.deduped is False
    assert second.deduped is True
    with testing_session() as db:
        task = db.query(Task).one()
        params = json.loads(task.params_json)
        input_files = json.loads(task.input_files_json)
    stored_path = Path(input_files["directory_path"])
    assert params["input_source"] == "upload"
    assert input_files["original_filename"] == "客户报价.txt"
    assert stored_path.is_file()
    assert stored_path.read_bytes() == b"hello 123"
    assert not list((tmp_path / "uploads" / "_tmp_uploads" / "word_count").glob("*"))


def test_word_count_upload_rejects_oversize_and_cleans_temp_file(tmp_path, monkeypatch):
    monkeypatch.setattr(queue_module.settings, "UPLOAD_DIR", str(tmp_path / "uploads"))
    monkeypatch.setattr(queue_module.settings, "WORD_COUNT_UPLOAD_MAX_MB", 1)
    service = queue_module.TaskQueueService()
    upload = UploadFile(filename="too-large.txt", file=io.BytesIO(b"x" * (1024 * 1024 + 1)))

    async def scenario():
        return await service.submit_word_count_upload_task(file=upload, ocr_mode="off")

    with pytest.raises(queue_module.UploadSizeLimitError, match="请改用共享路径统计"):
        anyio.run(scenario)

    temp_dir = tmp_path / "uploads" / "_tmp_uploads" / "word_count"
    assert not list(temp_dir.glob("*"))


def test_word_count_upload_endpoint_and_size_error(monkeypatch):
    captured = {}

    async def fake_submit_word_count_upload_task(**kwargs):
        captured.update(kwargs)
        return SimpleNamespace(task_id="upload-task", deduped=False)

    monkeypatch.setattr(
        task_controller.task_queue_service,
        "submit_word_count_upload_task",
        fake_submit_word_count_upload_task,
    )
    upload = UploadFile(filename="报价.docx", file=io.BytesIO(b"content"))
    result = anyio.run(
        task_controller.submit_word_count_upload,
        upload,
        "on",
        "google/gemini-3.1-flash-lite",
    )

    assert result["status"] == "ACCEPTED"
    assert captured["file"] is upload
    assert captured["ocr_mode"] == "on"
    assert captured["ocr_model"] == "google/gemini-3.1-flash-lite"

    async def reject_upload(**kwargs):
        raise queue_module.UploadSizeLimitError("文件超过上传限制 1 MB")

    monkeypatch.setattr(task_controller.task_queue_service, "submit_word_count_upload_task", reject_upload)
    with pytest.raises(task_controller.HTTPException) as exc_info:
        anyio.run(task_controller.submit_word_count_upload, upload, "off", None)
    assert exc_info.value.status_code == 413
