# -*- coding: utf-8 -*-

from __future__ import annotations

import zipfile
from pathlib import Path

import anyio
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches

from app.controller import task as task_controller
from app.service.english_variant_service import convert_text, get_converter
from app.service import office_text_transform_service
from app.service.office_text_transform_service import transform_office_file
from app.service.task_queue_service import TaskQueueService
from scripts.build_english_variant_dictionary import (
    DEFAULT_OUTPUT,
    DEFAULT_SOURCE,
    build_dictionary,
    compile_dictionary,
)


def test_dictionary_compiler_matches_committed_runtime_dictionary() -> None:
    payload = compile_dictionary(DEFAULT_SOURCE)
    assert payload["stats"] == {
        "raw_pairs": 2269,
        "unique_pairs": 1923,
        "british_to_american_rules": 1895,
        "british_to_american_ambiguous": 14,
        "american_to_british_rules": 1863,
        "american_to_british_ambiguous": 30,
    }
    assert build_dictionary(DEFAULT_SOURCE, DEFAULT_OUTPUT, check=True)


def test_converter_handles_phrases_case_boundaries_and_ambiguity() -> None:
    result = convert_text(
        "The AEROPLANE accessorised an american football. An aeroplanet, a check and a meter.",
        "american",
    )
    assert result["converted_text"] == (
        "The AIRPLANE accessorized an football. An aeroplanet, a check and a meter."
    )
    assert result["replacement_count"] == 3
    assert result["ambiguous_hit_count"] == 0

    british = convert_text("The airplane accessorized football; a check and a meter.", "british")
    assert british["converted_text"].startswith(
        "The aeroplane accessorised american football"
    )
    assert {item["term"] for item in british["ambiguous_hits"]} == {"check", "meter"}
    assert british["ambiguous_hit_count"] == 2


def test_dictionary_hash_participates_in_task_fingerprint() -> None:
    converter = get_converter()
    files = [{"role": "input", "filename": "sample.docx", "size": 3, "sha256": "abc"}]
    first = TaskQueueService.build_request_fingerprint(
        "english_variant",
        {"target_style": "british", "dictionary_sha256": converter.source_sha256},
        files,
    )
    second = TaskQueueService.build_request_fingerprint(
        "english_variant",
        {"target_style": "british", "dictionary_sha256": "changed"},
        files,
    )
    assert first != second


def test_docx_conversion_preserves_runs_tables_headers_and_footers(tmp_path: Path) -> None:
    source = tmp_path / "source.docx"
    output = tmp_path / "output.docx"
    document = Document()
    paragraph = document.add_paragraph()
    first = paragraph.add_run("aero")
    first.bold = True
    paragraph.add_run("plane and accessorised")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "airplane"
    document.sections[0].header.paragraphs[0].text = "aeroplane"
    document.sections[0].footer.paragraphs[0].text = "accessorised"
    document.save(source)

    summary = transform_office_file(source, output, "american")
    converted = Document(output)
    assert converted.paragraphs[0].text == "airplane and accessorized"
    assert converted.paragraphs[0].runs[0].bold is True
    assert converted.tables[0].cell(0, 0).text == "airplane"
    assert converted.sections[0].header.paragraphs[0].text == "airplane"
    assert converted.sections[0].footer.paragraphs[0].text == "accessorized"
    assert summary["replacement_count"] == 4


def test_docx_note_parts_are_converted(tmp_path: Path) -> None:
    document_path = tmp_path / "notes.docx"
    Document().save(document_path)
    footnotes_xml = b'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:footnotes xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:footnote w:id="1"><w:p><w:r><w:t>aero</w:t></w:r><w:r><w:t>plane</w:t></w:r></w:p></w:footnote>
</w:footnotes>'''
    with zipfile.ZipFile(document_path, "a") as package:
        package.writestr("word/footnotes.xml", footnotes_xml)

    summary = office_text_transform_service.ConversionSummary(get_converter(), "american")
    office_text_transform_service._transform_docx_note_parts(document_path, summary)

    with zipfile.ZipFile(document_path) as package:
        converted = package.read("word/footnotes.xml").decode("utf-8")
    assert "airplane" in converted
    assert summary.to_dict()["replacement_count"] == 1


def test_xlsx_conversion_skips_formulas_and_preserves_style(tmp_path: Path) -> None:
    source = tmp_path / "source.xlsx"
    output = tmp_path / "output.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = "The aeroplane was accessorised."
    sheet["A1"].font = Font(bold=True, color="FF0000")
    sheet["B1"] = "=LEN(A1)"
    sheet.merge_cells("C1:D1")
    sheet["C1"] = "airplane"
    workbook.save(source)

    summary = transform_office_file(source, output, "american")
    converted = load_workbook(output, data_only=False)
    sheet = converted.active
    assert sheet["A1"].value == "The airplane was accessorized."
    assert sheet["A1"].font.bold is True
    assert sheet["A1"].font.color.rgb == "00FF0000"
    assert sheet["B1"].value == "=LEN(A1)"
    assert sheet["C1"].value == "airplane"
    assert summary["replacement_count"] == 2


def test_pptx_conversion_handles_runs_tables_and_notes(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "output.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    paragraph = text_box.text_frame.paragraphs[0]
    paragraph.add_run().text = "aero"
    paragraph.add_run().text = "plane"
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "accessorised"
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "aeroplane"
    presentation.save(source)

    summary = transform_office_file(source, output, "american")
    converted = Presentation(output)
    slide = converted.slides[0]
    assert slide.shapes[0].text == "airplane"
    assert slide.shapes[1].table.cell(0, 0).text == "accessorized"
    assert slide.notes_slide.notes_text_frame.text == "airplane"
    assert summary["replacement_count"] == 3


def test_text_api_and_config_expose_dictionary_metadata() -> None:
    config = anyio.run(task_controller.get_english_variant_config)
    assert config["dictionary_version"] == "260505"
    assert config["stats"]["british_to_american_ambiguous"] == 14

    body = task_controller.EnglishVariantTextBody(
        text="The airplane.",
        target_style="british",
    )
    result = anyio.run(task_controller.run_english_variant_text, body)
    assert result["converted_text"] == "The aeroplane."


def test_output_file_is_registered_for_download() -> None:
    files = TaskQueueService._extract_output_files(
        "english_variant",
        {"output_file": "outputs/english_variant/1/sample_british.docx"},
        None,
        "sample.docx",
    )
    assert files == [
        {
            "name": "sample_british.docx",
            "path": "outputs/english_variant/1/sample_british.docx",
            "type": "output",
        }
    ]


def test_legacy_formats_reuse_existing_libreoffice_converters(
    tmp_path: Path,
    monkeypatch,
) -> None:
    calls: list[tuple[str, str]] = []

    def fake_converter(source: Path, target: Path) -> Path:
        calls.append((source.suffix.lower(), target.suffix.lower()))
        target.write_bytes(b"converted")
        return target

    monkeypatch.setattr(
        office_text_transform_service,
        "convert_doc_to_docx_via_libreoffice",
        fake_converter,
    )
    monkeypatch.setattr(
        office_text_transform_service,
        "convert_spreadsheet_to_xlsx_via_libreoffice",
        fake_converter,
    )
    monkeypatch.setattr(
        office_text_transform_service,
        "convert_presentation_to_pptx_via_libreoffice",
        fake_converter,
    )

    expected = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}
    for source_suffix, output_suffix in expected.items():
        source = tmp_path / f"sample{source_suffix}"
        source.write_bytes(b"legacy")
        normalized = office_text_transform_service._normalize_office_input(source, tmp_path)
        assert normalized.suffix == output_suffix

    assert calls == list(expected.items())
