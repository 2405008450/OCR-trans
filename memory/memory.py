import os
import re
import shutil
import threading
import queue
from datetime import datetime
import pandas as pd
from lxml import etree
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import OpenAI
import tkinter as tk
from tkinter import ttk, filedialog, scrolledtext, messagebox
import websockets
import asyncio
import time
import json
from dotenv import load_dotenv

# 加载项目根目录的 .env 统一配置
_project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
load_dotenv(os.path.join(_project_root, ".env"))

# ==========================================
# === ⚙️ 全局配置 ===
# ==========================================
ROW_BUCKET = 20_000
API_KEY = os.getenv("OPENROUTER_API_KEY", "")  # 从项目根 .env 统一管理
BASE_URL = os.getenv("OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1")



# 每块≤25000字：2份=50k, 3份=75k, 4份=100k, 5份=125k, 6份=150k, 7份=175k, 8份=200k
THRESHOLD_2_PARTS = 25000
THRESHOLD_3_PARTS = 50000
THRESHOLD_4_PARTS = 75000
THRESHOLD_5_PARTS = 100000
THRESHOLD_6_PARTS = 125000
THRESHOLD_7_PARTS = 150000
THRESHOLD_8_PARTS = 175000
BUFFER_CHARS = 2000
OUTPUT_DIR = "Result_Output"

# ==========================================
# === 🤖 可用模型配置 ===
# ==========================================
# OpenRouter 模型
OPENROUTER_MODELS = {
    "Google Gemini 2.5 Flash": {
        "id": "google/gemini-2.5-flash",
        "description": "建议先检查文章是否有目录，先将目录删除再处理",
        "max_output": 65536,
        "provider": "openrouter"
    },
    "Google Gemini 2.5 Pro": {
        "id": "google/gemini-2.5-pro",
        "description": "PPT推荐-增强，速度稍慢，100万上下文，65K输出",
        "max_output": 65536,
        "provider": "openrouter"
    },
    "Google: Gemini 3 Pro Preview": {
        "id": "google/gemini-3-pro-preview",
        "description": "最强推理，100万上下文，65K输出",
        "max_output": 65536,
        "provider": "openrouter"
    },
}

# 当前可用模型（根据提供商切换）
AVAILABLE_MODELS = OPENROUTER_MODELS.copy()
DEFAULT_MODEL = "Google Gemini 2.5 Flash"
DEFAULT_PROVIDER = "openrouter"  # 路智深已屏蔽，默认使用 OpenRouter



CHAPTER_PATTERNS = [
    r'^第[一二三四五六七八九十百千\d]+[章节篇部]', r'^Chapter\s*\d+', r'^CHAPTER\s*\d+',
    r'^\d+[\.、]\s*\S+', r'^[一二三四五六七八九十]+[、.]\s*\S+',
    r'^Part\s*\d+', r'^PART\s*\d+', r'^Section\s*\d+',
]

ORIG_PATTERNS = ['原文', '中文', 'source', 'original', 'chinese', 'cn', '源文']
TRANS_PATTERNS = ['译文', '英文', 'target', 'translation', 'english', 'en', '翻译']

# ==========================================
# === 🌍 多语言配置 ===
# ==========================================
SUPPORTED_LANGUAGES = {
    "中文": {
        "code": "zh",
        "english_name": "Chinese",
        "char_pattern": r'[\u4e00-\u9fa5]',  # 中文字符
        "word_based": False,  # 按字符计数
        "description": "中文（简体/繁体）"
    },
    "英语": {
        "code": "en",
        "english_name": "English",
        "char_pattern": r'\b[a-zA-Z]+\b',  # 英文单词
        "word_based": True,  # 按词计数
        "description": "English"
    },
    "西班牙语": {
        "code": "es",
        "english_name": "Spanish",
        "char_pattern": r'\b[a-zA-ZáéíóúüñÁÉÍÓÚÜÑ]+\b',
        "word_based": True,
        "description": "Español"
    },
    "葡语": {
        "code": "pt",
        "english_name": "Portuguese",
        "char_pattern": r'\b[a-zA-ZáéíóúâêôãõçÁÉÍÓÚÂÊÔÃÕÇ]+\b',
        "word_based": True,
        "description": "Português"
    },
    "日语": {
        "code": "ja",
        "english_name": "Japanese",
        "char_pattern": r'[\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FAF]',  # 平假名+片假名+汉字
        "word_based": False,  # 按字符计数
        "description": "日本語"
    },
    "俄语": {
        "code": "ru",
        "english_name": "Russian",
        "char_pattern": r'\b[а-яА-ЯёЁ]+\b',  # 西里尔字母
        "word_based": True,
        "description": "Русский"
    },
    "韩语": {
        "code": "ko",
        "english_name": "Korean",
        "char_pattern": r'[\uAC00-\uD7AF\u1100-\u11FF]',  # 韩文音节+字母
        "word_based": False,  # 按字符计数
        "description": "한국어"
    },
    "阿语": {
        "code": "ar",
        "english_name": "Arabic",
        "char_pattern": r'[\u0600-\u06FF\u0750-\u077F]+',  # 阿拉伯字母
        "word_based": True,
        "description": "العربية"
    },
    "法语": {
        "code": "fr",
        "english_name": "French",
        "char_pattern": r'\b[a-zA-ZàâäéèêëïîôùûüÿœæçÀÂÄÉÈÊËÏÎÔÙÛÜŸŒÆÇ]+\b',
        "word_based": True,
        "description": "Français"
    },
    "波兰语": {
        "code": "pl",
        "english_name": "Polish",
        "char_pattern": r'\b[a-zA-ZąćęłńóśźżĄĆĘŁŃÓŚŹŻ]+\b',
        "word_based": True,
        "description": "Polski"
    },
    "意大利语": {
        "code": "it",
        "english_name": "Italian",
        "char_pattern": r'\b[a-zA-ZàèéìíîòóùúÀÈÉÌÍÎÒÓÙÚ]+\b',
        "word_based": True,
        "description": "Italiano"
    },
    "德语": {
        "code": "de",
        "english_name": "German",
        "char_pattern": r'\b[a-zA-ZäöüßÄÖÜ]+\b',
        "word_based": True,
        "description": "Deutsch"
    },
}

# 默认语言设置
DEFAULT_SOURCE_LANG = "中文"
DEFAULT_TARGET_LANG = "英语"


# ==========================================
# === 提示词（保持原版不变）===
# ==========================================
def get_ppt_alignment_prompt(source_lang="中文", target_lang="英语"):
    """生成动态的PPT对齐提示词"""
    return f"""
你是一个「PPT 双语文本对齐器」。

当前任务：将 {source_lang} 原文与 {target_lang} 译文进行对齐。

═══════════════════════════════════════
【核心铁律 - 违反任意一条即为失败】
═══════════════════════════════════════
1. 禁止编造：左侧只能填原文中实际存在的文字，右侧只能填译文中实际存在的文字
2. 禁止占位符：绝对不允许出现 XXX、---、...、[空]、[无]、N/A 等任何占位符号
3. 禁止跨页：幻灯片 N 的内容只能与幻灯片 N 的内容配对
4. 分隔符独立：「---- 幻灯片 N ----」必须单独成行，不参与 ||| 配对

═══════════════════════════════════════
【输出格式 - 严格遵守】
═══════════════════════════════════════
---- 幻灯片 1 ----
{source_lang}行A ||| {target_lang}行A
{source_lang}行B ||| {target_lang}行B
{source_lang}行C |||
||| {target_lang}行D
---- 幻灯片 2 ----
...

格式说明：
- 分隔行「---- 幻灯片 N ----」独立一行，前后不加 |||
- 匹配成功：原文 ||| 译文
- 原文无对应译文：原文 |||
- 译文无对应原文：||| 译文

═══════════════════════════════════════
【对齐策略 - 按优先级执行】
═══════════════════════════════════════

第一步：锚点识别（同页内）
  ├─ 数字/日期锚点：12.15 ↔ December 15，2024年 ↔ 2024
  ├─ 专名锚点：公司名、人名、产品名（即使一侧是音译/意译）
  ├─ 编号锚点：① ↔ 1)，1. ↔ (1)，注1 ↔ Note 1
  └─ 标题锚点：通常字体大、位置靠上、内容概括性强

第二步：脚注特殊处理
  ├─ 识别脚注区域：通常在页面底部，带上标数字或特殊标记
  ├─ 编号与内容绑定：「1 XXX说明文字」整体视为一个脚注单元
  ├─ 配对原则：原文脚注1 ↔ 译文脚注1（按编号，非按位置）
  └─ 不可混配：脚注编号不能与正文内容配对

第三步：语义配对
  ├─ 含义最接近的行 1:1 配对
  ├─ 允许合理拆分：一长句 ↔ 两短句（语义完整时）
  ├─ 允许合理合并：两短句 ↔ 一长句（语义完整时）
  └─ 每个文本片段只能使用一次

第四步：处理剩余
  ├─ 原文有、译文无 → 「原文 |||」
  ├─ 译文有、原文无 → 「||| 译文」
  └─ 绝不填充任何虚构内容

═══════════════════════════════════════
【执行检查清单】
═══════════════════════════════════════
输出前逐项确认：
□ 分隔符是否独立成行？
□ 是否存在任何 XXX/---/.../[空] 等占位符？
□ 脚注编号是否与对应编号配对（而非与正文配对）？
□ 每个原文片段是否只出现一次？
□ 每个译文片段是否只出现一次？
□ 是否有跨页配对？

只输出对齐结果，不输出任何解释说明。
"""


# 保留默认的提示词用于兼容
PPT_ALIGNMENT_SYSTEM_PROMPT = get_ppt_alignment_prompt()


def get_docx_alignment_prompt(source_lang="中文", target_lang="英语"):
    """生成动态的文档对齐提示词"""
    source_info = SUPPORTED_LANGUAGES.get(source_lang, SUPPORTED_LANGUAGES["中文"])
    target_info = SUPPORTED_LANGUAGES.get(target_lang, SUPPORTED_LANGUAGES["英语"])

    # 判断原文是否为中日韩等需要特殊断句的语言
    source_is_cjk = source_lang in ["中文", "日语", "韩语"]

    # 根据原文语言类型确定断句标点
    if source_is_cjk:
        punctuation_marks = "。！？"
        punctuation_desc = "。！？（全角中日韩句末标点）"
    else:
        punctuation_marks = ". ! ?"
        punctuation_desc = ". ! ?（半角西文句末标点）"

    sentence_rule = f"""2. 🎯 断句规则：{source_lang}主导，{target_lang}跟随（强优先级）

   第一步：只看{source_lang}，按以下标点断句
   - 断句标点：{punctuation_desc}
   - 遇到这些标点就断开，形成一个{source_lang}片段

   第二步：{target_lang}匹配{source_lang}
   - 在 Stream B 中找出与该{source_lang}片段语义对应的{target_lang}部分
   - {target_lang}可能是1句、半句、或多句，都没关系
   - 译文位置可能有错位，按语义匹配即可
   - 只要语义完全覆盖该{source_lang}片段即可

3. 质量检查
   - 左边的每个{source_lang}片段是否都以正确的句末标点（{punctuation_marks}）结尾？
   - 右边{target_lang}的语义是否完全对应左边{source_lang}？
   - 是否有任何内容被修改或自行翻译？（必须为否）"""

    return f"""你是一个「双流文本对齐同步器」(Dual-Stream Aligner)。

核心任务
将 Stream A ({source_lang}原文) 和 Stream B ({target_lang}译文) 进行精确对齐。

⛔ 绝对铁律
1. 来源锁定（最高优先级）
   - "|||" 左侧 = 必须100%来自 Stream A，一个字都不能改
   - "|||" 右侧 = 必须100%来自 Stream B，一个字都不能改
   - 「禁止」自己翻译或编造任何内容
   - 即使发现原文有错别字、译文有翻译错误，也必须原样保留

{sentence_rule}

4. 空行处理规则
   - 如果原文或译文中有空行（连续换行符），这些空行已经被适当压缩保留
   - 空行通常用于分隔段落或章节
   - 在对齐时，忽略空行的位置差异，专注于有内容的文本对齐
   - 不要输出空行对应的对齐结果（空|||空）

输出格式
{source_lang}片段 ||| 对应的{target_lang}内容
"""


# 保留一个默认的提示词用于兼容
DOCX_ALIGNMENT_SYSTEM_PROMPT = get_docx_alignment_prompt()


def get_split_row_prompt(source_lang="中文"):
    """生成动态的分句对齐提示词"""
    # 判断原文是否为中日韩等需要特殊断句的语言
    source_is_cjk = source_lang in ["中文", "日语", "韩语"]

    if source_is_cjk:
        punctuation_marks = "。！？"
        punctuation_desc = "。！？（全角句末标点）"
        abbreviation_rule = ""
        numbering_rule = "1. / 2. / 11. / 1）/ 2）/ ①② 等数字编号"
    else:
        punctuation_marks = ". ! ?"
        punctuation_desc = ". ! ?（半角句末标点）"
        numbering_rule = """- 数字编号：1. / 2. / 11. / (1) / (2) 等
   - 字母编号：A. / B. / C. / a. / b. / c. / (A) / (B) 等
   - 罗马数字：I. / II. / III. / i. / ii. / iii. 等
   - 混合编号：1.1 / 1.2 / A.1 / A.2 等（在最后一个点后断开）"""
        abbreviation_rule = """
⚠️ 缩写例外（这些句号不表示句末，不要在此断开）：
- 学术引用：et al. / etc. / e.g. / i.e. / vs. / cf. / ibid.
- 称谓：Mr. / Mrs. / Ms. / Dr. / Prof. / Jr. / Sr. / St.
- 公司：Inc. / Ltd. / Co. / Corp. / L.L.C.
- 国家/组织：U.S. / U.K. / U.N. / E.U.
- 时间：a.m. / p.m. / A.M. / P.M.
- 引用标记：No. / Vol. / Fig. / Ch. / Sec. / pp. / p.
- 其他：approx. / est. / max. / min. / avg.

💡 区分编号与缩写的方法：
- 编号特征：单独的数字或字母 + 句号，如 "1." "A." "II."
- 缩写特征：多个字母组成的词 + 句号，如 "et al." "Dr." "Inc."
- 编号后面通常跟着正文内容，缩写后面通常还有更多文字"""

    # 针对西方语言添加更详细的断句示例
    if not source_is_cjk:
        example_section = """
## 示例

### 示例1：数字编号 + 作者信息
输入：
【原文】11. Piacentini MG, et al. A randomized controlled trial on the effect of Solidago virgaurea extract.
【译文】11.Piacentini MG，等。一项关于一枝黄花提取物影响的随机对照试验。

正确输出（3行）：
11. ||| 11.
Piacentini MG, et al. ||| Piacentini MG，等。
A randomized controlled trial on the effect of Solidago virgaurea extract. ||| 一项关于一枝黄花提取物影响的随机对照试验。

### 示例2：字母编号
输入：
【原文】A. Introduction. B. Methods. C. Results.
【译文】A. 引言。B. 方法。C. 结果。

正确输出（6行）：
A. ||| A.
Introduction. ||| 引言。
B. ||| B.
Methods. ||| 方法。
C. ||| C.
Results. ||| 结果。

### 示例3：混合编号
输入：
【原文】1.1 Background. 1.2 Objectives.
【译文】1.1 背景。1.2 目标。

正确输出（4行）：
1.1 ||| 1.1
Background. ||| 背景。
1.2 ||| 1.2
Objectives. ||| 目标。

### 关键判断规则：
- "A." "B." "1." "11." "1.1" → 编号，后面断开
- "et al." "Dr." "U.S." "Inc." → 缩写，不在此断开
- 编号 = 单独数字/字母 + 点号
- 缩写 = 多字母词汇 + 点号
"""
    else:
        example_section = ""
        numbering_rule = "1. / 2. / 11. / 1）/ 2）/ ①② 等数字编号"

    return f"""你是一个精确的「单行分句对齐器」。

## 任务
根据【原文】的分割点，将【译文】对应拆分对齐。

## ⛔ 铁律

### 1. 分割点识别（必须在以下位置断开）

#### 1.1 编号（编号后必须断开，编号单独成行）
   {numbering_rule}

#### 1.2 句末标点（{punctuation_desc}）
   每个句号/感叹号/问号后都要断开
{abbreviation_rule}

### 2. 禁止修改内容
- 左侧必须100%来自原文，一字不改
- 右侧必须100%来自译文，一字不改
- 禁止翻译、编造、补充任何内容

### 3. 译文对齐规则
- 根据语义将译文拆分，匹配原文的每个分割单元
- 如果译文标点与原文不一致，按语义对齐

### 4. 分割原则
- 编号单独成行（如 "11." / "A." / "1.1"）
- 作者信息单独成行（如 "Piacentini MG, et al."）
- 每个完整句子单独成行
- 即使片段很短也要独立成行
- 输出行数 = 原文分割单元数
{example_section}
## 输出格式
每行一对，用 ||| 分隔：
原文片段1 ||| 译文对应部分1
原文片段2 ||| 译文对应部分2

只输出对齐结果，不要任何解释。"""


# 保留默认的提示词用于兼容
SPLIT_ROW_SYSTEM_PROMPT = get_split_row_prompt()


def get_table_cell_split_prompt(source_lang="中文"):
    """生成表格单元格细粒度分句提示词 - 专门用于中英方向表格处理"""
    source_is_cjk = source_lang in ["中文", "日语", "韩语"]

    if source_is_cjk:
        punctuation_marks = "。！？"
        punctuation_desc = "。！？（全角句末标点）"
    else:
        punctuation_marks = ". ! ?"
        punctuation_desc = ". ! ?（半角句末标点）"

    return f"""你是一个精确的「表格单元格分句对齐器」。

## 任务
对表格单元格内的原文和译文进行细粒度分句对齐。

## 分句规则（必须严格遵守）

### 1. 断句位置（只在以下位置断开）
- **句末标点**：{punctuation_desc}
- **换行符**：当原文中有空行分隔时，在空行处断开

### 2. 【重要】不是断句点的情况
- **序号标点不断句**：如"1."、"2."、"3."、"①"、"②"、"(1)"、"a."等序号后的点号，这是序号的一部分，不是句末
- **序号必须保留在句子开头**：序号和后面的内容属于同一句
- 只有真正表示句子结束的标点（{punctuation_marks}）才断句

### 3. 分句原则
- 以原文的分句结构为准
- 每个完整句子（以{punctuation_marks}结尾）独立成行
- 当遇到空行分隔的段落时，按段落断开
- 保持原文和译文的句子一一对应

### 4. 禁止修改内容
- 左侧必须100%来自原文，一字不改
- 右侧必须100%来自译文，一字不改
- 禁止翻译、编造、补充任何内容
- 保留所有标点符号和序号

### 5. 译文对齐规则
- 根据语义将译文拆分，匹配原文的每个句子
- 如果译文标点与原文不完全一致，按语义对齐

## 示例1（带序号的多条内容）

### 输入：
【原文】
1.这是第一条说明。需要注意这个问题。

2.这是第二条说明；

【译文】
1. This is the first instruction. Please note this issue.
2. This is the second instruction;

### 正确输出：
1.这是第一条说明。 ||| 1. This is the first instruction.
需要注意这个问题。 ||| Please note this issue.
2.这是第二条说明； ||| 2. This is the second instruction;

### 错误输出（绝对禁止）：
1. ||| 1.
这是第一条说明。 ||| This is the first instruction.
（错误原因：把序号"1."单独拆分了，序号必须和后面的内容在一起）

## 示例2（普通多句内容）

### 输入：
【原文】
这是第一句。这是第二句！
这是换行后的第三句？

【译文】
This is the first sentence. This is the second sentence!
This is the third sentence after newline?

### 正确输出：
这是第一句。 ||| This is the first sentence.
这是第二句！ ||| This is the second sentence!
这是换行后的第三句？ ||| This is the third sentence after newline?

## 输出格式
每行一对，用 ||| 分隔：
原文句子1 ||| 译文对应部分1
原文句子2 ||| 译文对应部分2

只输出对齐结果，不要任何解释。"""


# 保留默认的表格分句提示词用于兼容
TABLE_CELL_SPLIT_PROMPT = get_table_cell_split_prompt()


# ==========================================
# === 📝 日志管理器 ===
# ==========================================
class LogManager:
    def __init__(self):
        self.log_queue = queue.Queue()
        self.exception_queue = queue.Queue()
        self.stream_queue = queue.Queue()

    def log(self, message, level="INFO"):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_queue.put(f"[{timestamp}] [{level}] {message}")

    def log_exception(self, message, data=None):
        timestamp = datetime.now().strftime("%H:%M:%S")
        exception_msg = f"[{timestamp}] ⚠️ {message}"
        if data:
            exception_msg += f"\n    数据: {data}"
        self.exception_queue.put(exception_msg)

    def log_stream(self, content):
        self.stream_queue.put(content)


log_manager = LogManager()


# ==========================================
# === 🛠️ 工具函数 ===
# ==========================================
def get_file_type(file_path):
    """获取文件类型"""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.docx':
        return 'docx'
    elif ext == '.doc':
        return 'doc'
    elif ext == '.pptx':
        return 'pptx'
    elif ext in ['.xlsx', '.xls']:
        return 'excel'
    return 'unknown'


def identify_column(df, patterns):
    """根据模式识别列名"""
    for col in df.columns:
        col_lower = str(col).lower()
        for pattern in patterns:
            if pattern in col_lower:
                return col
    return None


def convert_doc_to_docx(doc_path, output_dir=None):
    """将 .doc 文件转换为 .docx 文件"""
    try:
        import pythoncom
        from win32com import client as win32_client
        pythoncom.CoInitialize()

        if output_dir is None:
            output_dir = os.path.dirname(doc_path)

        base_name = os.path.splitext(os.path.basename(doc_path))[0]
        docx_path = os.path.join(output_dir, f"{base_name}_converted.docx")
        doc_path = os.path.abspath(doc_path)
        docx_path = os.path.abspath(docx_path)

        word = win32_client.Dispatch("Word.Application")
        word.Visible = False
        try:
            doc = word.Documents.Open(doc_path)
            doc.SaveAs2(docx_path, FileFormat=16)
            doc.Close()
            log_manager.log(f"已将 .doc 转换为 .docx: {docx_path}")
            return docx_path
        finally:
            word.Quit()
            pythoncom.CoUninitialize()
    except ImportError:
        log_manager.log_exception("需要安装 pywin32 来处理 .doc 文件", "pip install pywin32")
    except Exception as e:
        log_manager.log_exception(f"转换 .doc 文件失败: {e}", doc_path)
    return None


def read_excel_file(file_path):
    """读取 Excel 文件"""
    try:
        df = pd.read_excel(file_path)
        log_manager.log(f"读取 Excel 文件: {file_path}")
        log_manager.log(f"列名: {list(df.columns)}, 行数: {len(df)}")
        return df
    except Exception as e:
        log_manager.log_exception(f"读取 Excel 文件失败: {e}", file_path)
        return None


def _parse_xlsx_with_zipfile(file_path):
    """使用 zipfile 和 lxml 直接解析 xlsx 文件，完全绕过 openpyxl

    xlsx 文件本质是一个 ZIP 包，包含多个 XML 文件：
    - xl/workbook.xml: 工作簿信息（包含工作表名称）
    - xl/sharedStrings.xml: 共享字符串表
    - xl/worksheets/sheet1.xml 等: 各工作表数据
    """
    import zipfile

    NS = {
        'main': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'rel': 'http://schemas.openxmlformats.org/package/2006/relationships'
    }

    with zipfile.ZipFile(file_path, 'r') as zf:
        # 1. 读取共享字符串表
        shared_strings = []
        if 'xl/sharedStrings.xml' in zf.namelist():
            with zf.open('xl/sharedStrings.xml') as f:
                tree = etree.parse(f)
                for si in tree.findall('.//main:si', NS):
                    # 获取所有文本内容（包括富文本）
                    texts = []
                    for t in si.iter('{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t'):
                        if t.text:
                            texts.append(t.text)
                    shared_strings.append(''.join(texts))

        # 2. 读取工作簿获取工作表信息
        with zf.open('xl/workbook.xml') as f:
            wb_tree = etree.parse(f)
            sheets_info = []
            for sheet in wb_tree.findall('.//main:sheet', NS):
                sheet_name = sheet.get('name')
                sheet_id = sheet.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                sheets_info.append((sheet_name, sheet_id))

        # 3. 读取工作表关系文件获取实际文件路径
        sheet_files = {}
        with zf.open('xl/_rels/workbook.xml.rels') as f:
            rels_tree = etree.parse(f)
            for rel in rels_tree.findall('.//rel:Relationship', NS):
                rid = rel.get('Id')
                target = rel.get('Target')
                if 'worksheet' in target.lower():
                    sheet_files[rid] = 'xl/' + target.lstrip('/')

        # 4. 读取每个工作表的数据
        sheets_dict = {}
        for sheet_name, sheet_rid in sheets_info:
            sheet_path = sheet_files.get(sheet_rid)
            if not sheet_path or sheet_path not in zf.namelist():
                # 尝试常见路径
                for i in range(1, 100):
                    alt_path = f'xl/worksheets/sheet{i}.xml'
                    if alt_path in zf.namelist():
                        sheet_path = alt_path
                        break

            if not sheet_path or sheet_path not in zf.namelist():
                sheets_dict[sheet_name] = pd.DataFrame()
                continue

            with zf.open(sheet_path) as f:
                sheet_tree = etree.parse(f)
                rows_data = {}
                max_col = 0

                for row in sheet_tree.findall('.//main:row', NS):
                    row_idx = int(row.get('r', 0)) - 1  # 转为0索引
                    if row_idx < 0:
                        continue

                    for cell in row.findall('main:c', NS):
                        cell_ref = cell.get('r', '')
                        cell_type = cell.get('t', '')

                        # 解析列索引
                        col_str = ''.join(c for c in cell_ref if c.isalpha())
                        col_idx = 0
                        for c in col_str.upper():
                            col_idx = col_idx * 26 + (ord(c) - ord('A') + 1)
                        col_idx -= 1  # 转为0索引
                        max_col = max(max_col, col_idx + 1)

                        # 获取单元格值
                        value_elem = cell.find('main:v', NS)
                        value = None
                        if value_elem is not None and value_elem.text:
                            if cell_type == 's':  # 共享字符串
                                try:
                                    idx = int(value_elem.text)
                                    value = shared_strings[idx] if idx < len(shared_strings) else ''
                                except (ValueError, IndexError):
                                    value = ''
                            elif cell_type == 'b':  # 布尔值
                                value = value_elem.text == '1'
                            elif cell_type == 'inlineStr':  # 内联字符串
                                is_elem = cell.find('main:is', NS)
                                if is_elem is not None:
                                    texts = []
                                    for t in is_elem.iter(
                                            '{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t'):
                                        if t.text:
                                            texts.append(t.text)
                                    value = ''.join(texts)
                            else:  # 数字或其他
                                try:
                                    if '.' in value_elem.text:
                                        value = float(value_elem.text)
                                    else:
                                        value = int(value_elem.text)
                                except ValueError:
                                    value = value_elem.text
                        else:
                            # 检查内联字符串
                            is_elem = cell.find('main:is', NS)
                            if is_elem is not None:
                                texts = []
                                for t in is_elem.iter('{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t'):
                                    if t.text:
                                        texts.append(t.text)
                                value = ''.join(texts)

                        if row_idx not in rows_data:
                            rows_data[row_idx] = {}
                        rows_data[row_idx][col_idx] = value

                # 转换为 DataFrame
                if rows_data:
                    max_row = max(rows_data.keys()) + 1
                    data = [[None] * max_col for _ in range(max_row)]
                    for r_idx, cols in rows_data.items():
                        for c_idx, val in cols.items():
                            if c_idx < max_col:
                                data[r_idx][c_idx] = val
                    df = pd.DataFrame(data)
                else:
                    df = pd.DataFrame()

                sheets_dict[sheet_name] = df

        return sheets_dict


def read_excel_all_sheets(file_path):
    """读取 Excel 文件的所有工作簿，返回 {sheet_name: DataFrame} 字典

    使用多种方法尝试读取，以应对不同格式的Excel文件：
    1. 首先尝试标准 pandas 读取
    2. 如果失败（如 InlineFont 错误），使用 zipfile+lxml 直接解析
    """
    log_manager.log(f"读取 Excel 文件: {file_path}")

    # 方法1：尝试标准 pandas 读取
    try:
        excel_file = pd.ExcelFile(file_path, engine='openpyxl')
        sheet_names = excel_file.sheet_names
        log_manager.log(f"发现 {len(sheet_names)} 个工作簿: {sheet_names}")

        sheets_dict = {}
        for sheet_name in sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name, header=None, engine='openpyxl')
            sheets_dict[sheet_name] = df
            log_manager.log(f"  工作簿 '{sheet_name}': {df.shape[0]} 行 x {df.shape[1]} 列")

        return sheets_dict
    except Exception as e:
        log_manager.log(f"标准读取失败，尝试 XML 解析方法: {e}")

    # 方法2：使用 zipfile + lxml 直接解析 xlsx（完全绕过 openpyxl）
    try:
        sheets_dict = _parse_xlsx_with_zipfile(file_path)
        sheet_names = list(sheets_dict.keys())
        log_manager.log(f"(XML解析) 发现 {len(sheet_names)} 个工作簿: {sheet_names}")

        for sheet_name, df in sheets_dict.items():
            log_manager.log(f"  工作簿 '{sheet_name}': {df.shape[0]} 行 x {df.shape[1]} 列")

        return sheets_dict
    except Exception as e2:
        log_manager.log(f"XML 解析方法失败，尝试 openpyxl 只读模式: {e2}")

    # 方法3：使用 openpyxl 直接读取（只读模式）
    try:
        from openpyxl import load_workbook

        wb = load_workbook(file_path, data_only=True, read_only=True)
        sheet_names = wb.sheetnames
        log_manager.log(f"(openpyxl只读) 发现 {len(sheet_names)} 个工作簿: {sheet_names}")

        sheets_dict = {}
        for sheet_name in sheet_names:
            ws = wb[sheet_name]
            data = []
            for row in ws.iter_rows():
                row_data = [cell.value for cell in row]
                data.append(row_data)

            if data:
                df = pd.DataFrame(data)
            else:
                df = pd.DataFrame()

            sheets_dict[sheet_name] = df
            log_manager.log(f"  工作簿 '{sheet_name}': {df.shape[0]} 行 x {df.shape[1]} 列")

        wb.close()
        return sheets_dict
    except Exception as e3:
        log_manager.log_exception(f"读取 Excel 文件失败（所有方法均失败）: {e3}", file_path)
        return None


# ==========================================
# === 文档读取函数 ===
# ==========================================
def get_all_content_elements(doc):
    """获取文档所有内容元素"""
    body_elements = []
    if hasattr(doc, 'element') and hasattr(doc.element, 'body'):
        for child in doc.element.body.iterchildren():
            if child.tag.endswith('p'):
                body_elements.append(Paragraph(child, doc))
            elif child.tag.endswith('tbl'):
                body_elements.append(Table(child, doc))
    return body_elements


def get_text_count_by_language(text, lang_name):
    """根据语言类型统计文本字数/词数"""
    if not text:
        return 0

    lang_info = SUPPORTED_LANGUAGES.get(lang_name)
    if lang_info:
        pattern = lang_info['char_pattern']
        return len(re.findall(pattern, text))

    # 兼容旧的 lang_type 参数
    if lang_name == 'Chinese':
        return len(re.findall(r'[\u4e00-\u9fa5]', text))
    else:
        return len(re.findall(r'\b[a-zA-Z0-9-]+\b', text))


def get_element_text_count(element, lang_type):
    """获取元素的字数/词数"""
    text = ""
    if isinstance(element, Paragraph):
        text = element.text
    elif isinstance(element, Table):
        for row in element.rows:
            for cell in row.cells:
                text += cell.text + " "

    return get_text_count_by_language(text, lang_type)


def read_full_docx(file_path):
    """读取完整的docx文件内容，包括脚注、尾注等"""
    try:
        doc = Document(file_path)
        full_text = []
        consecutive_empty = 0  # 记录连续空段落数

        # 1. 按文档顺序遍历所有元素（段落和表格）
        # 使用 get_all_content_elements 的方式来保持原始顺序
        if hasattr(doc, 'element') and hasattr(doc.element, 'body'):
            for child in doc.element.body.iterchildren():
                if child.tag.endswith('p'):
                    # 段落元素
                    para = Paragraph(child, doc)
                    if para.text.strip():
                        full_text.append(para.text)
                        consecutive_empty = 0  # 重置空段落计数
                    else:
                        # 空段落：最多保留两个连续空行（用于表示段落分隔）
                        consecutive_empty += 1
                        if consecutive_empty <= 2 and full_text:  # 只在有内容后才添加空行
                            full_text.append("")  # 保留空行标记
                elif child.tag.endswith('tbl'):
                    consecutive_empty = 0  # 遇到表格重置计数
                    # 表格元素
                    table = Table(child, doc)
                    seen_cells = set()
                    for row in table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text and cell_text not in seen_cells:
                                full_text.append(cell_text)
                                seen_cells.add(cell_text)
        else:
            # 备用方案：如果无法使用上述方法，回退到原来的方式
            # 1. 段落文本
            consecutive_empty = 0
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
                    consecutive_empty = 0
                else:
                    # 空段落：最多保留两个连续空行
                    consecutive_empty += 1
                    if consecutive_empty <= 2 and full_text:
                        full_text.append("")

            # 2. 表格处理
            for table in doc.tables:
                seen_cells = set()
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen_cells:
                            full_text.append(cell_text)
                            seen_cells.add(cell_text)

        # 3. 页眉页脚 (包括段落、表格、文本框)
        for section in doc.sections:
            # 页眉段落
            for p in section.header.paragraphs:
                if p.text.strip():
                    full_text.append(p.text)
            # 页眉表格
            for table in section.header.tables:
                seen_header_cells = set()
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen_header_cells:
                            full_text.append(cell_text)
                            seen_header_cells.add(cell_text)
            # 页眉文本框
            if section.header._element is not None:
                header_xml = etree.tostring(section.header._element, encoding='unicode')
                header_root = etree.fromstring(header_xml.encode('utf-8'))
                nsmap_header = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                for txbx in header_root.xpath('.//w:txbxContent', namespaces=nsmap_header):
                    txbx_text = ''.join(txbx.xpath('.//w:t/text()', namespaces=nsmap_header))
                    if txbx_text.strip():
                        full_text.append(txbx_text.strip())

            # 页脚段落
            for p in section.footer.paragraphs:
                if p.text.strip():
                    full_text.append(p.text)
            # 页脚表格
            for table in section.footer.tables:
                seen_footer_cells = set()
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen_footer_cells:
                            full_text.append(cell_text)
                            seen_footer_cells.add(cell_text)
            # 页脚文本框
            if section.footer._element is not None:
                footer_xml = etree.tostring(section.footer._element, encoding='unicode')
                footer_root = etree.fromstring(footer_xml.encode('utf-8'))
                nsmap_footer = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                for txbx in footer_root.xpath('.//w:txbxContent', namespaces=nsmap_footer):
                    txbx_text = ''.join(txbx.xpath('.//w:t/text()', namespaces=nsmap_footer))
                    if txbx_text.strip():
                        full_text.append(txbx_text.strip())

        # 4. 文本框（支持传统格式和 DrawingML 格式）
        if hasattr(doc.element, 'xml'):
            xml = doc.element.xml
            root = etree.fromstring(xml.encode('utf-8'))

            # 扩展的命名空间，支持多种文本框格式
            nsmap = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'wps': 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape',
                'wpg': 'http://schemas.microsoft.com/office/word/2010/wordprocessingGroup',
                'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
                'w14': 'http://schemas.microsoft.com/office/word/2010/wordml',
            }

            textbox_texts = set()  # 用于去重

            # 方式1: 传统文本框 (w:txbxContent) - 按段落提取保留换行
            try:
                textbox_containers = root.xpath('.//w:txbxContent', namespaces=nsmap)
                for container in textbox_containers:
                    # 按段落提取，保留段落间的换行
                    paragraphs = container.xpath('.//w:p', namespaces=nsmap)
                    para_texts = []
                    for p in paragraphs:
                        # 合并段落内所有 <w:t> 节点的文本
                        p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged_text = '\n'.join(para_texts)
                    if merged_text.strip():
                        textbox_texts.add(merged_text.strip())
            except:
                pass

            # 方式2: DrawingML 文本框 (wps:txbx) - 按段落提取保留换行
            try:
                textbox_containers = root.xpath('.//wps:txbx', namespaces=nsmap)
                for container in textbox_containers:
                    # DrawingML 使用 <a:p> 作为段落
                    paragraphs = container.xpath('.//a:p', namespaces=nsmap)
                    para_texts = []
                    for p in paragraphs:
                        # 合并段落内所有 <a:t> 节点的文本
                        p_text = ''.join([t.text for t in p.xpath('.//a:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged_text = '\n'.join(para_texts)
                    if merged_text.strip():
                        textbox_texts.add(merged_text.strip())
            except:
                pass

            # 方式3: DrawingML 形状中的文本 (a:txBody) - 按段落提取保留换行
            try:
                shape_containers = root.xpath('.//a:txBody', namespaces=nsmap)
                for container in shape_containers:
                    # DrawingML 使用 <a:p> 作为段落
                    paragraphs = container.xpath('.//a:p', namespaces=nsmap)
                    para_texts = []
                    for p in paragraphs:
                        # 合并段落内所有 <a:t> 节点的文本
                        p_text = ''.join([t.text for t in p.xpath('.//a:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged_text = '\n'.join(para_texts)
                    if merged_text.strip():
                        textbox_texts.add(merged_text.strip())
            except:
                pass

            # 方式4: Word 2010+ 格式的文本框 (wps:wsp//wps:txbx) - 按段落提取保留换行
            try:
                textbox_containers = root.xpath('.//wps:wsp//wps:txbx', namespaces=nsmap)
                for container in textbox_containers:
                    # 按段落提取，保留段落间的换行
                    paragraphs = container.xpath('.//w:p', namespaces=nsmap)
                    para_texts = []
                    for p in paragraphs:
                        # 合并段落内所有 <w:t> 节点的文本
                        p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged_text = '\n'.join(para_texts)
                    if merged_text.strip():
                        textbox_texts.add(merged_text.strip())
            except:
                pass

            # 添加所有找到的文本框内容
            for text in textbox_texts:
                full_text.append(text)

        # 5. 脚注、尾注
        if hasattr(doc, 'part'):
            for rel in doc.part.rels.values():
                ref = rel.target_ref
                nsmap = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

                if "footnotes" in ref:
                    try:
                        root = etree.fromstring(rel.target_part.blob)
                        for fn in root.xpath('.//w:footnote', namespaces=nsmap):
                            fn_type = fn.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type')
                            if fn_type in ['separator', 'continuationSeparator']:
                                continue
                            fn_id = fn.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}id')
                            fn_texts = []
                            for p in fn.xpath('.//w:p', namespaces=nsmap):
                                p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap) if t.text])
                                if p_text.strip():
                                    fn_texts.append(p_text.strip())
                            fn_text = ' '.join(fn_texts)
                            if fn_text.strip():
                                full_text.append(fn_text.strip())
                    except:
                        pass

                elif "endnotes" in ref:
                    try:
                        root = etree.fromstring(rel.target_part.blob)
                        for en in root.xpath('.//w:endnote', namespaces=nsmap):
                            en_type = en.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type')
                            if en_type in ['separator', 'continuationSeparator']:
                                continue
                            en_id = en.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}id')
                            en_texts = []
                            for p in en.xpath('.//w:p', namespaces=nsmap):
                                p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap) if t.text])
                                if p_text.strip():
                                    en_texts.append(p_text.strip())
                            en_text = ' '.join(en_texts)
                            if en_text.strip():
                                full_text.append(en_text.strip())
                    except:
                        pass

        result = "\n".join(full_text)

        # 压缩多余的连续空行（最多保留2个连续换行，即1个空行）
        import re
        result = re.sub(r'\n{4,}', '\n\n\n', result)  # 3个\n = 2个空行，已经足够表示段落分隔

        if not result.strip():
            log_manager.log_exception("文档内容为空！", f"文件: {file_path}")
        return result

    except Exception as e:
        log_manager.log_exception(f"读取内容失败: {e}", f"文件: {file_path}")
        return ""


def _iter_group_shapes(group_shape, base_top=0, base_left=0):
    """展开组合形状，返回 (top, left, text)"""
    for sub in group_shape.shapes:
        top = base_top + (sub.top or 0)
        left = base_left + (sub.left or 0)

        if sub.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_group_shapes(sub, top, left)
            continue

        if hasattr(sub, "has_text_frame") and sub.has_text_frame:
            # 按段落提取，保留换行
            para_texts = []
            for paragraph in sub.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    para_texts.append(para_text)
            if para_texts:
                txt = '\n'.join(para_texts)
                yield (top, left, txt)

        if hasattr(sub, "has_table") and sub.has_table:
            for r_idx, row in enumerate(sub.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    # 表格单元格也按段落提取
                    para_texts = []
                    for paragraph in cell.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            para_texts.append(para_text)
                    if para_texts:
                        txt = '\n'.join(para_texts)
                        yield (top + r_idx * 1_000, left + c_idx * 1_000, txt)


def _extract_slide_items(slide):
    """提取单页幻灯片的所有文本项，返回已排序的文本列表"""
    items = []
    for shape in slide.shapes:
        top = shape.top or 0
        left = shape.left or 0

        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            # 按段落提取，保留换行
            para_texts = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    para_texts.append(para_text)
            if para_texts:
                txt = '\n'.join(para_texts)
                items.append((top, left, txt))

        if hasattr(shape, "has_table") and shape.has_table:
            for r_idx, row in enumerate(shape.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    # 表格单元格也按段落提取
                    para_texts = []
                    for paragraph in cell.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            para_texts.append(para_text)
                    if para_texts:
                        txt = '\n'.join(para_texts)
                        items.append((top + r_idx * 1_000, left + c_idx * 1_000, txt))

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            items.extend(_iter_group_shapes(shape, top, left))

    items.sort(key=lambda t: (t[0] // ROW_BUCKET, t[1]))
    return [t[2] for t in items]


def read_full_pptx(file_path):
    """读取PPT，按显示顺序提取所有文本"""
    try:
        prs = Presentation(file_path)
        all_lines = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            all_lines.append(f"---- 幻灯片 {slide_idx} ----")
            all_lines.extend(_extract_slide_items(slide))
        result = "\n".join(all_lines)
        if not result.strip():
            log_manager.log_exception("PPT内容为空！", f"文件: {file_path}")
        return result
    except Exception as e:
        log_manager.log_exception(f"读取PPT失败: {e}", f"文件: {file_path}")
        return ""


def read_full_excel(file_path):
    """读取 Excel 文件的所有内容"""
    try:
        df = pd.read_excel(file_path)
        texts = []
        for col in df.columns:
            texts.append(f"[列: {col}]")
            for val in df[col].dropna():
                if str(val).strip():
                    texts.append(str(val).strip())
        result = "\n".join(texts)
        if not result.strip():
            log_manager.log_exception("Excel内容为空！", f"文件: {file_path}")
        return result
    except Exception as e:
        log_manager.log_exception(f"读取Excel失败: {e}", f"文件: {file_path}")
        return ""


def read_file_content(file_path):
    """统一的文件读取接口"""
    file_type = get_file_type(file_path)
    if file_type == 'docx':
        return read_full_docx(file_path)
    elif file_type == 'pptx':
        return read_full_pptx(file_path)
    elif file_type == 'excel':
        return read_full_excel(file_path)
    elif file_type == 'doc':
        log_manager.log_exception(".doc 文件需要先转换", file_path)
        return ""
    else:
        log_manager.log_exception(f"不支持的文件类型", file_path)
        return ""


# ==========================================
# === 文档分析与分割 ===
# ==========================================
def analyze_document_structure(doc_path, lang_name):
    """分析文档结构，支持多语言"""
    file_type = get_file_type(doc_path)

    if file_type == 'pptx':
        text = read_full_pptx(doc_path)
        count = get_text_count_by_language(text, lang_name)
        return count, 0
    elif file_type == 'excel':
        df = read_excel_file(doc_path)
        if df is None:
            return 0, 0
        text = df.to_string()
        count = get_text_count_by_language(text, lang_name)
        return count, len(df)
    elif file_type == 'doc':
        log_manager.log_exception(".doc 文件需要先转换为 .docx")
        return 0, 0
    else:
        doc = Document(doc_path)
        elements = get_all_content_elements(doc)
        total_count = sum(get_element_text_count(el, lang_name) for el in elements)
        return total_count, len(elements)


def extract_text_from_elements(elements, start_idx, end_idx):
    """从元素中提取文本"""
    texts = []
    for i in range(start_idx, min(end_idx, len(elements))):
        elem = elements[i]
        if isinstance(elem, Paragraph):
            if elem.text.strip():
                texts.append(elem.text.strip())
        elif isinstance(elem, Table):
            for row in elem.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    texts.append(" | ".join(row_texts))
    return "\n".join(texts)


def delete_elements_in_range(doc, start_idx, end_idx):
    """删除指定范围的元素"""
    all_elements = get_all_content_elements(doc)
    elements_to_delete = []
    total = len(all_elements)

    for i in range(total):
        if start_idx <= i < end_idx:
            if isinstance(all_elements[i], Paragraph):
                elements_to_delete.append(all_elements[i]._element)
            elif isinstance(all_elements[i], Table):
                elements_to_delete.append(all_elements[i]._element)

    for el in elements_to_delete:
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


def find_element_index_by_char_count(element_counts, target_chars):
    """二分查找：根据目标字数找到最接近的元素索引

    修复：原逻辑只返回累计字数 < target 的最后一个元素，
    当存在超大元素（如大表格）时，多个目标值会映射到同一索引。
    现在额外比较下一个元素，返回累计字数更接近 target 的那个。
    """
    left, right = 0, len(element_counts) - 1
    best_idx = 0
    while left <= right:
        mid = (left + right) // 2
        if element_counts[mid] < target_chars:
            best_idx = mid
            left = mid + 1
        else:
            right = mid - 1
    # 检查下一个元素是否更接近目标字数
    if best_idx + 1 < len(element_counts):
        diff_before = target_chars - element_counts[best_idx]
        diff_after = element_counts[best_idx + 1] - target_chars
        if diff_after < diff_before:
            return best_idx + 1
    return best_idx


def _find_buffer_end(element_counts, split_idx, buffer_chars, direction='right'):
    """基于字数精确计算缓冲区边界（而非用全局平均估算元素数）

    direction='right': 从 split_idx 向右扩展 buffer_chars 字，返回结束索引
    direction='left':  从 split_idx 向左扩展 buffer_chars 字，返回开始索引
    """
    n = len(element_counts)
    if n == 0:
        return split_idx

    if direction == 'right':
        base_chars = element_counts[split_idx] if split_idx < n else element_counts[-1]
        target = base_chars + buffer_chars
        for i in range(split_idx + 1, n):
            if element_counts[i] >= target:
                return i + 1  # +1 因为 end 是开区间
        return n
    else:  # left
        base_chars = element_counts[split_idx] if split_idx < n else element_counts[-1]
        target = base_chars - buffer_chars
        if target <= 0:
            return 0
        for i in range(split_idx - 1, -1, -1):
            if element_counts[i] <= target:
                return i
        return 0


def smart_split_with_buffer(src_path, num_parts, output_dir, lang_type, buffer_chars=2000,
                            split_element_ratios=None):
    """智能分割文档：按字数均分 + 缓冲区重叠

    Args:
        split_element_ratios: 可选，由主文档计算出的分割比例列表（元素位置占比）。
                              提供时，本文档按相同比例分割，确保原文/译文内容对齐。
    Returns:
        (generated_files, part_info, element_ratios)
        element_ratios: 理想分割点的元素位置比例，可传给另一文档以保持同步。
    """
    doc = Document(src_path)
    elements = get_all_content_elements(doc)
    base_name = os.path.splitext(os.path.basename(src_path))[0]

    # 计算每个元素的累计字数
    element_counts = []
    cumulative_count = 0
    for elem in elements:
        count = get_element_text_count(elem, lang_type)
        cumulative_count += count
        element_counts.append(cumulative_count)

    total_count = cumulative_count
    if total_count == 0:
        log_manager.log_exception("文档字数为0，无法分割")
        return [], [], []

    target_per_part = total_count // num_parts

    log_manager.log(f"总字数: {total_count:,}, 目标每份: {target_per_part:,}, 缓冲区: {buffer_chars:,} 字")

    # 计算理想分割点
    if split_element_ratios is not None:
        # 从主文档的元素比例映射到本文档的元素索引
        log_manager.log(f"使用主文档分割比例: {[f'{r:.4f}' for r in split_element_ratios]}")
        ideal_splits = []
        for ratio in split_element_ratios:
            idx = max(0, min(int(ratio * len(elements)), len(elements) - 1))
            ideal_splits.append(idx)
    else:
        # 自主计算（作为主文档）
        ideal_splits = []
        for i in range(1, num_parts):
            target_chars = target_per_part * i
            split_idx = find_element_index_by_char_count(element_counts, target_chars)
            ideal_splits.append(split_idx)

    # 确保分割点严格递增（避免大元素导致多个分割点重叠）
    for i in range(1, len(ideal_splits)):
        if ideal_splits[i] <= ideal_splits[i - 1]:
            ideal_splits[i] = ideal_splits[i - 1] + 1
    # 确保不越界
    for i in range(len(ideal_splits)):
        ideal_splits[i] = min(ideal_splits[i], len(elements) - 1)

    # 计算元素位置比例（供另一文档使用）
    element_ratios = [idx / len(elements) for idx in ideal_splits] if elements else []

    log_manager.log(f"理想分割点索引: {ideal_splits}")
    for i, idx in enumerate(ideal_splits):
        chars_at_split = element_counts[idx] if idx < len(element_counts) else total_count
        log_manager.log(f"  分割点{i + 1}: 元素[{idx}]/{len(elements)}, 累计字数: {chars_at_split:,}")

    # 生成带缓冲的分割范围（基于字数精确计算缓冲区，而非全局平均）
    split_ranges = []
    for part_idx in range(num_parts):
        if part_idx == 0:
            start = 0
            if ideal_splits:
                end = _find_buffer_end(element_counts, ideal_splits[0], buffer_chars, 'right')
            else:
                end = len(elements)
        elif part_idx == num_parts - 1:
            start = _find_buffer_end(element_counts, ideal_splits[-1], buffer_chars, 'left')
            end = len(elements)
        else:
            start = _find_buffer_end(element_counts, ideal_splits[part_idx - 1], buffer_chars, 'left')
            end = _find_buffer_end(element_counts, ideal_splits[part_idx], buffer_chars, 'right')

        # 安全裁剪
        start = max(0, min(start, len(elements) - 1))
        end = max(start + 1, min(end, len(elements)))
        split_ranges.append((start, end))

    for i, (s, e) in enumerate(split_ranges):
        part_chars = element_counts[min(e, len(element_counts)) - 1] - (
            element_counts[s - 1] if s > 0 else 0) if e > s else 0
        log_manager.log(f"  Part{i + 1}: 元素[{s}:{e}], 约 {part_chars:,} 字")

    # 生成分割后的文件
    generated_files = []
    part_info = []

    for i, (start_idx, end_idx) in enumerate(split_ranges):
        part_num = i + 1
        dest_filename = f"{base_name}_Part{part_num}.docx"
        dest_path = os.path.join(output_dir, dest_filename)

        shutil.copy2(src_path, dest_path)
        doc_copy = Document(dest_path)

        total_elems = len(get_all_content_elements(doc_copy))
        delete_elements_in_range(doc_copy, end_idx, total_elems + 5000)
        delete_elements_in_range(doc_copy, 0, start_idx)
        doc_copy.save(dest_path)

        first_text = extract_text_from_elements(elements, start_idx, min(start_idx + 3, end_idx))
        last_text = extract_text_from_elements(elements, max(start_idx, end_idx - 3), end_idx)

        part_info.append({
            'path': dest_path,
            'first_anchor': first_text[:200] if first_text else "",
            'last_anchor': last_text[-200:] if last_text else "",
            'start_idx': start_idx,
            'end_idx': end_idx
        })

        generated_files.append(dest_path)
        log_manager.log(f"生成: {dest_filename}")

    return generated_files, part_info, element_ratios


# ==========================================
# === 🔍 对齐质量检查器 ===
# ==========================================
class AlignmentChecker:
    @staticmethod
    def check_language_consistency(df, source_lang="中文", target_lang="英语"):
        """检查语言一致性，支持多语言"""
        issues = []

        source_info = SUPPORTED_LANGUAGES.get(source_lang, SUPPORTED_LANGUAGES["中文"])
        target_info = SUPPORTED_LANGUAGES.get(target_lang, SUPPORTED_LANGUAGES["英语"])

        source_pattern = source_info['char_pattern']
        target_pattern = target_info['char_pattern']

        for idx, row in df.iterrows():
            original = str(row.get('原文', ''))
            trans = str(row.get('译文', ''))

            if not original or not trans:
                continue

            # 计算原文中源语言字符的比例
            source_chars_in_original = len(re.findall(source_pattern, original))
            target_chars_in_original = len(re.findall(target_pattern, original))

            # 计算译文中目标语言字符的比例
            source_chars_in_trans = len(re.findall(source_pattern, trans))
            target_chars_in_trans = len(re.findall(target_pattern, trans))

            total_original = len(original) if original else 1
            total_trans = len(trans) if trans else 1

            # 检测是否存在语言错位
            source_ratio_original = source_chars_in_original / total_original
            target_ratio_trans = target_chars_in_trans / total_trans

            # 如果原文中源语言占比很低，而译文中源语言占比很高，可能是错位
            source_ratio_trans = source_chars_in_trans / total_trans
            target_ratio_original = target_chars_in_original / total_original

            if source_ratio_original < 0.2 and source_ratio_trans > 0.4:
                issues.append({
                    'row': idx + 1, 'type': '语言错位',
                    'detail': f'原文疑似{target_lang}，译文疑似{source_lang}',
                    'original_text': original, 'trans_text': trans
                })
            elif target_ratio_original > 0.4 and target_ratio_trans < 0.2:
                issues.append({
                    'row': idx + 1, 'type': '语言错位',
                    'detail': f'原文疑似{target_lang}，译文疑似{source_lang}',
                    'original_text': original, 'trans_text': trans
                })
        return issues

    @staticmethod
    def check_length_anomaly(df, threshold_ratio=5):
        """检查长度异常"""
        issues = []
        for idx, row in df.iterrows():
            original = str(row.get('原文', ''))
            trans = str(row.get('译文', ''))

            len_orig = len(original)
            len_trans = len(trans)

            if len_orig == 0 or len_trans == 0:
                continue

            ratio = max(len_orig, len_trans) / min(len_orig, len_trans)
            if ratio > threshold_ratio:
                issues.append({
                    'row': idx + 1, 'type': '长度异常',
                    'detail': f'长度比 {ratio:.1f}:1',
                    'original_text': original, 'trans_text': trans
                })
        return issues

    @staticmethod
    def full_check(df, source_lang="中文", target_lang="英语"):
        """执行完整检查"""
        all_issues = []
        all_issues.extend(AlignmentChecker.check_language_consistency(df, source_lang, target_lang))
        all_issues.extend(AlignmentChecker.check_length_anomaly(df))
        return all_issues


def save_issues_report(issues, output_path):
    if not issues:
        return
    report_data = [{
        '行号': issue.get('row', ''),
        '问题类型': issue.get('type', ''),
        '问题详情': issue.get('detail', ''),
        '原文': issue.get('original_text', ''),
        '译文': issue.get('trans_text', '')
    } for issue in issues]
    pd.DataFrame(report_data).to_excel(output_path, index=False)
    log_manager.log(f"问题报告已保存: {output_path}")


# ==========================================
# === 🤖 核心 AI 对齐 ===
# ==========================================
def call_openrouter_stream(system_prompt, user_prompt, model_id, max_output_tokens, filename=""):
    """OpenRouter API 流式调用"""
    client = OpenAI(base_url=BASE_URL, api_key=API_KEY)

    try:
        log_manager.log("请求 OpenRouter API...")
        stream = client.chat.completions.create(
            model=model_id,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.1,
            max_tokens=max_output_tokens,
            stream=True,
            timeout=600.0,
            extra_headers={"HTTP-Referer": "local-debug", "X-Title": "Doc-Aligner"},
        )

        full_response_text = ""
        log_manager.log("接收数据流...")
        log_manager.log_stream("\n" + "=" * 50 + f" {filename} " + "=" * 50 + "\n")

        for chunk in stream:
            if not chunk.choices:
                continue
            delta = chunk.choices[0].delta
            if hasattr(delta, "content") and delta.content:
                content = delta.content
                log_manager.log_stream(content)
                full_response_text += content

        log_manager.log_stream("\n" + "=" * 50 + " 输出结束 " + "=" * 50 + "\n")
        return full_response_text

    except Exception as e:
        log_manager.log_exception(f"OpenRouter API调用失败", str(e))
        import traceback
        log_manager.log_exception("详细堆栈", traceback.format_exc())
        return None




def call_llm_stream(system_prompt, user_prompt, model_id, filename=""):
    """统一的LLM流式调用 - 自动检测提供商"""
    # 获取模型的最大输出 tokens 和提供商
    max_output_tokens = 65536
    provider = "openrouter"  # 默认提供商

    for model_info in AVAILABLE_MODELS.values():
        if model_info['id'] == model_id:
            max_output_tokens = model_info['max_output']
            provider = model_info.get('provider', 'openrouter')
            break

    log_manager.log(f"使用提供商: {provider}")

    return call_openrouter_stream(system_prompt, user_prompt, model_id, max_output_tokens, filename)


def parse_alignment_response(response_text):
    """解析对齐响应"""
    if not response_text:
        return []

    cleaned_text = response_text.replace('\r\n', '\n').replace('\r', '\n')
    try:
        # 注释掉会导致多行合并的正则，保留每一行的独立性
        # cleaned_text = re.sub(r'\s*\|\|\|\s*\n+\s*', ' ||| ', cleaned_text)
        # cleaned_text = re.sub(r'\n+\s*\|\|\|\s*', ' ||| ', cleaned_text)
        cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
    except Exception as e:
        log_manager.log_exception(f"正则替换失败", str(e))

    lines = cleaned_text.splitlines()
    data = []
    pending_line = ""

    for line_num, line in enumerate(lines, 1):
        line = line.strip()
        if not line or line.startswith('```'):
            continue

        if pending_line:
            line = pending_line + " " + line
            pending_line = ""

        if "|||" in line:
            parts = line.split("|||", 1)
            if len(parts) >= 2:
                original = parts[0].strip()
                trans = parts[1].strip()
                # original = re.sub(r'^\d+\.\s*', '', original)
                # trans = re.sub(r'^\d+\.\s*', '', trans)
                if original or trans:
                    data.append({"原文": original, "译文": trans})
        else:
            if len(line) < 50:
                pending_line = line
            else:
                log_manager.log_exception(f"第 {line_num} 行：缺少分隔符", line[:100])

    return data


# ==========================================
# === 🔤 英文后处理分句 ===
# ==========================================
# 常见英文缩写（句号不表示句末）
ENGLISH_ABBREVIATIONS = {
    # 学术/引用
    'et al', 'etc', 'e.g', 'i.e', 'vs', 'cf', 'ibid', 'op', 'cit',
    # 称谓
    'mr', 'mrs', 'ms', 'dr', 'prof', 'jr', 'sr', 'st',
    # 公司/组织
    'inc', 'ltd', 'co', 'corp', 'llc', 'l.l.c',
    # 国家/地区缩写（完整形式）
    'u.s', 'u.k', 'u.n', 'e.u', 'u.s.a',
    # 时间/度量
    'a.m', 'p.m',
    # 引用标记
    'no', 'vol', 'fig', 'ch', 'sec', 'pp', 'approx', 'est', 'max', 'min', 'avg',
}

# 多字母缩写的结尾字母（用于检测 "U.S." "U.K." 等）
# 这些是常见多字母缩写中的最后一个字母
ABBREVIATION_ENDING_PATTERNS = {
    's': ['u.s', 'u.s.a'],  # U.S., U.S.A.
    'k': ['u.k'],  # U.K.
    'n': ['u.n'],  # U.N.
    'u': ['e.u'],  # E.U.
}


def is_abbreviation_period(text, period_pos):
    """判断句号是否属于缩写或连续编号的中间点（不是真正的句末）

    返回 True 的情况（不计入句子数量）：
    - 缩写：多字母词汇 + 句号，如 "et al." "Dr." "U.S."
    - 连续编号的中间点：如 "7.2.1" 中的前两个点
    - 编号与内容连写：如 "8.Al-Snafi" 中的点（无空格分隔）

    返回 False 的情况（计入句子数量，需要断开）：
    - 编号：单独的数字或字母 + 句号 + 空格，如 "1. " "A. " "II. "
    """
    if period_pos <= 0:
        return False

    # 获取句号前后的文本
    before = text[:period_pos]
    after = text[period_pos + 1:] if period_pos + 1 < len(text) else ""

    # 关键检查1：如果句号后面紧跟着数字，这是连续编号的中间点（如 7.2.1 中的点）
    if after and after[0].isdigit():
        return True

    # 关键检查2：如果句号后面紧跟着字母（无空格），这是编号与内容连写（如 "8.Al-Snafi"）
    # 这种情况不应该分割
    if after and after[0].isalpha():
        return True

    # 找到句号前的最后一个"词"
    words = before.split()
    if not words:
        return False

    last_token = words[-1]
    last_word_lower = last_token.lower().rstrip('.,;:')

    # 1. 检查是否是多字母缩写（如 "et al", "Dr", "Inc"）
    if last_word_lower in ENGLISH_ABBREVIATIONS:
        return True

    # 2. 检查是否是 X.Y. 缩写模式（如 U.S.）
    # 如果当前 token 中包含点号，可能是缩写的一部分
    if '.' in last_token:
        return True

    # 3. 检查是否是单个字母，且前面是 "X." 模式（如 "U.S." 中的 "S"）
    if len(last_word_lower) == 1 and last_word_lower.isalpha():
        if len(words) >= 2:
            prev_token = words[-2].lower()
            if re.match(r'^[a-z]\.$', prev_token):
                return True
        # 独立的单个字母（如 "A" "B"）→ 是编号，需要断开
        return False

    # 4. 检查是否是纯数字（如 "11"）→ 不是缩写，需要断开
    if last_word_lower.isdigit():
        return False

    # 5. 检查是否是罗马数字
    roman_pattern = r'^(i{1,3}|iv|vi{0,3}|ix|xi{0,3}|xiv|xvi{0,3}|xix|xxi{0,3})$'
    if re.match(roman_pattern, last_word_lower):
        return False  # 罗马数字是编号，需要断开

    # 6. 默认：不是已知的缩写
    return False


def count_real_sentences_english(text):
    """计算英文文本中的真实句子数量（排除缩写）"""
    if not text:
        return 0

    count = 0
    i = 0
    while i < len(text):
        if text[i] in '.!?':
            # 检查是否是缩写
            if text[i] == '.' and is_abbreviation_period(text, i):
                i += 1
                continue
            count += 1
        i += 1

    return max(1, count)


def has_numbering_pattern(text):
    """检测文本是否包含需要分割的编号模式

    核心规则：只要检测到"末尾是句号+空格"的编号模式，就需要分割

    需要分割：
    - "7.2. TITLE" → 检测到 "2. "，需要分割
    - "A. Introduction" → 检测到 "A. "，需要分割
    - "11. Author" → 检测到 "11. "，需要分割

    不需要分割：
    - "7.2.1 TITLE" → "7.2.1" 末尾是数字不是句号，不匹配
    """
    if not text:
        return False

    # 编号模式正则表达式
    # 关键：必须是 "句号+空格" 结尾，不能是 "数字+空格"
    numbering_patterns = [
        # 数字编号：1. / 11. / 7.2.（末尾必须是句号+空格）
        # 注意：7.2.1 不匹配，因为 1 后面是空格不是句号
        r'\d+\.\s',

        # 字母编号：A. / B. / a. / b.（单个字母+句号+空格）
        # 使用负向后瞻排除：
        # - 前面是字母的情况（如 "et al." 中的 "l"）
        # - 前面是句号的情况（如 "U.S." 中的 "S"）
        r'(?<![A-Za-z.])[A-Za-z]\.\s',

        # 带括号的编号：(1) / (A) / (a)
        r'\([0-9]+\)\s',
        r'\([A-Za-z]\)\s',

        # 罗马数字：I. / II. / III. / IV.（罗马数字+句号+空格）
        r'(?<![A-Za-z])[IVXivx]+\.\s',
    ]

    # 只要匹配到任一编号模式，就需要分割
    for pattern in numbering_patterns:
        if re.search(pattern, text):
            return True

    return False


def needs_english_post_split(orig_text, source_lang):
    """判断是否需要英文后处理分句

    触发条件（满足任一即触发）：
    1. 文本中包含编号模式（如 A. / 1. / I. 等）
    2. 文本中有多个真实句子（排除缩写后的句号数 > 1）
    """
    # 只对非CJK语言进行后处理
    if source_lang in ["中文", "日语", "韩语"]:
        return False

    if not orig_text:
        return False

    # 条件1：检测是否包含编号模式
    if has_numbering_pattern(orig_text):
        return True

    # 条件2：计算真实句子数量
    real_sentence_count = count_real_sentences_english(orig_text)
    if real_sentence_count > 1:
        return True

    return False


def post_process_english_split(data, model_id, source_lang="英语", target_lang="中文", enable_ai_split=True):
    """
    后处理：对英文原文进行进一步的细粒度分句

    Args:
        data: 解析后的对齐数据列表 [{"原文": ..., "译文": ...}, ...]
        model_id: AI模型ID（用于分句时调用）
        source_lang: 源语言
        target_lang: 目标语言
        enable_ai_split: 是否启用AI分句（True=调用AI，False=仅检测但不分割）

    Returns:
        处理后的数据列表
    """
    # 只对非CJK源语言进行后处理
    if source_lang in ["中文", "日语", "韩语"]:
        return data

    if not data:
        return data

    result = []
    split_count = 0

    for idx, row in enumerate(data):
        orig = row.get('原文', '')
        trans = row.get('译文', '')

        # 检查是否需要分句
        if needs_english_post_split(orig, source_lang):
            if enable_ai_split:
                # 调用AI进行分句
                log_manager.log(f"后处理分句: 第 {idx + 1} 行需要进一步细分")
                split_results = split_row_with_ai(orig, trans, model_id, f"后处理-{idx + 1}", source_lang)

                if split_results and len(split_results) > 1:
                    result.extend(split_results)
                    split_count += 1
                    log_manager.log(f"  ✅ 1 行 → {len(split_results)} 行")
                else:
                    # AI分句失败，保留原数据
                    result.append(row)
            else:
                # 不启用AI，只记录需要分句的行
                log_manager.log_exception(f"第 {idx + 1} 行可能需要进一步分句", orig[:100])
                result.append(row)
        else:
            result.append(row)

    if split_count > 0:
        log_manager.log(f"后处理分句完成: {split_count} 行被细分，总行数 {len(data)} → {len(result)}")

    return result


def run_llm_alignment(file_original_path, file_trans_path, output_excel_path, model_id,
                      anchor_info_orig=None, anchor_info_trans=None, system_prompt_override=None,
                      source_lang="中文", target_lang="英语", enable_post_split=True,
                      enable_table_separate_processing=False):
    """
    对齐函数 - 用于Word/PPT文档对齐

    参数:
        enable_table_separate_processing: 是否启用Word表格单独处理（按单元格位置匹配 + LLM细粒度分句）
                                          默认关闭，启用后会将Word表格按单元格位置匹配处理

    注意：Excel文件应使用 process_excel_dual_file_alignment() 函数处理，不要使用本函数
    """
    filename = os.path.basename(file_original_path)
    file_type = get_file_type(file_original_path)
    log_manager.log(f"正在 AI 对齐: {filename}")
    log_manager.log(f"使用模型: {model_id}")
    log_manager.log(f"语言对: {source_lang} → {target_lang}")

    all_data = []

    # 检查是否为 Word 文档且包含表格，且启用了表格单独处理
    if file_type == 'docx' and enable_table_separate_processing:
        orig_has_tables = has_docx_tables(file_original_path)
        trans_has_tables = has_docx_tables(file_trans_path)

        if orig_has_tables or trans_has_tables:
            log_manager.log("📊 检测到表格内容，启用表格单独处理模式")
            log_manager.log("   处理策略：按单元格绝对位置匹配 + LLM细粒度分句")

            # 1. 先处理表格部分（按单元格位置匹配 + LLM分句）
            table_results = process_docx_tables_alignment(
                file_original_path, file_trans_path, model_id,
                source_lang=source_lang, target_lang=target_lang
            )

            if table_results:
                log_manager.log(f"   表格处理结果: {len(table_results)} 行")
                # 移除来源列，只保留原文和译文
                for result in table_results:
                    all_data.append({
                        "原文": result.get("原文", ""),
                        "译文": result.get("译文", "")
                    })

            # 2. 再处理非表格部分
            text_original = read_docx_without_tables(file_original_path)
            text_trans = read_docx_without_tables(file_trans_path)

            if text_original and text_trans:
                log_manager.log("📝 处理非表格内容...")

                # 根据语言生成动态提示词
                if system_prompt_override:
                    system_prompt = system_prompt_override
                else:
                    system_prompt = get_docx_alignment_prompt(source_lang, target_lang)

                anchor_hint = ""
                if anchor_info_orig and anchor_info_trans:
                    anchor_hint = f"""
## 上下文提示
这是文档的一个片段：
- 原文开头: "{anchor_info_orig.get('first_anchor', '')[:100]}"
- 译文开头: "{anchor_info_trans.get('first_anchor', '')[:100]}"
"""

                user_prompt = f"""{anchor_hint}

<Stream_A_Original>
{text_original}
</Stream_A_Original>

<Stream_B_Translation>
{text_trans}
</Stream_B_Translation>

请严格按规则输出对齐结果：
"""

                response = call_llm_stream(system_prompt, user_prompt, model_id, f"{filename}-非表格")

                if response:
                    last_500 = response[-500:] if len(response) > 500 else response
                    if '|||' not in last_500:
                        log_manager.log_exception("⚠️ 非表格部分输出可能被截断")

                    paragraph_data = parse_alignment_response(response)

                    if paragraph_data:
                        # 对非CJK源语言进行后处理分句
                        if enable_post_split and source_lang not in ["中文", "日语", "韩语"]:
                            log_manager.log("检查非表格部分是否需要后处理分句...")
                            paragraph_data = post_process_english_split(
                                paragraph_data, model_id, source_lang, target_lang, enable_ai_split=True
                            )

                        log_manager.log(f"   非表格处理结果: {len(paragraph_data)} 行")
                        all_data.extend(paragraph_data)

            if not all_data:
                log_manager.log_exception("处理结果为空，文件处理失败")
                return False

            df = pd.DataFrame(all_data)

            # 质量检查
            log_manager.log("执行质量检查...")
            issues = AlignmentChecker.full_check(df, source_lang, target_lang)

            if issues:
                log_manager.log_exception(f"发现 {len(issues)} 个潜在问题（仅警告）")
                for issue in issues[:10]:
                    orig_text = issue.get('original_text', '')
                    orig_preview = orig_text[:80] if orig_text else ''
                    log_manager.log_exception(
                        f"行 {issue.get('row', '?')}: {issue.get('type', '')} - {issue.get('detail', '')}",
                        f"原文: {orig_preview}..."
                    )
                issue_path = output_excel_path.replace('.xlsx', '_问题报告.xlsx')
                save_issues_report(issues, issue_path)

            df.to_excel(output_excel_path, index=False)
            log_manager.log(f"✅ 已保存: {output_excel_path}（{len(df)} 行，含表格 + 段落）")
            return True

    # 原有逻辑：无表格或非Word文档的处理
    text_original = read_file_content(file_original_path)
    text_trans = read_file_content(file_trans_path)

    if not text_original or not text_trans:
        log_manager.log_exception("内容为空，跳过文件", f"原文: {file_original_path}")
        return False

    # 根据语言生成动态提示词
    if system_prompt_override:
        system_prompt = system_prompt_override
    else:
        system_prompt = get_docx_alignment_prompt(source_lang, target_lang)

    anchor_hint = ""
    if anchor_info_orig and anchor_info_trans:
        anchor_hint = f"""
## 上下文提示
这是文档的一个片段：
- 原文开头: "{anchor_info_orig.get('first_anchor', '')[:100]}"
- 译文开头: "{anchor_info_trans.get('first_anchor', '')[:100]}"
"""

    user_prompt = f"""{anchor_hint}

<Stream_A_Original>
{text_original}
</Stream_A_Original>

<Stream_B_Translation>
{text_trans}
</Stream_B_Translation>

请严格按规则输出对齐结果：
"""

    response = call_llm_stream(system_prompt, user_prompt, model_id, filename)

    # 检测截断
    if response:
        last_500 = response[-500:] if len(response) > 500 else response
        if '|||' not in last_500:
            log_manager.log_exception("⚠️ 输出可能被截断（最后500字符无分隔符）")

    data = parse_alignment_response(response)

    if not data:
        log_manager.log_exception("解析结果为空，文件处理失败")
        return False

    # 对非CJK源语言进行后处理分句（如英文原文需要按句号细分）
    if enable_post_split and source_lang not in ["中文", "日语", "韩语"]:
        log_manager.log("检查是否需要后处理分句...")
        data = post_process_english_split(data, model_id, source_lang, target_lang, enable_ai_split=True)

    df = pd.DataFrame(data)

    # 质量检查（传递语言参数）
    log_manager.log("执行质量检查...")
    issues = AlignmentChecker.full_check(df, source_lang, target_lang)

    if issues:
        log_manager.log_exception(f"发现 {len(issues)} 个潜在问题（仅警告）")
        for issue in issues[:10]:
            orig_text = issue.get('original_text', '')
            orig_preview = orig_text[:80] if orig_text else ''
            log_manager.log_exception(
                f"行 {issue.get('row', '?')}: {issue.get('type', '')} - {issue.get('detail', '')}",
                f"原文: {orig_preview}..."
            )
        issue_path = output_excel_path.replace('.xlsx', '_问题报告.xlsx')
        save_issues_report(issues, issue_path)

    df.to_excel(output_excel_path, index=False)
    log_manager.log(f"✅ 已保存: {output_excel_path}（{len(df)} 行）")
    return True


def needs_sentence_split(orig_text, trans_text, source_lang="中文"):
    """判断是否需要分句 - 根据原文语言检测相应的句末标点"""
    source_is_cjk = source_lang in ["中文", "日语", "韩语"]

    if source_is_cjk:
        # 中日韩语言使用全角句末标点
        terminators = ['。', '！', '？']
    else:
        # 西方语言使用半角句末标点
        terminators = ['. ', '! ', '? ', '.', '!', '?']

    count = sum(orig_text.count(t) for t in terminators)
    return count > 1


def split_row_with_ai(orig_text, trans_text, model_id, row_num, source_lang="中文"):
    """使用 AI 对单行进行分句对齐"""
    log_manager.log_stream(f"\n{'=' * 50}\n")
    log_manager.log_stream(f"📍 第 {row_num} 行\n")
    log_manager.log_stream(f"[原文] {orig_text[:100]}{'...' if len(orig_text) > 100 else ''}\n")
    log_manager.log_stream(f"[译文] {trans_text[:100]}{'...' if len(trans_text) > 100 else ''}\n")
    log_manager.log_stream(f"[AI输出] ")

    user_prompt = f"""请根据原文的标点分句，将译文对应拆分对齐：

【原文】
{orig_text}

【译文】
{trans_text}

输出分句对齐结果（以原文分句数量为准）："""

    # 根据原文语言获取对应的分句提示词
    system_prompt = get_split_row_prompt(source_lang)

    try:
        # 使用统一的 LLM 调用函数，自动检测提供商
        full_response = call_llm_stream(
            system_prompt,
            user_prompt,
            model_id,
            f"分句-第{row_num}行"
        )

        if not full_response:
            log_manager.log_exception(f"第 {row_num} 行 AI 返回为空")
            return None

        results = parse_alignment_response(full_response)

        if results:
            log_manager.log_stream(f"✅ 分句结果: 1 → {len(results)} 句\n")
            log_manager.log(f"第 {row_num} 行: 1 → {len(results)} 句")
            return results
        else:
            log_manager.log_exception(f"第 {row_num} 行 AI 返回解析失败", full_response[:200])
            return None

    except Exception as e:
        log_manager.log_exception(f"第 {row_num} 行 AI 处理失败", str(e))
        return None


def needs_table_cell_split(text, source_lang="中文"):
    """
    判断表格单元格内容是否需要分句 - 检测句末标点、换行、序号等

    触发LLM分句的条件（满足任一即可）：
    1. 有多个句末标点（。！？）
    2. 有"！"或"？"
    3. 有换行符（\n、\r、或多个连续空格模拟的换行）
    4. 有序号模式（如：1. 2. 或 （一）（二）或 ①② 等）
    """
    if not text or not text.strip():
        return False

    import re

    source_is_cjk = source_lang in ["中文", "日语", "韩语"]

    # 检测换行符（包括 \n、\r\n、\r）
    has_newline = '\n' in text or '\r' in text

    # 检测连续空格（2个或以上连续空格，可能是段落分隔）
    has_consecutive_spaces = '  ' in text

    # 检测句末标点
    if source_is_cjk:
        # 中文标点
        period_count = text.count('。')
        exclamation_count = text.count('！')
        question_count = text.count('？')
    else:
        # 英文标点
        period_count = text.count('.')
        exclamation_count = text.count('!')
        question_count = text.count('?')

    total_punct_count = period_count + exclamation_count + question_count

    # 检测序号模式
    # 中文序号：（一）（二）、（1）（2）、一、二、等
    # 数字序号：1. 2. 或 1、2、或 1) 2)
    # 圆圈序号：①②③
    has_numbered_list = bool(re.search(
        r'（[一二三四五六七八九十\d]+）|'  # （一）（二）（1）（2）
        r'[（\(]\d+[）\)]|'  # (1) (2)
        r'[①②③④⑤⑥⑦⑧⑨⑩]|'  # ①②③
        r'\d+\.(?!\d)\s*\S|'  # 1. 后面跟非数字（避免匹配小数如1.0、2.5）
        r'\d+、\s*\S|'  # 1、2、后面跟内容
        r'\d+[）\)]\s*\S|'  # 1) 2）后面跟内容
        r'^[一二三四五六七八九十]+[、．.]',  # 一、二、
        text
    ))

    # 触发条件（满足任一即可）：
    # 1. 有多个句末标点（需要分句）
    # 2. 有感叹号或问号
    # 3. 有换行符（段落分隔）
    # 4. 有连续空格（可能是段落分隔）
    # 5. 有序号列表模式
    return (
            total_punct_count > 1 or  # 多个句末标点
            exclamation_count > 0 or  # 有感叹号
            question_count > 0 or  # 有问号
            has_newline or  # 有换行
            has_consecutive_spaces or  # 有连续空格
            has_numbered_list  # 有序号列表
    )


def split_table_cell_with_ai(orig_text, trans_text, model_id, cell_ref, source_lang="中文"):
    """使用 AI 对表格单元格内容进行细粒度分句对齐"""
    log_manager.log_stream(f"\n{'=' * 50}\n")
    log_manager.log_stream(f"📍 单元格 {cell_ref}\n")
    log_manager.log_stream(f"[原文] {orig_text[:100]}{'...' if len(orig_text) > 100 else ''}\n")
    log_manager.log_stream(f"[译文] {trans_text[:100]}{'...' if len(trans_text) > 100 else ''}\n")
    log_manager.log_stream(f"[AI分句] ")

    user_prompt = f"""请根据原文的句末标点（。！？）和换行符进行分句，将译文对应拆分对齐：

【原文】
{orig_text}

【译文】
{trans_text}

输出分句对齐结果："""

    # 获取表格单元格分句提示词
    system_prompt = get_table_cell_split_prompt(source_lang)

    try:
        # 使用统一的 LLM 调用函数
        full_response = call_llm_stream(
            system_prompt,
            user_prompt,
            model_id,
            f"表格分句-{cell_ref}"
        )

        if not full_response:
            log_manager.log_exception(f"单元格 {cell_ref} AI 返回为空")
            return None

        results = parse_alignment_response(full_response)

        if results:
            log_manager.log_stream(f"✅ 分句结果: 1 → {len(results)} 句\n")
            log_manager.log(f"单元格 {cell_ref}: 1 → {len(results)} 句")
            return results
        else:
            log_manager.log_exception(f"单元格 {cell_ref} AI 返回解析失败", full_response[:200])
            return None

    except Exception as e:
        log_manager.log_exception(f"单元格 {cell_ref} AI 处理失败", str(e))
        return None


def extract_docx_tables_with_position(doc_path):
    """
    从 Word 文档中提取所有表格，保留单元格位置信息
    返回: [(table_idx, row_idx, col_idx, cell_text), ...]
    """
    try:
        doc = Document(doc_path)
        table_cells = []

        for table_idx, table in enumerate(doc.tables):
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    if cell_text:
                        table_cells.append({
                            'table_idx': table_idx,
                            'row_idx': row_idx,
                            'col_idx': col_idx,
                            'text': cell_text,
                            'cell_ref': f"Word表{table_idx + 1}-行{row_idx + 1}列{col_idx + 1}"
                        })

        return table_cells
    except Exception as e:
        log_manager.log_exception(f"提取表格失败", str(e))
        return []


def process_docx_tables_alignment(orig_docx_path, trans_docx_path, model_id,
                                  source_lang="中文", target_lang="英语"):
    """
    处理 Word 文档中的表格 - 按单元格绝对位置匹配后进行细粒度分句

    处理流程：
    1. 提取原文和译文文档中的所有表格单元格（保留位置信息）
    2. 按照 (表格索引, 行索引, 列索引) 进行位置匹配
    3. 对匹配后的单元格内容调用 LLM 进行细粒度分句
    4. 分句规则：句末为"。","？","！"以及换行时断句

    返回: [(原文, 译文), ...] 格式的对齐结果列表
    """
    log_manager.log("=" * 60)
    log_manager.log("📊 Word 表格处理（按单元格位置匹配 + LLM细粒度分句）")
    log_manager.log("=" * 60)
    log_manager.log(f"语言对: {source_lang} → {target_lang}")

    # 提取原文和译文的表格单元格
    orig_cells = extract_docx_tables_with_position(orig_docx_path)
    trans_cells = extract_docx_tables_with_position(trans_docx_path)

    log_manager.log(f"原文表格单元格数: {len(orig_cells)}")
    log_manager.log(f"译文表格单元格数: {len(trans_cells)}")

    if not orig_cells:
        log_manager.log("原文文档中没有表格内容")
        return []

    # 建立译文单元格的位置索引
    trans_cell_map = {}
    for cell in trans_cells:
        key = (cell['table_idx'], cell['row_idx'], cell['col_idx'])
        trans_cell_map[key] = cell

    all_results = []
    total_cells_processed = 0
    total_cells_split = 0

    log_manager.log_stream("\n" + "=" * 60 + "\n")
    log_manager.log_stream(f"📄 开始处理 Word 文档表格（非Excel）\n")
    log_manager.log_stream(f"📝 共 {len(orig_cells)} 个 Word 表格单元格待处理\n")
    log_manager.log_stream("=" * 60 + "\n")

    for orig_cell in orig_cells:
        key = (orig_cell['table_idx'], orig_cell['row_idx'], orig_cell['col_idx'])
        cell_ref = orig_cell['cell_ref']
        orig_text = orig_cell['text']

        # 查找对应位置的译文单元格
        trans_cell = trans_cell_map.get(key)
        trans_text = trans_cell['text'] if trans_cell else ""

        total_cells_processed += 1

        # 跳过原文和译文都为空的情况
        if not orig_text and not trans_text:
            continue

        # 判断是否需要细粒度分句
        if orig_text and trans_text and needs_table_cell_split(orig_text, source_lang):
            log_manager.log(f"  {cell_ref}: 需要细粒度分句")
            split_results = split_table_cell_with_ai(
                orig_text, trans_text, model_id, cell_ref, source_lang
            )

            if split_results and len(split_results) > 1:
                # 分句成功，添加所有分句结果
                for result in split_results:
                    result['来源'] = cell_ref
                all_results.extend(split_results)
                total_cells_split += 1
            else:
                # 分句失败或只有一句，保留原内容
                all_results.append({
                    "原文": orig_text,
                    "译文": trans_text,
                    "来源": cell_ref
                })
        else:
            # 不需要分句，直接添加
            if orig_text or trans_text:
                all_results.append({
                    "原文": orig_text,
                    "译文": trans_text,
                    "来源": cell_ref
                })
                if orig_text and trans_text:
                    log_manager.log_stream(f"[{cell_ref}] 直接配对\n")

    log_manager.log_stream("\n" + "=" * 60 + "\n")
    log_manager.log_stream(f"✅ Word 文档表格处理完成！\n")
    log_manager.log_stream(f"   处理 Word 单元格数: {total_cells_processed}\n")
    log_manager.log_stream(f"   AI分句处理数: {total_cells_split}\n")
    log_manager.log_stream(f"   最终输出行数: {len(all_results)}\n")
    log_manager.log_stream("=" * 60 + "\n")

    log_manager.log(f"✅ 表格处理完成: {total_cells_processed} 个单元格 → {len(all_results)} 行")

    return all_results


def read_docx_without_tables(file_path):
    """读取 Word 文档内容，排除表格部分（表格单独处理）"""
    try:
        doc = Document(file_path)
        full_text = []

        if hasattr(doc, 'element') and hasattr(doc.element, 'body'):
            for child in doc.element.body.iterchildren():
                if child.tag.endswith('p'):
                    # 段落元素
                    para = Paragraph(child, doc)
                    if para.text.strip():
                        full_text.append(para.text)
                # 跳过表格元素 (tbl)，表格单独处理

        # 页眉页脚 (包括段落、表格、文本框)
        for section in doc.sections:
            # 页眉段落
            for p in section.header.paragraphs:
                if p.text.strip():
                    full_text.append(p.text)
            # 页眉表格
            for table in section.header.tables:
                seen_header_cells = set()
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen_header_cells:
                            full_text.append(cell_text)
                            seen_header_cells.add(cell_text)
            # 页眉文本框
            if section.header._element is not None:
                header_xml = etree.tostring(section.header._element, encoding='unicode')
                header_root = etree.fromstring(header_xml.encode('utf-8'))
                nsmap_header = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                for txbx in header_root.xpath('.//w:txbxContent', namespaces=nsmap_header):
                    txbx_text = ''.join(txbx.xpath('.//w:t/text()', namespaces=nsmap_header))
                    if txbx_text.strip():
                        full_text.append(txbx_text.strip())

            # 页脚段落
            for p in section.footer.paragraphs:
                if p.text.strip():
                    full_text.append(p.text)
            # 页脚表格
            for table in section.footer.tables:
                seen_footer_cells = set()
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen_footer_cells:
                            full_text.append(cell_text)
                            seen_footer_cells.add(cell_text)
            # 页脚文本框
            if section.footer._element is not None:
                footer_xml = etree.tostring(section.footer._element, encoding='unicode')
                footer_root = etree.fromstring(footer_xml.encode('utf-8'))
                nsmap_footer = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                for txbx in footer_root.xpath('.//w:txbxContent', namespaces=nsmap_footer):
                    txbx_text = ''.join(txbx.xpath('.//w:t/text()', namespaces=nsmap_footer))
                    if txbx_text.strip():
                        full_text.append(txbx_text.strip())

        # 文本框处理（保留原有逻辑）
        if hasattr(doc.element, 'xml'):
            xml = doc.element.xml
            root = etree.fromstring(xml.encode('utf-8'))

            nsmap = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'wps': 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape',
            }

            textbox_texts = set()

            # 传统文本框
            try:
                textbox_containers = root.xpath('.//w:txbxContent', namespaces=nsmap)
                for container in textbox_containers:
                    paragraphs = container.xpath('.//w:p', namespaces=nsmap)
                    para_texts = []
                    for p in paragraphs:
                        p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged_text = '\n'.join(para_texts)
                    if merged_text.strip():
                        textbox_texts.add(merged_text.strip())
            except:
                pass

            for text in textbox_texts:
                full_text.append(text)

        # 合并文本，过滤空行避免原文译文空行不对应导致的对齐问题
        result = "\n".join(full_text)
        # 把连续多个换行符合并成单个换行符
        import re
        result = re.sub(r'\n{2,}', '\n', result)
        return result.strip()
    except Exception as e:
        log_manager.log_exception(f"读取文档失败（排除表格）", str(e))
        return ""


def has_docx_tables(file_path):
    """检查 Word 文档是否包含表格"""
    try:
        doc = Document(file_path)
        return len(doc.tables) > 0
    except:
        return False


# ==========================================
# === 📊 Excel 双文件处理（新增功能）===
# ==========================================
def process_excel_dual_file_alignment(orig_excel_path, trans_excel_path, output_excel_path, model_id,
                                      source_lang="中文", target_lang="英语"):
    """
    处理两个对应的 Excel 表格文件
    - 支持多工作簿
    - 按单元格位置对应
    - 需要分句时调用 AI 分割
    """
    log_manager.log("=" * 60)
    log_manager.log("📊 Excel 双文件对齐模式（按单元格位置对应）")
    log_manager.log("=" * 60)
    log_manager.log(f"语言对: {source_lang} → {target_lang}")
    log_manager.log(f"原文文件: {orig_excel_path}")
    log_manager.log(f"译文文件: {trans_excel_path}")

    # 读取所有工作簿
    orig_sheets = read_excel_all_sheets(orig_excel_path)
    trans_sheets = read_excel_all_sheets(trans_excel_path)

    if orig_sheets is None or trans_sheets is None:
        log_manager.log_exception("无法读取 Excel 文件")
        return False

    # 获取共同的工作簿名称
    orig_sheet_names = set(orig_sheets.keys())
    trans_sheet_names = set(trans_sheets.keys())
    common_sheets = orig_sheet_names & trans_sheet_names

    if not common_sheets:
        # 如果工作簿名称不匹配，尝试按顺序对应
        log_manager.log("工作簿名称不匹配，尝试按顺序对应...")
        orig_list = list(orig_sheets.keys())
        trans_list = list(trans_sheets.keys())
        sheet_pairs = list(zip(orig_list, trans_list))
        log_manager.log(f"按顺序配对: {sheet_pairs}")
    else:
        sheet_pairs = [(name, name) for name in sorted(common_sheets)]
        log_manager.log(f"共同工作簿: {list(common_sheets)}")

        # 检查是否有不匹配的工作簿
        only_in_orig = orig_sheet_names - trans_sheet_names
        only_in_trans = trans_sheet_names - orig_sheet_names
        if only_in_orig:
            log_manager.log_exception(f"以下工作簿仅存在于原文文件中: {only_in_orig}")
        if only_in_trans:
            log_manager.log_exception(f"以下工作簿仅存在于译文文件中: {only_in_trans}")

    all_results = []
    total_cells_processed = 0
    total_rows_split = 0

    log_manager.log_stream("\n" + "=" * 60 + "\n")
    log_manager.log_stream(f"📊 开始处理 Excel 双文件对齐\n")
    log_manager.log_stream(f"📝 共 {len(sheet_pairs)} 个工作簿待处理\n")
    log_manager.log_stream("=" * 60 + "\n")

    for sheet_idx, (orig_sheet_name, trans_sheet_name) in enumerate(sheet_pairs):
        log_manager.log(
            f"\n处理工作簿 {sheet_idx + 1}/{len(sheet_pairs)}: '{orig_sheet_name}' <-> '{trans_sheet_name}'")
        log_manager.log_stream(f"\n{'─' * 40}\n")
        log_manager.log_stream(f"📑 工作簿: {orig_sheet_name}\n")
        log_manager.log_stream(f"{'─' * 40}\n")

        df_orig = orig_sheets[orig_sheet_name]
        df_trans = trans_sheets[trans_sheet_name]

        # 获取两个表格的最大行列数
        max_rows = max(df_orig.shape[0], df_trans.shape[0])
        max_cols = max(df_orig.shape[1], df_trans.shape[1])

        log_manager.log(f"  原文表格: {df_orig.shape[0]} 行 x {df_orig.shape[1]} 列")
        log_manager.log(f"  译文表格: {df_trans.shape[0]} 行 x {df_trans.shape[1]} 列")
        log_manager.log(f"  处理范围: {max_rows} 行 x {max_cols} 列")

        # 遍历所有单元格位置
        cell_num = 0
        for row_idx in range(max_rows):
            for col_idx in range(max_cols):
                # 获取原文单元格内容
                orig_text = ""
                if row_idx < df_orig.shape[0] and col_idx < df_orig.shape[1]:
                    cell_value = df_orig.iloc[row_idx, col_idx]
                    if pd.notna(cell_value):
                        orig_text = str(cell_value).strip()

                # 获取译文单元格内容
                trans_text = ""
                if row_idx < df_trans.shape[0] and col_idx < df_trans.shape[1]:
                    cell_value = df_trans.iloc[row_idx, col_idx]
                    if pd.notna(cell_value):
                        trans_text = str(cell_value).strip()

                # 跳过两边都为空的单元格
                if not orig_text and not trans_text:
                    continue

                cell_num += 1
                total_cells_processed += 1
                cell_ref = f"[{orig_sheet_name}] R{row_idx + 1}C{col_idx + 1}"

                # 判断是否需要分句（根据原文语言检测标点、换行、序号等）
                if orig_text and trans_text and needs_table_cell_split(orig_text, source_lang):
                    log_manager.log(f"  {cell_ref}: 需要细粒度分句")
                    split_results = split_table_cell_with_ai(
                        orig_text, trans_text, model_id, cell_ref, source_lang
                    )

                    if split_results and len(split_results) > 1:
                        # 分句成功，添加所有分句结果
                        for result in split_results:
                            result['来源'] = cell_ref
                        all_results.extend(split_results)
                        total_rows_split += 1
                    else:
                        # 分句失败或只有一句，保留原内容
                        all_results.append({
                            "原文": orig_text,
                            "译文": trans_text,
                            "来源": cell_ref
                        })
                else:
                    # 不需要分句，直接添加
                    if orig_text or trans_text:
                        all_results.append({
                            "原文": orig_text,
                            "译文": trans_text,
                            "来源": cell_ref
                        })
                        if orig_text and trans_text:
                            log_manager.log_stream(f"[{cell_ref}] 直接配对\n")

        log_manager.log(f"  工作簿 '{orig_sheet_name}' 处理完成，有效单元格: {cell_num}")

    if not all_results:
        log_manager.log_exception("处理结果为空")
        return False

    # 创建结果 DataFrame
    result_df = pd.DataFrame(all_results)

    # 可选：移除"来源"列（如果不需要显示）
    # 如果需要保留来源信息用于调试，可以注释掉下面这行
    if '来源' in result_df.columns:
        result_df_output = result_df[['原文', '译文']].copy()
    else:
        result_df_output = result_df

    # 质量检查（传递语言参数）
    log_manager.log("执行质量检查...")
    issues = AlignmentChecker.full_check(result_df_output, source_lang, target_lang)
    if issues:
        log_manager.log_exception(f"发现 {len(issues)} 个潜在问题")
        issue_path = output_excel_path.replace('.xlsx', '_问题报告.xlsx')
        save_issues_report(issues, issue_path)

    # 保存结果
    result_df_output.to_excel(output_excel_path, index=False)

    # 输出统计
    log_manager.log_stream("\n" + "=" * 60 + "\n")
    log_manager.log_stream(f"✅ Excel 双文件对齐完成！\n")
    log_manager.log_stream(f"   处理工作簿数: {len(sheet_pairs)}\n")
    log_manager.log_stream(f"   有效单元格数: {total_cells_processed}\n")
    log_manager.log_stream(f"   AI分句处理数: {total_rows_split}\n")
    log_manager.log_stream(f"   最终输出行数: {len(result_df_output)}\n")
    log_manager.log_stream("=" * 60 + "\n")

    log_manager.log(f"✅ 处理完成: {total_cells_processed} 个单元格 → {len(result_df_output)} 行")
    log_manager.log(f"📁 已保存: {output_excel_path}")

    return True


def process_excel_alignment(original_excel_path, trans_excel_path, output_excel_path, model_id,
                            source_lang="中文", target_lang="英语"):
    """
    处理 Excel 双文件对齐（兼容旧接口）
    自动检测并调用新的双文件处理函数
    """
    # 使用新的双文件处理函数
    return process_excel_dual_file_alignment(original_excel_path, trans_excel_path,
                                             output_excel_path, model_id, source_lang, target_lang)


def merge_and_deduplicate_excels(excel_paths, final_output_path, source_lang="中文", target_lang="英语"):
    """合并Excel文件，高亮单列重复行，去除完全相同的行"""
    from openpyxl import load_workbook
    from openpyxl.styles import PatternFill

    log_manager.log(f"合并 {len(excel_paths)} 个文件...")
    log_manager.log(f"待合并文件列表:")
    for i, path in enumerate(excel_paths):
        exists = "✓ 存在" if os.path.exists(path) else "✗ 不存在"
        log_manager.log(f"  {i + 1}. {os.path.basename(path)} [{exists}]")

    dfs = []
    for path in excel_paths:
        if os.path.exists(path):
            try:
                df = pd.read_excel(path)
                if not df.empty:
                    dfs.append(df)
                    log_manager.log(f"  ✅ 读取成功: {os.path.basename(path)} ({len(df)} 行)")
                else:
                    log_manager.log_exception(f"文件为空", os.path.basename(path))
            except Exception as e:
                log_manager.log_exception(f"读取失败: {e}", path)
        else:
            log_manager.log_exception(f"文件不存在", path)

    if not dfs:
        log_manager.log_exception("没有可合并的数据")
        return None

    combined_df = pd.concat(dfs, ignore_index=True)
    total_before = len(combined_df)
    log_manager.log(f"合并后总行数: {total_before}")

    combined_df.drop_duplicates(subset=['原文', '译文'], keep='first', inplace=True)
    total_after_full_dedup = len(combined_df)
    full_dup_removed = total_before - total_after_full_dedup
    log_manager.log(f"去除完全重复: {total_before} -> {total_after_full_dedup} 行 (移除 {full_dup_removed} 行)")

    combined_df = combined_df.reset_index(drop=True)

    orig_duplicated = combined_df.duplicated(subset=['原文'], keep=False)
    trans_duplicated = combined_df.duplicated(subset=['译文'], keep=False)

    highlight_mask = orig_duplicated | trans_duplicated
    highlight_rows = combined_df[highlight_mask].index.tolist()

    orig_dup_count = orig_duplicated.sum()
    trans_dup_count = trans_duplicated.sum()
    log_manager.log(f"单列重复检测: 原文重复 {orig_dup_count} 行, 译文重复 {trans_dup_count} 行")

    issues = AlignmentChecker.full_check(combined_df, source_lang, target_lang)
    if issues:
        log_manager.log_exception(f"发现 {len(issues)} 个潜在问题")
        issue_path = final_output_path.replace('.xlsx', '_问题报告.xlsx')
        save_issues_report(issues, issue_path)

    combined_df.to_excel(final_output_path, index=False)

    try:
        wb = load_workbook(final_output_path)
        ws = wb.active

        yellow_fill = PatternFill(start_color='FFFF00', end_color='FFFF00', fill_type='solid')

        for idx in highlight_rows:
            excel_row = idx + 2
            for col in range(1, 3):
                ws.cell(row=excel_row, column=col).fill = yellow_fill

        wb.save(final_output_path)
        log_manager.log(f"高亮标记: {len(highlight_rows)} 行（缓冲区重叠导致的单列重复）")

    except Exception as e:
        log_manager.log_exception(f"高亮处理失败", str(e))

    log_manager.log(f"✅ 最终结果: {final_output_path} ({len(combined_df)} 行)")
    return final_output_path


# ==========================================
# === 🖥️ GUI 界面 ===
# ==========================================
class DocumentAlignerGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("📄 多语对照记忆工具 v4.3")
        self.root.geometry("1400x950")
        self.root.minsize(1100, 750)

        self.supported_filetypes = [
            ("支持的文件", "*.docx *.doc *.pptx *.xlsx *.xls"),
            ("Word文档", "*.docx *.doc"),
            ("PowerPoint", "*.pptx"),
            ("Excel表格", "*.xlsx *.xls"),
            ("所有文件", "*.*")
        ]
        self.excel_filetypes = [("Excel表格", "*.xlsx *.xls"), ("所有文件", "*.*")]

        self.running = True
        self.processing_stopped = False

        self.setup_styles()
        self.create_widgets()
        self.update_logs()

        self.root.protocol("WM_DELETE_WINDOW", self.on_closing)

    def setup_styles(self):
        style = ttk.Style()
        style.configure("Title.TLabel", font=("Microsoft YaHei", 14, "bold"))
        style.configure("Section.TLabelframe.Label", font=("Microsoft YaHei", 10, "bold"))
        style.configure("Info.TLabel", font=("Microsoft YaHei", 9))

    def create_widgets(self):
        main_container = ttk.Frame(self.root, padding="10")
        main_container.pack(fill=tk.BOTH, expand=True)

        # 语言选择
        lang_frame = ttk.LabelFrame(main_container, text="🌍 语言设置", padding="10")
        lang_frame.pack(fill=tk.X, pady=(0, 10))

        lang_options = list(SUPPORTED_LANGUAGES.keys())

        ttk.Label(lang_frame, text="原文语言:").grid(row=0, column=0, sticky=tk.W, padx=5, pady=5)
        self.source_lang_var = tk.StringVar(value=DEFAULT_SOURCE_LANG)
        self.source_lang_combo = ttk.Combobox(lang_frame, textvariable=self.source_lang_var,
                                              values=lang_options, width=15, state="readonly")
        self.source_lang_combo.grid(row=0, column=1, padx=5, pady=5, sticky=tk.W)

        ttk.Label(lang_frame, text="译文语言:").grid(row=0, column=2, sticky=tk.W, padx=(20, 5), pady=5)
        self.target_lang_var = tk.StringVar(value=DEFAULT_TARGET_LANG)
        self.target_lang_combo = ttk.Combobox(lang_frame, textvariable=self.target_lang_var,
                                              values=lang_options, width=15, state="readonly")
        self.target_lang_combo.grid(row=0, column=3, padx=5, pady=5, sticky=tk.W)

        # 语言说明
        self.lang_info_var = tk.StringVar(value="")
        self.lang_info_label = ttk.Label(lang_frame, textvariable=self.lang_info_var, foreground="#666666")
        self.lang_info_label.grid(row=0, column=4, padx=20, pady=5, sticky=tk.W)

        # 后处理分句选项（针对英文等西方语言）
        self.post_split_var = tk.BooleanVar(value=True)
        self.post_split_check = ttk.Checkbutton(
            lang_frame,
            text="启用后处理细粒度分句（英文等西方语言按句号细分，自动排除缩写）",
            variable=self.post_split_var
        )
        self.post_split_check.grid(row=1, column=0, columnspan=5, sticky=tk.W, padx=5, pady=(5, 0))

        # 后处理分句说明
        post_split_info = ttk.Label(
            lang_frame,
            text="💡 当源语言为英文等西方语言时，会在LLM对齐后进一步按句号细分",
            foreground="#0066cc"
        )
        post_split_info.grid(row=2, column=0, columnspan=5, sticky=tk.W, padx=5, pady=(2, 0))

        self.source_lang_combo.bind("<<ComboboxSelected>>", self.on_lang_changed)
        self.target_lang_combo.bind("<<ComboboxSelected>>", self.on_lang_changed)
        self.on_lang_changed()  # 初始化显示

        # 处理模式选择
        mode_frame = ttk.LabelFrame(main_container, text="📋 处理模式", padding="10")
        mode_frame.pack(fill=tk.X, pady=(0, 10))

        self.mode_var = tk.StringVar(value="dual_file")
        ttk.Radiobutton(mode_frame, text="常规模式（Word/PPT/Excel：分别选择原文和译文文件）",
                        variable=self.mode_var, value="dual_file",
                        command=self.on_mode_changed).pack(anchor=tk.W, pady=2)

        # 模式说明
        self.mode_info_label = ttk.Label(mode_frame,
                                         text="💡 双文件模式支持 Excel：两个表格按单元格位置对应，支持多工作簿",
                                         foreground="#0066cc")
        self.mode_info_label.pack(anchor=tk.W, pady=(5, 0))

        # 文件选择
        self.file_frame = ttk.LabelFrame(main_container, text="📁 文件选择", padding="10")
        self.file_frame.pack(fill=tk.X, pady=(0, 10))

        self.original_label = ttk.Label(self.file_frame, text="原文文件 (中文):")
        self.original_label.grid(row=0, column=0, sticky=tk.W, pady=5)
        self.original_path_var = tk.StringVar()
        self.original_entry = ttk.Entry(self.file_frame, textvariable=self.original_path_var, width=80)
        self.original_entry.grid(row=0, column=1, padx=5, pady=5, sticky=tk.EW)
        self.original_button = ttk.Button(self.file_frame, text="浏览...", command=self.browse_original)
        self.original_button.grid(row=0, column=2, padx=5)

        self.trans_label = ttk.Label(self.file_frame, text="译文文件 (英文):")
        self.trans_label.grid(row=1, column=0, sticky=tk.W, pady=5)
        self.trans_path_var = tk.StringVar()
        self.trans_entry = ttk.Entry(self.file_frame, textvariable=self.trans_path_var, width=80)
        self.trans_entry.grid(row=1, column=1, padx=5, pady=5, sticky=tk.EW)
        self.trans_button = ttk.Button(self.file_frame, text="浏览...", command=self.browse_trans)
        self.trans_button.grid(row=1, column=2, padx=5)

        ttk.Label(self.file_frame, text="输出目录:").grid(row=2, column=0, sticky=tk.W, pady=5)
        self.output_path_var = tk.StringVar(value=os.path.abspath(OUTPUT_DIR))
        ttk.Entry(self.file_frame, textvariable=self.output_path_var, width=80).grid(row=2, column=1, padx=5, pady=5,
                                                                                     sticky=tk.EW)
        ttk.Button(self.file_frame, text="浏览...", command=self.browse_output).grid(row=2, column=2, padx=5)

        self.file_frame.columnconfigure(1, weight=1)

        # 更新文件标签（在标签创建后）
        self.update_file_labels()

        # 模型选择
        model_frame = ttk.LabelFrame(main_container, text="🤖 模型选择", padding="10")
        model_frame.pack(fill=tk.X, pady=(0, 10))

        # 提供商切换
        provider_frame = ttk.Frame(model_frame)
        provider_frame.grid(row=0, column=0, columnspan=4, sticky=tk.W, pady=(0, 10))

        ttk.Label(provider_frame, text="API 提供商:", font=("Microsoft YaHei", 9, "bold")).pack(side=tk.LEFT,
                                                                                                padx=(0, 10))

        self.provider_var = tk.StringVar(value=DEFAULT_PROVIDER)

        self.openrouter_radio = ttk.Radiobutton(provider_frame, text="🌐 OpenRouter",
                                                variable=self.provider_var, value="openrouter",
                                                command=self.on_provider_changed)
        self.openrouter_radio.pack(side=tk.LEFT, padx=5)

        self.provider_info_label = ttk.Label(provider_frame, text="", foreground="#0066cc")
        self.provider_info_label.pack(side=tk.LEFT, padx=20)

        # 模型选择
        ttk.Label(model_frame, text="选择模型:").grid(row=1, column=0, sticky=tk.W, pady=5)

        self.model_var = tk.StringVar(value=DEFAULT_MODEL)
        self.model_combo = ttk.Combobox(model_frame, textvariable=self.model_var,
                                        values=list(AVAILABLE_MODELS.keys()), width=30, state="readonly")
        self.model_combo.grid(row=1, column=1, padx=5, pady=5, sticky=tk.W)
        self.model_combo.bind("<<ComboboxSelected>>", self.on_model_selected)

        self.model_info_var = tk.StringVar()
        ttk.Label(model_frame, textvariable=self.model_info_var, foreground="#666666").grid(row=1, column=2, padx=20,
                                                                                            pady=5, sticky=tk.W)

        # 初始化提供商信息
        detail_frame = ttk.Frame(model_frame)
        detail_frame.grid(row=2, column=0, columnspan=3, sticky=tk.EW, pady=(10, 0))

        ttk.Label(detail_frame, text="模型ID:", font=("Microsoft YaHei", 9, "bold")).grid(row=0, column=0, sticky=tk.W)
        self.model_id_var = tk.StringVar()
        ttk.Label(detail_frame, textvariable=self.model_id_var, foreground="#0066cc").grid(row=0, column=1, sticky=tk.W,
                                                                                           padx=10)

        ttk.Label(detail_frame, text="最大输出:", font=("Microsoft YaHei", 9, "bold")).grid(row=0, column=2,
                                                                                            sticky=tk.W, padx=(30, 0))
        self.max_output_var = tk.StringVar()
        ttk.Label(detail_frame, textvariable=self.max_output_var, foreground="#009900").grid(row=0, column=3,
                                                                                             sticky=tk.W, padx=10)

        # 初始化提供商和模型信息（在所有变量创建之后）
        self.on_provider_changed()

        # 配置选项
        config_frame = ttk.LabelFrame(main_container, text="⚙️ 处理选项（仅对 DOCX 生效）", padding="10")
        config_frame.pack(fill=tk.X, pady=(0, 10))

        ttk.Label(config_frame, text="分割阈值(字数):").grid(row=0, column=0, sticky=tk.W, padx=5)

        threshold_frame = ttk.Frame(config_frame)
        threshold_frame.grid(row=0, column=1, columnspan=3, sticky=tk.W)

        ttk.Label(threshold_frame, text="2份:").pack(side=tk.LEFT)
        self.threshold_2_var = tk.StringVar(value=str(THRESHOLD_2_PARTS))
        ttk.Entry(threshold_frame, textvariable=self.threshold_2_var, width=8).pack(side=tk.LEFT, padx=(2, 15))

        ttk.Label(threshold_frame, text="3份:").pack(side=tk.LEFT)
        self.threshold_3_var = tk.StringVar(value=str(THRESHOLD_3_PARTS))
        ttk.Entry(threshold_frame, textvariable=self.threshold_3_var, width=8).pack(side=tk.LEFT, padx=(2, 15))

        ttk.Label(threshold_frame, text="4份:").pack(side=tk.LEFT)
        self.threshold_4_var = tk.StringVar(value=str(THRESHOLD_4_PARTS))
        ttk.Entry(threshold_frame, textvariable=self.threshold_4_var, width=8).pack(side=tk.LEFT, padx=(2, 15))

        ttk.Label(threshold_frame, text="5份:").pack(side=tk.LEFT)
        self.threshold_5_var = tk.StringVar(value=str(THRESHOLD_5_PARTS))
        ttk.Entry(threshold_frame, textvariable=self.threshold_5_var, width=8).pack(side=tk.LEFT, padx=(2, 15))
        ttk.Label(threshold_frame, text="6份:").pack(side=tk.LEFT)
        self.threshold_6_var = tk.StringVar(value=str(THRESHOLD_6_PARTS))
        ttk.Entry(threshold_frame, textvariable=self.threshold_6_var, width=8).pack(side=tk.LEFT, padx=(2, 15))
        ttk.Label(threshold_frame, text="7份:").pack(side=tk.LEFT)
        self.threshold_7_var = tk.StringVar(value=str(THRESHOLD_7_PARTS))
        ttk.Entry(threshold_frame, textvariable=self.threshold_7_var, width=8).pack(side=tk.LEFT, padx=(2, 15))
        ttk.Label(threshold_frame, text="8份:").pack(side=tk.LEFT)
        self.threshold_8_var = tk.StringVar(value=str(THRESHOLD_8_PARTS))
        ttk.Entry(threshold_frame, textvariable=self.threshold_8_var, width=8).pack(side=tk.LEFT, padx=(2, 15))

        ttk.Label(threshold_frame, text="缓冲区:").pack(side=tk.LEFT, padx=(20, 0))
        self.buffer_var = tk.StringVar(value=str(BUFFER_CHARS))
        ttk.Entry(threshold_frame, textvariable=self.buffer_var, width=6).pack(side=tk.LEFT, padx=(2, 0))
        ttk.Label(threshold_frame, text="字").pack(side=tk.LEFT)

        # 控制按钮
        control_frame = ttk.Frame(main_container)
        control_frame.pack(fill=tk.X, pady=10)

        self.start_button = ttk.Button(control_frame, text="🚀 开始处理", command=self.start_processing)
        self.start_button.pack(side=tk.LEFT, padx=5)

        self.stop_button = ttk.Button(control_frame, text="⏹ 停止", command=self.stop_processing, state=tk.DISABLED)
        self.stop_button.pack(side=tk.LEFT, padx=5)

        ttk.Button(control_frame, text="🗑 清空日志", command=self.clear_logs).pack(side=tk.LEFT, padx=5)
        ttk.Button(control_frame, text="📂 打开输出目录", command=self.open_output_dir).pack(side=tk.LEFT, padx=5)

        self.progress_var = tk.DoubleVar()
        ttk.Progressbar(control_frame, variable=self.progress_var, maximum=100, length=300).pack(side=tk.RIGHT, padx=5)

        self.status_label = ttk.Label(control_frame, text="就绪")
        self.status_label.pack(side=tk.RIGHT, padx=10)

        # 日志区域
        paned = ttk.PanedWindow(main_container, orient=tk.VERTICAL)
        paned.pack(fill=tk.BOTH, expand=True)

        stream_frame = ttk.LabelFrame(paned, text="🌊 模型输出流", padding="5")
        self.stream_text = scrolledtext.ScrolledText(stream_frame, height=12, wrap=tk.WORD,
                                                     font=("Consolas", 9), bg="#1e1e1e", fg="#00ff00")
        self.stream_text.pack(fill=tk.BOTH, expand=True)
        paned.add(stream_frame, weight=2)

        bottom_paned = ttk.PanedWindow(paned, orient=tk.HORIZONTAL)
        paned.add(bottom_paned, weight=1)

        log_frame = ttk.LabelFrame(bottom_paned, text="📋 运行日志", padding="5")
        self.log_text = scrolledtext.ScrolledText(log_frame, height=10, wrap=tk.WORD,
                                                  font=("Microsoft YaHei", 9), bg="#f5f5f5")
        self.log_text.pack(fill=tk.BOTH, expand=True)
        bottom_paned.add(log_frame, weight=1)

        exception_frame = ttk.LabelFrame(bottom_paned, text="⚠️ 异常/警告", padding="5")
        self.exception_text = scrolledtext.ScrolledText(exception_frame, height=10, wrap=tk.WORD,
                                                        font=("Microsoft YaHei", 9), bg="#fff0f0", fg="#cc0000")
        self.exception_text.pack(fill=tk.BOTH, expand=True)
        bottom_paned.add(exception_frame, weight=1)

    def on_lang_changed(self, event=None):
        """语言选择变化时更新显示"""
        source_lang = self.source_lang_var.get()
        target_lang = self.target_lang_var.get()

        source_info = SUPPORTED_LANGUAGES.get(source_lang, {})
        target_info = SUPPORTED_LANGUAGES.get(target_lang, {})

        source_desc = source_info.get('description', source_lang)
        target_desc = target_info.get('description', target_lang)

        self.lang_info_var.set(f"{source_desc} → {target_desc}")

        # 更新文件选择标签
        self.update_file_labels()

    def update_file_labels(self):
        """更新文件选择标签"""
        source_lang = self.source_lang_var.get()
        target_lang = self.target_lang_var.get()

        # 检查标签是否已创建
        if hasattr(self, 'original_label'):
            self.original_label.config(text=f"原文文件 ({source_lang}):")
        if hasattr(self, 'trans_label'):
            self.trans_label.config(text=f"译文文件 ({target_lang}):")

    def on_mode_changed(self):
        source_lang = self.source_lang_var.get()
        target_lang = self.target_lang_var.get()
        self.original_label.config(text=f"原文文件 ({source_lang}):")
        self.trans_label.config(text=f"译文文件 ({target_lang}):")
        self.mode_info_label.config(text="💡 双文件模式支持 Excel：两个表格按单元格位置对应，支持多工作簿")

    def on_provider_changed(self, event=None):
        """切换 API 提供商"""
        global AVAILABLE_MODELS
        provider = self.provider_var.get()

        AVAILABLE_MODELS = OPENROUTER_MODELS.copy()
        default_model = "Google Gemini 2.5 Flash"
        self.provider_info_label.config(text="HTTP 协议，Gemini 模型")

        # 更新模型下拉框
        self.model_combo['values'] = list(AVAILABLE_MODELS.keys())

        # 设置默认模型
        if default_model in AVAILABLE_MODELS:
            self.model_var.set(default_model)
        elif AVAILABLE_MODELS:
            self.model_var.set(list(AVAILABLE_MODELS.keys())[0])

        self.update_model_info()
        self.update_model_details()
        log_manager.log(f"切换到提供商: {provider}")

    def on_model_selected(self, event=None):
        self.update_model_info()
        self.update_model_details()

    def update_model_info(self):
        model_name = self.model_var.get()
        if model_name in AVAILABLE_MODELS:
            self.model_info_var.set(AVAILABLE_MODELS[model_name]['description'])

    def update_model_details(self):
        model_name = self.model_var.get()
        if model_name in AVAILABLE_MODELS:
            info = AVAILABLE_MODELS[model_name]
            self.model_id_var.set(info['id'])
            self.max_output_var.set(f"{info['max_output']:,} tokens")

    def browse_original(self):
        path = filedialog.askopenfilename(title="选择原文文件", filetypes=self.supported_filetypes)
        if path:
            self.original_path_var.set(path)

    def browse_trans(self):
        path = filedialog.askopenfilename(title="选择译文文件", filetypes=self.supported_filetypes)
        if path:
            self.trans_path_var.set(path)

    def browse_output(self):
        path = filedialog.askdirectory(title="选择输出目录")
        if path:
            self.output_path_var.set(path)

    def open_output_dir(self):
        output_dir = self.output_path_var.get()
        if os.path.exists(output_dir):
            os.startfile(output_dir)
        else:
            messagebox.showwarning("提示", "输出目录不存在")

    def clear_logs(self):
        self.log_text.delete(1.0, tk.END)
        self.exception_text.delete(1.0, tk.END)
        self.stream_text.delete(1.0, tk.END)

    def update_logs(self):
        if not self.running:
            return

        while not log_manager.log_queue.empty():
            try:
                msg = log_manager.log_queue.get_nowait()
                self.log_text.insert(tk.END, msg + "\n")
                self.log_text.see(tk.END)
            except queue.Empty:
                break

        while not log_manager.exception_queue.empty():
            try:
                msg = log_manager.exception_queue.get_nowait()
                self.exception_text.insert(tk.END, msg + "\n")
                self.exception_text.see(tk.END)
            except queue.Empty:
                break

        while not log_manager.stream_queue.empty():
            try:
                msg = log_manager.stream_queue.get_nowait()
                scrollbar_position = self.stream_text.yview()[1]
                is_at_bottom = scrollbar_position > 0.95
                self.stream_text.insert(tk.END, msg)
                if is_at_bottom:
                    self.stream_text.see(tk.END)
            except queue.Empty:
                break

        self.root.after(100, self.update_logs)

    def validate_inputs(self):
        if not API_KEY:
            messagebox.showerror("错误", "请在代码中填入 API_KEY！")
            return False

        if not self.original_path_var.get().strip():
            messagebox.showerror("错误", "请选择文件！")
            return False

        if not os.path.exists(self.original_path_var.get()):
            messagebox.showerror("错误", "文件不存在！")
            return False

        if not self.trans_path_var.get().strip():
            messagebox.showerror("错误", "请选择译文文件！")
            return False
        if not os.path.exists(self.trans_path_var.get()):
            messagebox.showerror("错误", "译文文件不存在！")
            return False

        orig_type = get_file_type(self.original_path_var.get())
        trans_type = get_file_type(self.trans_path_var.get())

        if orig_type == 'unknown':
            messagebox.showerror("错误", "原文文件类型不支持！")
            return False
        if trans_type == 'unknown':
            messagebox.showerror("错误", "译文文件类型不支持！")
            return False

        # 统一类型判断
        orig_type_normalized = 'word' if orig_type in ['doc', 'docx'] else orig_type
        trans_type_normalized = 'word' if trans_type in ['doc', 'docx'] else trans_type

        if orig_type_normalized != trans_type_normalized:
            messagebox.showerror("错误", "原文和译文文件类型必须相同！")
            return False

        return True

    def start_processing(self):
        if not self.validate_inputs():
            return

        self.start_button.config(state=tk.DISABLED)
        self.stop_button.config(state=tk.NORMAL)
        self.status_label.config(text="处理中...")
        self.progress_var.set(0)
        self.processing_stopped = False

        self.processing_thread = threading.Thread(target=self.run_processing, daemon=True)
        self.processing_thread.start()

    def stop_processing(self):
        self.processing_stopped = True
        self.status_label.config(text="正在停止...")
        log_manager.log("用户请求停止...")

    def run_processing(self):
        """主处理流程"""
        try:
            mode = self.mode_var.get()
            original_path = self.original_path_var.get().strip()
            output_dir = self.output_path_var.get().strip()

            # 获取语言设置
            source_lang = self.source_lang_var.get()
            target_lang = self.target_lang_var.get()
            enable_post_split = self.post_split_var.get()

            model_name = self.model_var.get()
            model_id = AVAILABLE_MODELS[model_name]['id']

            base_name = os.path.splitext(os.path.basename(original_path))[0]
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

            task_dir = os.path.join(output_dir, f"{base_name}_{timestamp}")
            temp_dir = os.path.join(task_dir, "中间文件")

            os.makedirs(task_dir, exist_ok=True)
            os.makedirs(temp_dir, exist_ok=True)

            log_manager.log("=" * 60)
            log_manager.log("📄 多语对照记忆工具 v4.3")
            log_manager.log("=" * 60)
            log_manager.log(f"处理模式: 双文件模式")
            log_manager.log(f"语言对: {source_lang} → {target_lang}")
            log_manager.log(f"后处理分句: {'启用' if enable_post_split else '禁用'}")
            log_manager.log(f"模型: {model_name}")
            log_manager.log(f"输出目录: {task_dir}")
            log_manager.log("=" * 60)

            # 双文件模式
            trans_path = self.trans_path_var.get().strip()
            file_type = get_file_type(original_path)

            try:
                threshold_2 = int(self.threshold_2_var.get())
                threshold_3 = int(self.threshold_3_var.get())
                threshold_4 = int(self.threshold_4_var.get())
                threshold_5 = int(self.threshold_5_var.get())
                threshold_6 = int(self.threshold_6_var.get())
                threshold_7 = int(self.threshold_7_var.get())
                threshold_8 = int(self.threshold_8_var.get())
                buffer_chars = int(self.buffer_var.get())
            except ValueError:
                threshold_2, threshold_3, threshold_4 = THRESHOLD_2_PARTS, THRESHOLD_3_PARTS, THRESHOLD_4_PARTS
                threshold_5, threshold_6, threshold_7, threshold_8 = THRESHOLD_5_PARTS, THRESHOLD_6_PARTS, THRESHOLD_7_PARTS, THRESHOLD_8_PARTS
                buffer_chars = BUFFER_CHARS

            log_manager.log(f"文件类型: {file_type.upper()}")

            # 处理 .doc 文件
            if file_type == 'doc':
                log_manager.log("检测到 .doc 文件，正在转换为 .docx...")
                self.progress_var.set(5)

                converted_orig = convert_doc_to_docx(original_path, temp_dir)
                converted_trans = convert_doc_to_docx(trans_path, temp_dir)

                if converted_orig is None or converted_trans is None:
                    log_manager.log_exception("无法转换 .doc 文件")
                    self.root.after(0, lambda: self.status_label.config(text="转换失败"))
                    return

                original_path = converted_orig
                trans_path = converted_trans
                file_type = 'docx'

            # 处理 Excel 双文件模式（新增功能）
            if file_type == 'excel':
                log_manager.log("「阶段一」处理 Excel 表格文件（双文件模式 - 按单元格位置对应）...")
                self.progress_var.set(20)

                excel_name = f"{base_name}_对齐结果.xlsx"
                out_path = os.path.join(task_dir, excel_name)

                success = process_excel_dual_file_alignment(original_path, trans_path, out_path, model_id,
                                                            source_lang=source_lang, target_lang=target_lang)

                self.progress_var.set(100)

                if success:
                    log_manager.log("🎉 处理完成！")
                    log_manager.log(f"📁 最终结果: {out_path}")
                    self.root.after(0, lambda: self.status_label.config(text="完成！"))
                    self.root.after(0, lambda: messagebox.showinfo("完成", f"处理完成！\n\n最终结果:\n{out_path}"))
                else:
                    self.root.after(0, lambda: self.status_label.config(text="处理失败"))
                return

            # Word/PPT 处理
            log_manager.log("「阶段一」文档分析...")
            self.progress_var.set(10)

            count_a, _ = analyze_document_structure(original_path, source_lang)
            count_b, _ = analyze_document_structure(trans_path, target_lang)

            # 根据语言类型显示字/词
            source_info = SUPPORTED_LANGUAGES.get(source_lang, {})
            target_info = SUPPORTED_LANGUAGES.get(target_lang, {})
            source_unit = "字" if not source_info.get('word_based', True) else "词"
            target_unit = "字" if not target_info.get('word_based', True) else "词"

            log_manager.log(f"原文: {count_a:,} {source_unit}")
            log_manager.log(f"译文: {count_b:,} {target_unit}")

            if file_type == 'pptx':
                split_parts = 1
                log_manager.log("PPT文件：不进行分割")
            else:
                max_count = max(count_a, count_b)
                split_parts = 1
                if max_count > threshold_8:
                    split_parts = 8
                elif max_count > threshold_7:
                    split_parts = 7
                elif max_count > threshold_6:
                    split_parts = 6
                elif max_count > threshold_5:
                    split_parts = 5
                elif max_count > threshold_4:
                    split_parts = 4
                elif max_count > threshold_3:
                    split_parts = 3
                elif max_count > threshold_2:
                    split_parts = 2
                log_manager.log(f"分割策略: {split_parts} 份")

            self.progress_var.set(20)

            tasks_queue = []
            generated_excel_paths = []

            if split_parts > 1 and file_type == 'docx':
                log_manager.log(f"「阶段二」文件分割（缓冲区: {buffer_chars} 字）")

                # 先分割原文，得到分割比例；再用相同比例分割译文，确保内容对齐
                log_manager.log("分割原文（主文档，自主计算分割点）...")
                files_a, part_info_a, split_ratios = smart_split_with_buffer(
                    original_path, split_parts, temp_dir, source_lang, buffer_chars)
                log_manager.log("分割译文（从文档，使用原文的分割比例）...")
                files_b, part_info_b, _ = smart_split_with_buffer(
                    trans_path, split_parts, temp_dir, target_lang, buffer_chars,
                    split_element_ratios=split_ratios)

                for i in range(len(files_a)):
                    excel_name = f"Part{i + 1}_对齐结果.xlsx"
                    out_path = os.path.join(temp_dir, excel_name)
                    tasks_queue.append({
                        'original': files_a[i],
                        'trans': files_b[i],
                        'output': out_path,
                        'anchor_orig': part_info_a[i] if part_info_a else None,
                        'anchor_trans': part_info_b[i] if part_info_b else None,
                        'source_lang': source_lang,
                        'target_lang': target_lang,
                    })
                    generated_excel_paths.append(out_path)
            else:
                excel_name = f"{base_name}_对齐结果.xlsx"
                out_path = os.path.join(task_dir, excel_name)
                # PPT 使用动态生成的提示词
                ppt_prompt = get_ppt_alignment_prompt(source_lang, target_lang) if file_type == 'pptx' else None
                tasks_queue.append({
                    'original': original_path,
                    'trans': trans_path,
                    'output': out_path,
                    'anchor_orig': None,
                    'anchor_trans': None,
                    'system_prompt_override': ppt_prompt,
                    'source_lang': source_lang,
                    'target_lang': target_lang,
                })
                generated_excel_paths.append(out_path)

            self.progress_var.set(30)

            log_manager.log("「阶段三」AI 对齐")
            log_manager.log(f"待处理任务数: {len(tasks_queue)}")
            for i, task in enumerate(tasks_queue):
                log_manager.log(f"  任务 {i + 1}: {os.path.basename(task['output'])}")

            progress_per_task = 50 / len(tasks_queue) if tasks_queue else 50

            for idx, task in enumerate(tasks_queue):
                if self.processing_stopped:
                    log_manager.log("处理已停止")
                    break

                log_manager.log(f"")
                log_manager.log(f"{'=' * 40}")
                log_manager.log(f"处理任务 {idx + 1}/{len(tasks_queue)}: {os.path.basename(task['output'])}")
                log_manager.log(f"原文文件: {os.path.basename(task['original'])}")
                log_manager.log(f"译文文件: {os.path.basename(task['trans'])}")

                success = run_llm_alignment(
                    task['original'],
                    task['trans'],
                    task['output'],
                    model_id,
                    anchor_info_orig=task['anchor_orig'],
                    anchor_info_trans=task['anchor_trans'],
                    system_prompt_override=task.get('system_prompt_override'),
                    source_lang=task.get('source_lang', source_lang),
                    target_lang=task.get('target_lang', target_lang),
                    enable_post_split=enable_post_split
                )

                # 检查处理结果
                output_file = task['output']
                if success:
                    if os.path.exists(output_file):
                        log_manager.log(f"✅ 任务 {idx + 1} 成功: {os.path.basename(output_file)}")
                    else:
                        log_manager.log_exception(f"⚠️ 任务 {idx + 1} 返回成功但文件不存在: {output_file}")
                        if output_file in generated_excel_paths:
                            generated_excel_paths.remove(output_file)
                else:
                    log_manager.log_exception(f"❌ 任务 {idx + 1} 失败: {os.path.basename(output_file)}")
                    if output_file in generated_excel_paths:
                        generated_excel_paths.remove(output_file)

                self.progress_var.set(30 + (idx + 1) * progress_per_task)

            # 显示成功生成的文件列表
            log_manager.log(f"")
            log_manager.log(f"成功生成的文件数: {len(generated_excel_paths)}")
            for path in generated_excel_paths:
                log_manager.log(f"  - {os.path.basename(path)}")

            final_path = None
            if not self.processing_stopped:
                if split_parts > 1 and len(generated_excel_paths) > 0:
                    log_manager.log("")
                    log_manager.log("「阶段四」合并与去重")
                    final_path = os.path.join(task_dir, f"「最终结果」{base_name}_对齐.xlsx")
                    merge_and_deduplicate_excels(generated_excel_paths, final_path, source_lang, target_lang)
                else:
                    final_path = generated_excel_paths[0] if generated_excel_paths else None

            self.progress_var.set(100)

            if self.processing_stopped:
                self.root.after(0, lambda: self.status_label.config(text="已停止"))
            else:
                log_manager.log("🎉 处理完成！")
                if final_path and os.path.exists(final_path):
                    log_manager.log(f"📁 最终结果: {final_path}")

                self.root.after(0, lambda: self.status_label.config(text="完成！"))
                self.root.after(0, lambda: messagebox.showinfo(
                    "完成",
                    f"处理完成！\n\n最终结果:\n{final_path}\n\n任务目录:\n{task_dir}"
                ))

        except Exception as e:
            log_manager.log_exception(f"处理出错", str(e))
            import traceback
            log_manager.log_exception("详细错误", traceback.format_exc())
            self.root.after(0, lambda: self.status_label.config(text="错误"))

        finally:
            self.root.after(0, self.reset_buttons)

    def reset_buttons(self):
        self.start_button.config(state=tk.NORMAL)
        self.stop_button.config(state=tk.DISABLED)

    def on_closing(self):
        self.running = False
        self.processing_stopped = True
        self.root.destroy()


# ==========================================
# === 🎬 主程序 ===
# ==========================================
def main():
    root = tk.Tk()
    app = DocumentAlignerGUI(root)
    root.mainloop()


if __name__ == "__main__":
    main()