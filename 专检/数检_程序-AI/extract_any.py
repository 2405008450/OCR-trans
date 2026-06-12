"""
多格式文本提取器
支持: .docx / .xlsx / .pdf / .pptx
统一返回 list[Segment]
"""
from pathlib import Path
from dataclasses import dataclass


@dataclass
class Segment:
    source: str
    source_label: str
    text: str
    para_index: int = -1
    row_context: str = ""


def extract(path: str | Path) -> list[Segment]:
    path = Path(path)
    ext = path.suffix.lower()
    if ext == ".docx":
        return _from_docx(path)
    elif ext == ".xlsx":
        return _from_xlsx(path)
    elif ext == ".pdf":
        return _from_pdf(path)
    elif ext == ".pptx":
        return _from_pptx(path)
    else:
        raise ValueError(f"不支持的格式: {ext}")


# ── DOCX ──────────────────────────────────────────────────────────────
def _from_docx(path: Path) -> list[Segment]:
    from content import extract_docx_in_order, SOURCE_LABELS
    raw = extract_docx_in_order(path)
    return [
        Segment(
            source=seg.source,
            source_label=SOURCE_LABELS.get(seg.source, seg.source),
            text=seg.text,
        )
        for seg in raw
    ]


# ── XLSX ──────────────────────────────────────────────────────────────
def _from_xlsx(path: Path) -> list[Segment]:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    segments: list[Segment] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                val = cell.value
                if val is None:
                    continue
                text = str(val).strip()
                if text:
                    segments.append(Segment(
                        source="cell",
                        source_label=sheet.title,
                        text=text,
                    ))
    wb.close()
    return segments


# ── PPTX ──────────────────────────────────────────────────────────────
def _from_pptx(path: Path) -> list[Segment]:
    """
    按幻灯片顺序提取文本。
    每张幻灯片：标题优先，然后正文占位符，最后其他文本框和表格。
    source_label 格式：第N页 [标题前20字]
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    segments: list[Segment] = []

    for slide_no, slide in enumerate(prs.slides, start=1):
        # 取幻灯片标题
        slide_title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text_frame.text.strip()
        label = f"第{slide_no}页" + (f" {slide_title[:20]}" if slide_title else "")

        # 按视觉位置排序：标题优先，其余按行分组后列内排序
        # 行分组：按 top 排序，相邻形状 top 差值在行容差内归为同一行
        # 行内再按列分组（left 差值在列容差内视为同列），列内按 top 排，列间按 left 排
        ROW_TOLERANCE = 457200   # ~1.2cm，容忍预排版的垂直微偏
        COL_TOLERANCE = 36000    # ~0.1cm，容忍水平亚像素偏差

        def _group_into_rows(shape_list):
            by_top = sorted(shape_list, key=lambda s: s.top or 0)
            rows: list[list] = []
            for s in by_top:
                top = s.top or 0
                placed = False
                for row in rows:
                    row_top = min(rs.top or 0 for rs in row)
                    if abs(top - row_top) <= ROW_TOLERANCE:
                        row.append(s)
                        placed = True
                        break
                if not placed:
                    rows.append([s])
            # 行内按列分组后排序：列内按 top，列间按 left
            for row in rows:
                col_groups: list[list] = []
                for s in sorted(row, key=lambda s: s.left or 0):
                    left = s.left or 0
                    placed = False
                    for col in col_groups:
                        col_left = min(cs.left or 0 for cs in col)
                        if abs(left - col_left) <= COL_TOLERANCE:
                            col.append(s)
                            placed = True
                            break
                    if not placed:
                        col_groups.append([s])
                for col in col_groups:
                    col.sort(key=lambda s: s.top or 0)
                row[:] = [s for col in col_groups for s in col]
            return rows

        title_shapes = [s for s in slide.shapes if s == slide.shapes.title]
        other_shapes = [s for s in slide.shapes if s != slide.shapes.title]
        rows = _group_into_rows(other_shapes)
        shapes = title_shapes + [s for row in rows for s in row]

        for shape in shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        segments.append(Segment(
                            source="body",
                            source_label=label,
                            text=text,
                        ))
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text_frame.text.strip()
                        if text:
                            segments.append(Segment(
                                source="table",
                                source_label=label,
                                text=text,
                            ))

    return segments


# ── PDF ───────────────────────────────────────────────────────────────
def _from_pdf(path: Path) -> list[Segment]:
    """
    用 pymupdf (fitz) 的 TextBlock 提取 PDF 文本。
    TextBlock 按坐标聚合，比按行提取更接近段落语义。
    """
    import fitz

    segments: list[Segment] = []
    doc = fitz.open(str(path))

    for i, page in enumerate(doc, start=1):
        label = f"第{i}页"
        blocks = page.get_text("blocks", sort=True)
        for b in blocks:
            if b[6] != 0:  # 跳过图片块
                continue
            text = b[4].replace("\n", " ").strip()
            if not text:
                continue
            segments.append(Segment(source="body", source_label=label, text=text))

    doc.close()
    return segments


def align_pdf_segments(src: list[Segment], tgt: list[Segment]) -> list[tuple]:
    """
    按页分组后页内顺序对齐，多出的块插入 None 占位。
    适用于 PDF 和 PPTX 的跨语言对齐。
    """
    from collections import defaultdict

    def page_no(label: str) -> int:
        try:
            return int(label.split()[0].replace("第", "").replace("页", ""))
        except (ValueError, IndexError):
            return 0

    src_by_page: dict[int, list[Segment]] = defaultdict(list)
    tgt_by_page: dict[int, list[Segment]] = defaultdict(list)
    for s in src:
        src_by_page[page_no(s.source_label)].append(s)
    for s in tgt:
        tgt_by_page[page_no(s.source_label)].append(s)

    pairs: list[tuple] = []
    for p in sorted(set(src_by_page) | set(tgt_by_page)):
        sp = src_by_page.get(p, [])
        tp = tgt_by_page.get(p, [])
        for i in range(max(len(sp), len(tp))):
            pairs.append((
                sp[i] if i < len(sp) else None,
                tp[i] if i < len(tp) else None,
            ))

    return pairs
