"""PPTX / PPT 演示文稿解析器"""
import sys
from pathlib import Path


def parse_pptx(file_path: str, mode: str = "clean") -> str:
    """
    解析 PPTX 文件，提取幻灯片中的文本内容。
    需要安装 python-pptx: pip install python-pptx

    Args:
        file_path: 文件路径
        mode: "clean" / "structured"

    Returns:
        解析后的文本
    """
    p = Path(file_path)
    if not p.exists():
        print(f"❌ 文件不存在: {p}", file=sys.stderr)
        return None

    try:
        from pptx import Presentation
    except ImportError:
        print("❌ 需要安装 python-pptx: pip install python-pptx", file=sys.stderr)
        return None

    try:
        prs = Presentation(str(p))
        slides_text = []

        for i, slide in enumerate(prs.slides, 1):
            texts = []
            # 递归提取所有形状（包括组合形状内部）
            _extract_shapes(slide.shapes, texts)

            # 备注页
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"[备注] {notes}")

            if texts:
                if mode == "structured":
                    slides_text.append(f"=== 幻灯片 {i} ===\n" + "\n".join(texts))
                else:
                    slides_text.extend(texts)

        return "\n".join(slides_text)
    except Exception as e:
        print(f"❌ 解析 PPTX 失败: {e}", file=sys.stderr)
        return None


def _extract_shapes(shapes, texts: list):
    """递归提取所有形状中的文本（包括组合形状、表格）"""
    for shape in shapes:
        # 组合形状：递归进入
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                _extract_shapes(shape.shapes, texts)
            except Exception:
                pass

        # 文本框
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)

        # 表格
        if shape.has_table:
            for row in shape.table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    texts.append("\t".join(row_texts))


if __name__ == "__main__":
    r=parse_pptx(r"C:\Users\Administrator\Desktop\多语种标点检查\测试文件\测试文件.pptx")
    print(r)
    # import sys
    # if len(sys.argv) > 1:
    #     print(parse_pptx(sys.argv[1]))
