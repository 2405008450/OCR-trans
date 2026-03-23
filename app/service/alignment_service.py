"""
多语对照记忆工具 - Web Service
将原文与译文文档（DOCX/PPTX/Excel）通过 LLM 进行句级对齐，输出 Excel。
"""

import os
import re
import shutil
import asyncio
import uuid
import traceback
import threading
import importlib.util
from concurrent.futures import Executor
from datetime import datetime
from pathlib import Path
from typing import Optional

import pandas as pd
from lxml import etree
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.config import settings
from app.service.gemini_service import ensure_gemini_route_configured, generate_text

# ── 全局配置 ──────────────────────────────────────────────
ROW_BUCKET = 20_000
_gemini_route_local = threading.local()


def _set_current_gemini_route(route: str) -> None:
    _gemini_route_local.route = route


def _get_current_gemini_route() -> str:
    return getattr(_gemini_route_local, "route", settings.GEMINI_DEFAULT_ROUTE)


THRESHOLD_MAP = {
    2: 25_000, 3: 50_000, 4: 75_000, 5: 100_000,
    6: 125_000, 7: 150_000, 8: 175_000,
}
BUFFER_CHARS = 2000
OUTPUT_DIR = settings.OUTPUT_DIR

AVAILABLE_MODELS = {
    "Google Gemini 2.5 Flash": {
        "id": "google/gemini-2.5-flash",
        "description": "建议先检查文章是否有目录，先将目录删除再处理",
        "max_output": 65536,
    },
    "Google gemini-3-flash-preview": {
        "id": "google/gemini-3-flash-preview",
        "description": "建议先检查文章是否有目录，先将目录删除再处理",
        "max_output": 65536,
    },
    "Google Gemini 2.5 Pro": {
        "id": "google/gemini-2.5-pro",
        "description": "PPT推荐-增强，速度稍慢，100万上下文，65K输出",
        "max_output": 65536,
    },
    "Google: Gemini 3 Pro Preview": {
        "id": "google/gemini-3-pro-preview",
        "description": "最强推理，100万上下文，65K输出",
        "max_output": 65536,
    },
        "Google: google/gemini-3.1-pro-preview": {
        "id": "google/gemini-3.1-pro-preview",
        "description": "最强推理，100万上下文，65K输出",
        "max_output": 65536,
    },
}
DEFAULT_MODEL = "Google gemini-3-flash-preview"

CHAPTER_PATTERNS = [
    r'^第[一二三四五六七八九十百千\d]+[章节篇部]', r'^Chapter\s*\d+', r'^CHAPTER\s*\d+',
    r'^\d+[\.、]\s*\S+', r'^[一二三四五六七八九十]+[、.]\s*\S+',
    r'^Part\s*\d+', r'^PART\s*\d+', r'^Section\s*\d+',
]

SUPPORTED_LANGUAGES = {
    "中文": {"code": "zh", "english_name": "Chinese", "char_pattern": r'[\u4e00-\u9fa5]', "word_based": False, "description": "中文（简体/繁体）"},
    "英语": {"code": "en", "english_name": "English", "char_pattern": r'\b[a-zA-Z]+\b', "word_based": True, "description": "English"},
    "西班牙语": {"code": "es", "english_name": "Spanish", "char_pattern": r'\b[a-zA-ZáéíóúüñÁÉÍÓÚÜÑ]+\b', "word_based": True, "description": "Español"},
    "葡语": {"code": "pt", "english_name": "Portuguese", "char_pattern": r'\b[a-zA-ZáéíóúâêôãõçÁÉÍÓÚÂÊÔÃÕÇ]+\b', "word_based": True, "description": "Português"},
    "日语": {"code": "ja", "english_name": "Japanese", "char_pattern": r'[\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FAF]', "word_based": False, "description": "日本語"},
    "俄语": {"code": "ru", "english_name": "Russian", "char_pattern": r'\b[а-яА-ЯёЁ]+\b', "word_based": True, "description": "Русский"},
    "韩语": {"code": "ko", "english_name": "Korean", "char_pattern": r'[\uAC00-\uD7AF\u1100-\u11FF]', "word_based": False, "description": "한국어"},
    "阿语": {"code": "ar", "english_name": "Arabic", "char_pattern": r'[\u0600-\u06FF\u0750-\u077F]+', "word_based": True, "description": "العربية"},
    "法语": {"code": "fr", "english_name": "French", "char_pattern": r'\b[a-zA-ZàâäéèêëïîôùûüÿœæçÀÂÄÉÈÊËÏÎÔÙÛÜŸŒÆÇ]+\b', "word_based": True, "description": "Français"},
    "波兰语": {"code": "pl", "english_name": "Polish", "char_pattern": r'\b[a-zA-ZąćęłńóśźżĄĆĘŁŃÓŚŹŻ]+\b', "word_based": True, "description": "Polski"},
    "意大利语": {"code": "it", "english_name": "Italian", "char_pattern": r'\b[a-zA-ZàèéìíîòóùúÀÈÉÌÍÎÒÓÙÚ]+\b', "word_based": True, "description": "Italiano"},
    "德语": {"code": "de", "english_name": "German", "char_pattern": r'\b[a-zA-ZäöüßÄÖÜ]+\b', "word_based": True, "description": "Deutsch"},
}

# ── 进度追踪（线程安全）────────────────────────────────────
_task_progress: dict = {}
_progress_lock = threading.Lock()
_memory_module = None
# 当前执行任务的 task_id，用于将 memory 的 log_stream 写入该任务的 stream_log
_stream_task_id = threading.local()
_log_patch_installed = False


def _get_memory_module():
    global _memory_module
    if _memory_module is not None:
        return _memory_module

    memory_file = Path(__file__).resolve().parents[2] / "memory" / "memory.py"
    spec = importlib.util.spec_from_file_location("memory_legacy_module", str(memory_file))
    if spec is None or spec.loader is None:
        raise RuntimeError(f"无法加载 memory.py: {memory_file}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    _memory_module = module
    return _memory_module


def _install_log_patches():
    """一次性安装日志补丁，通过 threading.local 路由到正确的任务（线程安全）"""
    global _log_patch_installed
    if _log_patch_installed:
        return

    memory_module = _get_memory_module()
    _orig_log_stream = memory_module.log_manager.log_stream
    _orig_log = memory_module.log_manager.log
    _orig_log_exception = memory_module.log_manager.log_exception

    def _safe_log_stream(content):
        _orig_log_stream(content)
        tid = getattr(_stream_task_id, "task_id", None)
        if tid:
            with _progress_lock:
                if tid in _task_progress:
                    cur = _task_progress[tid].get("stream_log", "")
                    _task_progress[tid]["stream_log"] = cur + (content if isinstance(content, str) else str(content))

    def _safe_log(message, level="INFO"):
        _orig_log(message, level)
        print(f"[memory-log] {message}")
        tid = getattr(_stream_task_id, "task_id", None)
        if tid:
            with _progress_lock:
                if tid in _task_progress:
                    ts = datetime.now().strftime("%H:%M:%S")
                    line = f"[{ts}] {message}\n"
                    cur = _task_progress[tid].get("stream_log", "")
                    _task_progress[tid]["stream_log"] = cur + line

    def _safe_log_exception(message, data=None):
        _orig_log_exception(message, data)
        print(f"[memory-ERR] {message}" + (f" | {str(data)[:300]}" if data else ""))
        tid = getattr(_stream_task_id, "task_id", None)
        if tid:
            with _progress_lock:
                if tid in _task_progress:
                    ts = datetime.now().strftime("%H:%M:%S")
                    line = f"[{ts}] ⚠️ {message}"
                    if data:
                        line += f" | {str(data)[:200]}"
                    cur = _task_progress[tid].get("stream_log", "")
                    _task_progress[tid]["stream_log"] = cur + line + "\n"

    memory_module.log_manager.log_stream = _safe_log_stream
    memory_module.log_manager.log = _safe_log
    memory_module.log_manager.log_exception = _safe_log_exception
    _log_patch_installed = True
    print("[alignment] ✅ 日志补丁已永久安装（线程安全模式）")


def _update_progress(task_id: str, progress: int, message: str, **extra):
    data = {
        "status": "processing",
        "progress": progress,
        "message": message,
        **extra,
    }
    with _progress_lock:
        # 保留已有 stream_log，避免被覆盖
        if task_id in _task_progress and "stream_log" in _task_progress[task_id]:
            data.setdefault("stream_log", _task_progress[task_id]["stream_log"])
        _task_progress[task_id] = data


def _complete_task(task_id: str, *, result: dict = None, error: str = None):
    with _progress_lock:
        if error:
            stream_log = _task_progress.get(task_id, {}).get("stream_log", "")
            _task_progress[task_id] = {
                "status": "failed", "progress": 100, "message": "处理失败", "error": error,
                "stream_log": stream_log,
            }
        else:
            stream_log = _task_progress.get(task_id, {}).get("stream_log", "")
            if result is not None:
                result = dict(result)
                result["stream_log"] = stream_log
            _task_progress[task_id] = {"status": "done", "progress": 100, "message": "处理完成", "result": result}


def get_alignment_progress(task_id: str) -> Optional[dict]:
    with _progress_lock:
        data = _task_progress.get(task_id)
        return dict(data) if data else None


# ── Prompt 生成 ───────────────────────────────────────────
def _get_ppt_alignment_prompt(source_lang="中文", target_lang="英语"):
    return f"""
你是一个「PPT 双语文本对齐器」。

当前任务：将 {source_lang} 原文与 {target_lang} 译文进行对齐。

═══════════════════════════════════════
【核心铁律 - 违反任意一条即为失败】
═══════════════════════════════════════
1. 禁止编造：左侧只能填原文中实际存在的文字，右侧只能填译文中实际存在的文字
2. 禁止占位符：绝对不允许出现 XXX、---、...、[空]、[无]、N/A 等任何占位符号
3. 禁止跨页：幻灯片 N 的内容只能与幻灯片 N 的内容配对
4. 分隔符独立：「---- 幻灯片 N ----」必须单独成行，不参与 ||| 配对

═══════════════════════════════════════
【输出格式 - 严格遵守】
═══════════════════════════════════════
---- 幻灯片 1 ----
{source_lang}行A ||| {target_lang}行A
{source_lang}行B ||| {target_lang}行B
{source_lang}行C |||
||| {target_lang}行D
---- 幻灯片 2 ----
...

格式说明：
- 分隔行「---- 幻灯片 N ----」独立一行，前后不加 |||
- 匹配成功：原文 ||| 译文
- 原文无对应译文：原文 |||
- 译文无对应原文：||| 译文

═══════════════════════════════════════
【对齐策略 - 按优先级执行】
═══════════════════════════════════════

第一步：锚点识别（同页内）
  ├─ 数字/日期锚点：12.15 ↔ December 15，2024年 ↔ 2024
  ├─ 专名锚点：公司名、人名、产品名（即使一侧是音译/意译）
  ├─ 编号锚点：① ↔ 1)，1. ↔ (1)，注1 ↔ Note 1
  └─ 标题锚点：通常字体大、位置靠上、内容概括性强

第二步：脚注特殊处理
  ├─ 识别脚注区域：通常在页面底部，带上标数字或特殊标记
  ├─ 编号与内容绑定：「1 XXX说明文字」整体视为一个脚注单元
  ├─ 配对原则：原文脚注1 ↔ 译文脚注1（按编号，非按位置）
  └─ 不可混配：脚注编号不能与正文内容配对

第三步：语义配对
  ├─ 含义最接近的行 1:1 配对
  ├─ 允许合理拆分：一长句 ↔ 两短句（语义完整时）
  ├─ 允许合理合并：两短句 ↔ 一长句（语义完整时）
  └─ 每个文本片段只能使用一次

第四步：处理剩余
  ├─ 原文有、译文无 → 「原文 |||」
  ├─ 译文有、原文无 → 「||| 译文」
  └─ 绝不填充任何虚构内容

═══════════════════════════════════════
【执行检查清单】
═══════════════════════════════════════
输出前逐项确认：
□ 分隔符是否独立成行？
□ 是否存在任何 XXX/---/.../[空] 等占位符？
□ 脚注编号是否与对应编号配对（而非与正文配对）？
□ 每个原文片段是否只出现一次？
□ 每个译文片段是否只出现一次？
□ 是否有跨页配对？

只输出对齐结果，不输出任何解释说明。
"""


def _get_docx_alignment_prompt(source_lang="中文", target_lang="英语"):
    source_info = SUPPORTED_LANGUAGES.get(source_lang, SUPPORTED_LANGUAGES["中文"])
    target_info = SUPPORTED_LANGUAGES.get(target_lang, SUPPORTED_LANGUAGES["英语"])

    source_is_cjk = source_lang in ["中文", "日语", "韩语"]
    if source_is_cjk:
        punctuation_marks = "。！？"
        punctuation_desc = "。！？（全角中日韩句末标点）"
    else:
        punctuation_marks = ". ! ?"
        punctuation_desc = ". ! ?（半角西文句末标点）"

    sentence_rule = f"""2. 🎯 断句规则：{source_lang}主导，{target_lang}跟随（强优先级）

   第一步：只看{source_lang}，按以下标点断句
   - 断句标点：{punctuation_desc}
   - 遇到这些标点就断开，形成一个{source_lang}片段

   第二步：{target_lang}匹配{source_lang}
   - 在 Stream B 中找出与该{source_lang}片段语义对应的{target_lang}部分
   - {target_lang}可能是1句、半句、或多句，都没关系
   - 译文位置可能有错位，按语义匹配即可
   - 只要语义完全覆盖该{source_lang}片段即可

3. 质量检查
   - 左边的每个{source_lang}片段是否都以正确的句末标点（{punctuation_marks}）结尾？
   - 右边{target_lang}的语义是否完全对应左边{source_lang}？
   - 是否有任何内容被修改或自行翻译？（必须为否）"""

    return f"""你是一个「双流文本对齐同步器」(Dual-Stream Aligner)。

核心任务
将 Stream A ({source_lang}原文) 和 Stream B ({target_lang}译文) 进行精确对齐。

⛔ 绝对铁律
1. 来源锁定（最高优先级）
   - "|||" 左侧 = 必须100%来自 Stream A，一个字都不能改
   - "|||" 右侧 = 必须100%来自 Stream B，一个字都不能改
   - 「禁止」自己翻译或编造任何内容
   - 即使发现原文有错别字、译文有翻译错误，也必须原样保留

{sentence_rule}

4. 空行处理规则
   - 如果原文或译文中有空行（连续换行符），这些空行已经被适当压缩保留
   - 空行通常用于分隔段落或章节
   - 在对齐时，忽略空行的位置差异，专注于有内容的文本对齐
   - 不要输出空行对应的对齐结果（空|||空）

输出格式
{source_lang}片段 ||| 对应的{target_lang}内容
"""


def _get_split_row_prompt(source_lang="中文"):
    source_is_cjk = source_lang in ["中文", "日语", "韩语"]
    if source_is_cjk:
        punctuation_marks = "。！？"
        punctuation_desc = "。！？（全角句末标点）"
        abbreviation_rule = ""
        numbering_rule = "1. / 2. / 11. / 1）/ 2）/ ①② 等数字编号"
    else:
        punctuation_marks = ". ! ?"
        punctuation_desc = ". ! ?（半角句末标点）"
        numbering_rule = """- 数字编号：1. / 2. / 11. / (1) / (2) 等
   - 字母编号：A. / B. / C. / a. / b. / c. / (A) / (B) 等
   - 罗马数字：I. / II. / III. / i. / ii. / iii. 等
   - 混合编号：1.1 / 1.2 / A.1 / A.2 等（在最后一个点后断开）"""
        abbreviation_rule = """
⚠️ 缩写例外（这些句号不表示句末，不要在此断开）：
- 学术引用：et al. / etc. / e.g. / i.e. / vs. / cf. / ibid.
- 称谓：Mr. / Mrs. / Ms. / Dr. / Prof. / Jr. / Sr. / St.
- 公司：Inc. / Ltd. / Co. / Corp. / L.L.C.
- 国家/组织：U.S. / U.K. / U.N. / E.U.
- 时间：a.m. / p.m. / A.M. / P.M.
- 引用标记：No. / Vol. / Fig. / Ch. / Sec. / pp. / p.
- 其他：approx. / est. / max. / min. / avg.

💡 区分编号与缩写的方法：
- 编号特征：单独的数字或字母 + 句号，如 "1." "A." "II."
- 缩写特征：多个字母组成的词 + 句号，如 "et al." "Dr." "Inc."
- 编号后面通常跟着正文内容，缩写后面通常还有更多文字"""

    if not source_is_cjk:
        example_section = """
## 示例

### 示例1：数字编号 + 作者信息
输入：
【原文】11. Piacentini MG, et al. A randomized controlled trial on the effect of Solidago virgaurea extract.
【译文】11.Piacentini MG，等。一项关于一枝黄花提取物影响的随机对照试验。

正确输出（3行）：
11. ||| 11.
Piacentini MG, et al. ||| Piacentini MG，等。
A randomized controlled trial on the effect of Solidago virgaurea extract. ||| 一项关于一枝黄花提取物影响的随机对照试验。

### 示例2：字母编号
输入：
【原文】A. Introduction. B. Methods. C. Results.
【译文】A. 引言。B. 方法。C. 结果。

正确输出（6行）：
A. ||| A.
Introduction. ||| 引言。
B. ||| B.
Methods. ||| 方法。
C. ||| C.
Results. ||| 结果。

### 示例3：混合编号
输入：
【原文】1.1 Background. 1.2 Objectives.
【译文】1.1 背景。1.2 目标。

正确输出（4行）：
1.1 ||| 1.1
Background. ||| 背景。
1.2 ||| 1.2
Objectives. ||| 目标。

### 关键判断规则：
- "A." "B." "1." "11." "1.1" → 编号，后面断开
- "et al." "Dr." "U.S." "Inc." → 缩写，不在此断开
- 编号 = 单独数字/字母 + 点号
- 缩写 = 多字母词汇 + 点号
"""
    else:
        example_section = ""
        numbering_rule = "1. / 2. / 11. / 1）/ 2）/ ①② 等数字编号"

    return f"""你是一个精确的「单行分句对齐器」。

## 任务
根据【原文】的分割点，将【译文】对应拆分对齐。

## ⛔ 铁律

### 1. 分割点识别（必须在以下位置断开）

#### 1.1 编号（编号后必须断开，编号单独成行）
   {numbering_rule}

#### 1.2 句末标点（{punctuation_desc}）
   每个句号/感叹号/问号后都要断开
{abbreviation_rule}

### 2. 禁止修改内容
- 左侧必须100%来自原文，一字不改
- 右侧必须100%来自译文，一字不改
- 禁止翻译、编造、补充任何内容

### 3. 译文对齐规则
- 根据语义将译文拆分，匹配原文的每个分割单元
- 如果译文标点与原文不一致，按语义对齐

### 4. 分割原则
- 编号单独成行（如 "11." / "A." / "1.1"）
- 作者信息单独成行（如 "Piacentini MG, et al."）
- 每个完整句子单独成行
- 即使片段很短也要独立成行
- 输出行数 = 原文分割单元数
{example_section}
## 输出格式
每行一对，用 ||| 分隔：
原文片段1 ||| 译文对应部分1
原文片段2 ||| 译文对应部分2

只输出对齐结果，不要任何解释。"""


def _get_table_cell_split_prompt(source_lang="中文"):
    source_is_cjk = source_lang in ["中文", "日语", "韩语"]

    if source_is_cjk:
        punctuation_marks = "。！？"
        punctuation_desc = "。！？（全角句末标点）"
    else:
        punctuation_marks = ". ! ?"
        punctuation_desc = ". ! ?（半角句末标点）"

    return f"""你是一个精确的「表格单元格分句对齐器」。

## 任务
对表格单元格内的原文和译文进行细粒度分句对齐。

## 分句规则（必须严格遵守）

### 1. 断句位置（只在以下位置断开）
- **句末标点**：{punctuation_desc}
- **换行符**：当原文中有空行分隔时，在空行处断开

### 2. 【重要】不是断句点的情况
- **序号标点不断句**：如"1."、"2."、"3."、"①"、"②"、"(1)"、"a."等序号后的点号，这是序号的一部分，不是句末
- **序号必须保留在句子开头**：序号和后面的内容属于同一句
- 只有真正表示句子结束的标点（{punctuation_marks}）才断句

### 3. 分句原则
- 以原文的分句结构为准
- 每个完整句子（以{punctuation_marks}结尾）独立成行
- 当遇到空行分隔的段落时，按段落断开
- 保持原文和译文的句子一一对应

### 4. 禁止修改内容
- 左侧必须100%来自原文，一字不改
- 右侧必须100%来自译文，一字不改
- 禁止翻译、编造、补充任何内容
- 保留所有标点符号和序号

### 5. 译文对齐规则
- 根据语义将译文拆分，匹配原文的每个句子
- 如果译文标点与原文不完全一致，按语义对齐

## 示例1（带序号的多条内容）

### 输入：
【原文】
1.这是第一条说明。需要注意这个问题。

2.这是第二条说明；

【译文】
1. This is the first instruction. Please note this issue.
2. This is the second instruction;

### 正确输出：
1.这是第一条说明。 ||| 1. This is the first instruction.
需要注意这个问题。 ||| Please note this issue.
2.这是第二条说明； ||| 2. This is the second instruction;

### 错误输出（绝对禁止）：
1. ||| 1.
这是第一条说明。 ||| This is the first instruction.
（错误原因：把序号"1."单独拆分了，序号必须和后面的内容在一起）

## 示例2（普通多句内容）

### 输入：
【原文】
这是第一句。这是第二句！
这是换行后的第三句？

【译文】
This is the first sentence. This is the second sentence!
This is the third sentence after newline?

### 正确输出：
这是第一句。 ||| This is the first sentence.
这是第二句！ ||| This is the second sentence!
这是换行后的第三句？ ||| This is the third sentence after newline?

## 输出格式
每行一对，用 ||| 分隔：
原文句子1 ||| 译文对应部分1
原文句子2 ||| 译文对应部分2

只输出对齐结果，不要任何解释。"""


# ── 文件工具 ──────────────────────────────────────────────
def _get_file_type(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    return {'.docx': 'docx', '.pptx': 'pptx', '.xlsx': 'excel', '.xls': 'excel'}.get(ext, 'unknown')


def _get_text_count(text: str, lang_name: str) -> int:
    if not text:
        return 0
    lang_info = SUPPORTED_LANGUAGES.get(lang_name)
    if lang_info:
        return len(re.findall(lang_info['char_pattern'], text))
    return len(re.findall(r'\b[a-zA-Z0-9-]+\b', text))


def _iter_body_elements(body, doc):
    """递归遍历 body 下所有段落和表格，包括 sdt（内容控件）内部的元素"""
    for child in body.iterchildren():
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'p':
            yield ('p', Paragraph(child, doc))
        elif tag == 'tbl':
            yield ('tbl', Table(child, doc))
        elif tag == 'sdt':
            for sub in child.iterchildren():
                sub_tag = sub.tag.split('}')[-1] if '}' in sub.tag else sub.tag
                if sub_tag == 'sdtContent':
                    yield from _iter_body_elements(sub, doc)
                elif sub_tag == 'p':
                    yield ('p', Paragraph(sub, doc))
                elif sub_tag == 'tbl':
                    yield ('tbl', Table(sub, doc))


def _get_all_content_elements(doc):
    elements = []
    if hasattr(doc, 'element') and hasattr(doc.element, 'body'):
        for elem_type, elem in _iter_body_elements(doc.element.body, doc):
            elements.append(elem)
    if not elements:
        for para in doc.paragraphs:
            elements.append(para)
        for table in doc.tables:
            elements.append(table)
    return elements


def _get_element_text_count(element, lang_name: str) -> int:
    text = ""
    if isinstance(element, Paragraph):
        text = element.text
    elif isinstance(element, Table):
        for row in element.rows:
            for cell in row.cells:
                text += cell.text + " "
    return _get_text_count(text, lang_name)


def _read_full_docx(file_path: str) -> str:
    """读取完整的docx文件内容，包括段落、表格、页眉页脚、文本框、脚注尾注
    （与 memory.py 原版 read_full_docx 保持一致）
    """
    try:
        abs_path = os.path.abspath(file_path)
        print(f"[alignment-read] 读取 DOCX: {abs_path} (大小: {os.path.getsize(abs_path)} bytes)")
        doc = Document(abs_path)
        full_text = []
        consecutive_empty = 0

        # 1. 按文档顺序遍历所有元素（段落和表格），含 sdt 内容控件内元素
        if hasattr(doc, 'element') and hasattr(doc.element, 'body'):
            for elem_type, elem in _iter_body_elements(doc.element.body, doc):
                if elem_type == 'p':
                    if elem.text.strip():
                        full_text.append(elem.text)
                        consecutive_empty = 0
                    else:
                        consecutive_empty += 1
                        if consecutive_empty <= 2 and full_text:
                            full_text.append("")
                elif elem_type == 'tbl':
                    consecutive_empty = 0
                    seen_cells = set()
                    for row in elem.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text and cell_text not in seen_cells:
                                full_text.append(cell_text)
                                seen_cells.add(cell_text)
        else:
            consecutive_empty = 0
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
                    consecutive_empty = 0
                else:
                    consecutive_empty += 1
                    if consecutive_empty <= 2 and full_text:
                        full_text.append("")
            for table in doc.tables:
                seen_cells = set()
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen_cells:
                            full_text.append(cell_text)
                            seen_cells.add(cell_text)

        # 2. 页眉页脚（段落、表格、文本框）
        for section in doc.sections:
            for p in section.header.paragraphs:
                if p.text.strip():
                    full_text.append(p.text)
            for table in section.header.tables:
                seen_header_cells = set()
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen_header_cells:
                            full_text.append(cell_text)
                            seen_header_cells.add(cell_text)
            if section.header._element is not None:
                header_xml = etree.tostring(section.header._element, encoding='unicode')
                header_root = etree.fromstring(header_xml.encode('utf-8'))
                nsmap_header = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                for txbx in header_root.xpath('.//w:txbxContent', namespaces=nsmap_header):
                    txbx_text = ''.join(txbx.xpath('.//w:t/text()', namespaces=nsmap_header))
                    if txbx_text.strip():
                        full_text.append(txbx_text.strip())

            for p in section.footer.paragraphs:
                if p.text.strip():
                    full_text.append(p.text)
            for table in section.footer.tables:
                seen_footer_cells = set()
                for row in table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text and cell_text not in seen_footer_cells:
                            full_text.append(cell_text)
                            seen_footer_cells.add(cell_text)
            if section.footer._element is not None:
                footer_xml = etree.tostring(section.footer._element, encoding='unicode')
                footer_root = etree.fromstring(footer_xml.encode('utf-8'))
                nsmap_footer = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
                for txbx in footer_root.xpath('.//w:txbxContent', namespaces=nsmap_footer):
                    txbx_text = ''.join(txbx.xpath('.//w:t/text()', namespaces=nsmap_footer))
                    if txbx_text.strip():
                        full_text.append(txbx_text.strip())

        # 3. 文本框（传统格式 + DrawingML 格式）
        if hasattr(doc.element, 'xml'):
            xml = doc.element.xml
            root = etree.fromstring(xml.encode('utf-8'))

            nsmap = {
                'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'wps': 'http://schemas.microsoft.com/office/word/2010/wordprocessingShape',
                'wpg': 'http://schemas.microsoft.com/office/word/2010/wordprocessingGroup',
                'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
                'w14': 'http://schemas.microsoft.com/office/word/2010/wordml',
            }

            textbox_texts = set()

            # 方式1: 传统文本框 (w:txbxContent)
            try:
                for container in root.xpath('.//w:txbxContent', namespaces=nsmap):
                    para_texts = []
                    for p in container.xpath('.//w:p', namespaces=nsmap):
                        p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged = '\n'.join(para_texts)
                    if merged.strip():
                        textbox_texts.add(merged.strip())
            except Exception:
                pass

            # 方式2: DrawingML 文本框 (wps:txbx)
            try:
                for container in root.xpath('.//wps:txbx', namespaces=nsmap):
                    para_texts = []
                    for p in container.xpath('.//a:p', namespaces=nsmap):
                        p_text = ''.join([t.text for t in p.xpath('.//a:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged = '\n'.join(para_texts)
                    if merged.strip():
                        textbox_texts.add(merged.strip())
            except Exception:
                pass

            # 方式3: DrawingML 形状文本 (a:txBody)
            try:
                for container in root.xpath('.//a:txBody', namespaces=nsmap):
                    para_texts = []
                    for p in container.xpath('.//a:p', namespaces=nsmap):
                        p_text = ''.join([t.text for t in p.xpath('.//a:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged = '\n'.join(para_texts)
                    if merged.strip():
                        textbox_texts.add(merged.strip())
            except Exception:
                pass

            # 方式4: Word 2010+ 文本框 (wps:wsp//wps:txbx)
            try:
                for container in root.xpath('.//wps:wsp//wps:txbx', namespaces=nsmap):
                    para_texts = []
                    for p in container.xpath('.//w:p', namespaces=nsmap):
                        p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap) if t.text])
                        if p_text.strip():
                            para_texts.append(p_text.strip())
                    merged = '\n'.join(para_texts)
                    if merged.strip():
                        textbox_texts.add(merged.strip())
            except Exception:
                pass

            for text in textbox_texts:
                full_text.append(text)

        # 4. 脚注、尾注
        if hasattr(doc, 'part'):
            nsmap_fn = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
            for rel in doc.part.rels.values():
                ref = rel.target_ref

                if "footnotes" in ref:
                    try:
                        fn_root = etree.fromstring(rel.target_part.blob)
                        for fn in fn_root.xpath('.//w:footnote', namespaces=nsmap_fn):
                            fn_type = fn.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type')
                            if fn_type in ['separator', 'continuationSeparator']:
                                continue
                            fn_texts = []
                            for p in fn.xpath('.//w:p', namespaces=nsmap_fn):
                                p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap_fn) if t.text])
                                if p_text.strip():
                                    fn_texts.append(p_text.strip())
                            fn_text = ' '.join(fn_texts)
                            if fn_text.strip():
                                full_text.append(fn_text.strip())
                    except Exception:
                        pass

                elif "endnotes" in ref:
                    try:
                        en_root = etree.fromstring(rel.target_part.blob)
                        for en in en_root.xpath('.//w:endnote', namespaces=nsmap_fn):
                            en_type = en.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}type')
                            if en_type in ['separator', 'continuationSeparator']:
                                continue
                            en_texts = []
                            for p in en.xpath('.//w:p', namespaces=nsmap_fn):
                                p_text = ''.join([t.text for t in p.xpath('.//w:t', namespaces=nsmap_fn) if t.text])
                                if p_text.strip():
                                    en_texts.append(p_text.strip())
                            en_text = ' '.join(en_texts)
                            if en_text.strip():
                                full_text.append(en_text.strip())
                    except Exception:
                        pass

        result = "\n".join(full_text)
        result = re.sub(r'\n{4,}', '\n\n\n', result)
        print(f"[alignment-read]   读取完成: {len(full_text)} 段, 总长 {len(result)} 字符")
        if not result.strip():
            print(f"[alignment-read]   警告: 文档内容为空！")
        return result
    except Exception as e:
        print(f"[alignment-read] 读取 DOCX 失败: {e}")
        traceback.print_exc()
        return ""


def _iter_group_shapes(group_shape, base_top=0, base_left=0):
    for sub in group_shape.shapes:
        top = base_top + (sub.top or 0)
        left = base_left + (sub.left or 0)
        if sub.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_group_shapes(sub, top, left)
            continue
        if hasattr(sub, "has_text_frame") and sub.has_text_frame:
            para_texts = [p.text.strip() for p in sub.text_frame.paragraphs if p.text.strip()]
            if para_texts:
                yield (top, left, '\n'.join(para_texts))
        if hasattr(sub, "has_table") and sub.has_table:
            for r_idx, row in enumerate(sub.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    para_texts = [p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip()]
                    if para_texts:
                        yield (top + r_idx * 1000, left + c_idx * 1000, '\n'.join(para_texts))


def _extract_slide_items(slide):
    items = []
    for shape in slide.shapes:
        top = shape.top or 0
        left = shape.left or 0
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            para_texts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
            if para_texts:
                items.append((top, left, '\n'.join(para_texts)))
        if hasattr(shape, "has_table") and shape.has_table:
            for r_idx, row in enumerate(shape.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    para_texts = [p.text.strip() for p in cell.text_frame.paragraphs if p.text.strip()]
                    if para_texts:
                        items.append((top + r_idx * 1000, left + c_idx * 1000, '\n'.join(para_texts)))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            items.extend(_iter_group_shapes(shape, top, left))
    items.sort(key=lambda t: (t[0] // ROW_BUCKET, t[1]))
    return [t[2] for t in items]


def _read_full_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        lines = []
        for idx, slide in enumerate(prs.slides, 1):
            lines.append(f"---- 幻灯片 {idx} ----")
            lines.extend(_extract_slide_items(slide))
        return "\n".join(lines)
    except Exception as e:
        print(f"读取 PPTX 失败: {e}")
        return ""


def _read_full_excel(file_path: str) -> str:
    try:
        df = pd.read_excel(file_path)
        texts = []
        for col in df.columns:
            texts.append(f"[列: {col}]")
            for val in df[col].dropna():
                if str(val).strip():
                    texts.append(str(val).strip())
        return "\n".join(texts)
    except Exception as e:
        print(f"读取 Excel 失败: {e}")
        return ""


def _read_file_content(file_path: str) -> str:
    abs_path = os.path.abspath(file_path)
    ft = _get_file_type(abs_path)
    if ft == 'docx':
        return _read_full_docx(abs_path)
    elif ft == 'pptx':
        return _read_full_pptx(abs_path)
    elif ft == 'excel':
        return _read_full_excel(abs_path)
    print(f"[alignment-read] 不支持的文件类型: {ft} ({abs_path})")
    return ""


# ── Excel 多工作簿读取 ───────────────────────────────────
def _read_excel_all_sheets(file_path: str):
    try:
        excel_file = pd.ExcelFile(file_path, engine='openpyxl')
        sheets = {}
        for name in excel_file.sheet_names:
            sheets[name] = pd.read_excel(file_path, sheet_name=name, header=None, engine='openpyxl')
        return sheets
    except Exception:
        pass
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, data_only=True, read_only=True)
        sheets = {}
        for name in wb.sheetnames:
            ws = wb[name]
            data = [[cell.value for cell in row] for row in ws.iter_rows()]
            sheets[name] = pd.DataFrame(data) if data else pd.DataFrame()
        wb.close()
        return sheets
    except Exception as e:
        print(f"读取 Excel 失败: {e}")
        return None


# ── 文档分析与分割 ────────────────────────────────────────
def _analyze_document(doc_path: str, lang_name: str):
    abs_path = os.path.abspath(doc_path)
    ft = _get_file_type(abs_path)
    if ft == 'pptx':
        text = _read_full_pptx(abs_path)
        return _get_text_count(text, lang_name), 0
    elif ft == 'excel':
        text = _read_full_excel(abs_path)
        return _get_text_count(text, lang_name), 0
    else:
        doc = Document(abs_path)
        elements = _get_all_content_elements(doc)
        total = sum(_get_element_text_count(el, lang_name) for el in elements)
        return total, len(elements)


def _find_element_index_by_chars(element_counts, target_chars):
    left, right = 0, len(element_counts) - 1
    best = 0
    while left <= right:
        mid = (left + right) // 2
        if element_counts[mid] < target_chars:
            best = mid
            left = mid + 1
        else:
            right = mid - 1
    if best + 1 < len(element_counts):
        if element_counts[best + 1] - target_chars < target_chars - element_counts[best]:
            return best + 1
    return best


def _find_buffer_end(element_counts, split_idx, buffer_chars, direction='right'):
    n = len(element_counts)
    if n == 0:
        return split_idx
    if direction == 'right':
        base = element_counts[split_idx] if split_idx < n else element_counts[-1]
        target = base + buffer_chars
        for i in range(split_idx + 1, n):
            if element_counts[i] >= target:
                return i + 1
        return n
    else:
        base = element_counts[split_idx] if split_idx < n else element_counts[-1]
        target = base - buffer_chars
        if target <= 0:
            return 0
        for i in range(split_idx - 1, -1, -1):
            if element_counts[i] <= target:
                return i
        return 0


def _extract_text_from_elements(elements, start_idx, end_idx):
    texts = []
    for i in range(start_idx, min(end_idx, len(elements))):
        elem = elements[i]
        if isinstance(elem, Paragraph) and elem.text.strip():
            texts.append(elem.text.strip())
        elif isinstance(elem, Table):
            for row in elem.rows:
                row_texts = [c.text.strip() for c in row.cells if c.text.strip()]
                if row_texts:
                    texts.append(" | ".join(row_texts))
    return "\n".join(texts)


def _delete_elements_in_range(doc, start_idx, end_idx):
    all_elements = _get_all_content_elements(doc)
    to_delete = []
    for i in range(len(all_elements)):
        if start_idx <= i < end_idx:
            if isinstance(all_elements[i], (Paragraph, Table)):
                to_delete.append(all_elements[i]._element)
    for el in to_delete:
        parent = el.getparent()
        if parent is not None:
            parent.remove(el)


def _smart_split_with_buffer(src_path, num_parts, output_dir, lang_type, buffer_chars=2000,
                             split_element_ratios=None):
    src_path = os.path.abspath(src_path)
    doc = Document(src_path)
    elements = _get_all_content_elements(doc)
    base_name = os.path.splitext(os.path.basename(src_path))[0]

    element_counts = []
    cumulative = 0
    for elem in elements:
        cumulative += _get_element_text_count(elem, lang_type)
        element_counts.append(cumulative)

    total_count = cumulative
    if total_count == 0:
        return [], [], []

    target_per_part = total_count // num_parts

    if split_element_ratios is not None:
        ideal_splits = [max(0, min(int(r * len(elements)), len(elements) - 1)) for r in split_element_ratios]
    else:
        ideal_splits = []
        for i in range(1, num_parts):
            idx = _find_element_index_by_chars(element_counts, target_per_part * i)
            ideal_splits.append(idx)

    for i in range(1, len(ideal_splits)):
        if ideal_splits[i] <= ideal_splits[i - 1]:
            ideal_splits[i] = ideal_splits[i - 1] + 1
    for i in range(len(ideal_splits)):
        ideal_splits[i] = min(ideal_splits[i], len(elements) - 1)

    element_ratios = [idx / len(elements) for idx in ideal_splits] if elements else []

    split_ranges = []
    for part_idx in range(num_parts):
        if part_idx == 0:
            start = 0
            end = _find_buffer_end(element_counts, ideal_splits[0], buffer_chars, 'right') if ideal_splits else len(elements)
        elif part_idx == num_parts - 1:
            start = _find_buffer_end(element_counts, ideal_splits[-1], buffer_chars, 'left')
            end = len(elements)
        else:
            start = _find_buffer_end(element_counts, ideal_splits[part_idx - 1], buffer_chars, 'left')
            end = _find_buffer_end(element_counts, ideal_splits[part_idx], buffer_chars, 'right')
        start = max(0, min(start, len(elements) - 1))
        end = max(start + 1, min(end, len(elements)))
        split_ranges.append((start, end))

    generated_files = []
    part_info = []
    for i, (s, e) in enumerate(split_ranges):
        dest = os.path.join(output_dir, f"{base_name}_Part{i + 1}.docx")
        shutil.copy2(src_path, dest)
        doc_copy = Document(dest)
        total_elems = len(_get_all_content_elements(doc_copy))
        _delete_elements_in_range(doc_copy, e, total_elems + 5000)
        _delete_elements_in_range(doc_copy, 0, s)
        doc_copy.save(dest)
        first_text = _extract_text_from_elements(elements, s, min(s + 3, e))
        last_text = _extract_text_from_elements(elements, max(s, e - 3), e)
        part_info.append({
            'path': dest,
            'first_anchor': first_text[:200] if first_text else "",
            'last_anchor': last_text[-200:] if last_text else "",
        })
        generated_files.append(dest)

    return generated_files, part_info, element_ratios


# ── LLM 调用 ─────────────────────────────────────────────
def _call_llm(system_prompt: str, user_prompt: str, model_id: str, max_output: int = 65536) -> Optional[str]:
    route = _get_current_gemini_route()
    try:
        print(f"[alignment-llm] route={route}, model={model_id}, sys_len={len(system_prompt)}, user_len={len(user_prompt)}")
        full = generate_text(
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            model=model_id,
            route=route,
            temperature=0.1,
            max_output_tokens=max_output,
            timeout=600.0,
        )
        print(f"[alignment-llm] ??, ?? {len(full)} ??")
        return full
    except Exception as e:
        print(f"[alignment-llm] ????: {e}")
        import traceback as _tb
        _tb.print_exc()
        return None

def _get_model_max_output(model_id: str) -> int:
    for info in AVAILABLE_MODELS.values():
        if info.get("id") == model_id:
            return int(info.get("max_output", 65536))
    return 65536


# ── 解析对齐响应 ──────────────────────────────────────────
def _parse_alignment_response(response_text: str) -> list:
    if not response_text:
        return []
    cleaned = response_text.replace('\r\n', '\n').replace('\r', '\n')
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    lines = cleaned.splitlines()
    data = []
    pending = ""
    for line in lines:
        line = line.strip()
        if not line or line.startswith('```'):
            continue
        if pending:
            line = pending + " " + line
            pending = ""
        if "|||" in line:
            parts = line.split("|||", 1)
            original = parts[0].strip()
            trans = parts[1].strip() if len(parts) > 1 else ""
            if original or trans:
                data.append({"原文": original, "译文": trans})
        else:
            if len(line) < 50:
                pending = line
    return data


# ── 英文后处理分句 ────────────────────────────────────────
ENGLISH_ABBREVIATIONS = {
    'et al', 'etc', 'e.g', 'i.e', 'vs', 'cf', 'ibid', 'op', 'cit',
    'mr', 'mrs', 'ms', 'dr', 'prof', 'jr', 'sr', 'st',
    'inc', 'ltd', 'co', 'corp', 'llc', 'l.l.c',
    'u.s', 'u.k', 'u.n', 'e.u', 'u.s.a',
    'a.m', 'p.m',
    'no', 'vol', 'fig', 'ch', 'sec', 'pp', 'approx', 'est', 'max', 'min', 'avg',
}


def _is_abbreviation_period(text, pos):
    if pos <= 0:
        return False
    before = text[:pos]
    after = text[pos + 1:] if pos + 1 < len(text) else ""

    if after and after[0].isdigit():
        return True
    if after and after[0].isalpha():
        return True

    words = before.split()
    if not words:
        return False

    last_token = words[-1]
    last_word_lower = last_token.lower().rstrip('.,;:')

    if last_word_lower in ENGLISH_ABBREVIATIONS:
        return True
    if '.' in last_token:
        return True

    if len(last_word_lower) == 1 and last_word_lower.isalpha():
        if len(words) >= 2:
            prev_token = words[-2].lower()
            if re.match(r'^[a-z]\.$', prev_token):
                return True
        return False

    if last_word_lower.isdigit():
        return False

    roman_pattern = r'^(i{1,3}|iv|vi{0,3}|ix|xi{0,3}|xiv|xvi{0,3}|xix|xxi{0,3})$'
    if re.match(roman_pattern, last_word_lower):
        return False

    return False


def _count_real_sentences(text):
    if not text:
        return 0
    count = 0
    for i, ch in enumerate(text):
        if ch in '.!?':
            if ch == '.' and _is_abbreviation_period(text, i):
                continue
            count += 1
    return max(1, count)


def _needs_post_split(orig_text, source_lang):
    if source_lang in ["中文", "日语", "韩语"]:
        return False
    if not orig_text:
        return False
    numbering_patterns = [
        r'\d+\.\s',
        r'(?<![A-Za-z.])[A-Za-z]\.\s',
        r'\([0-9]+\)\s',
        r'\([A-Za-z]\)\s',
        r'(?<![A-Za-z])[IVXivx]+\.\s',
    ]
    for p in numbering_patterns:
        if re.search(p, orig_text):
            return True
    return _count_real_sentences(orig_text) > 1


def _split_row_with_ai(orig, trans, model_id, source_lang, row_num="后处理"):
    prompt = f"""请根据原文的标点分句，将译文对应拆分对齐：

【原文】
{orig}

【译文】
{trans}

输出分句对齐结果（以原文分句数量为准）："""
    resp = _call_llm(
        _get_split_row_prompt(source_lang),
        prompt,
        model_id,
        max_output=_get_model_max_output(model_id),
    )
    if resp:
        results = _parse_alignment_response(resp)
        if results:
            return results
    return None


def _post_process_split(data, model_id, source_lang):
    if source_lang in ["中文", "日语", "韩语"]:
        return data
    result = []
    for idx, row in enumerate(data):
        orig = row.get('原文', '')
        trans = row.get('译文', '')
        if _needs_post_split(orig, source_lang):
            split_results = _split_row_with_ai(orig, trans, model_id, source_lang, f"后处理-{idx + 1}")
            if split_results and len(split_results) > 1:
                result.extend(split_results)
            else:
                result.append(row)
        else:
            result.append(row)
    return result


# ── 质量检查 ──────────────────────────────────────────────
def _quality_check(df, source_lang, target_lang):
    issues = []
    source_info = SUPPORTED_LANGUAGES.get(source_lang, SUPPORTED_LANGUAGES["中文"])
    target_info = SUPPORTED_LANGUAGES.get(target_lang, SUPPORTED_LANGUAGES["英语"])
    source_pattern = source_info['char_pattern']
    target_pattern = target_info['char_pattern']

    for idx, row in df.iterrows():
        orig = str(row.get('原文', ''))
        trans = str(row.get('译文', ''))
        if not orig or not trans:
            continue

        source_chars_in_original = len(re.findall(source_pattern, orig))
        target_chars_in_original = len(re.findall(target_pattern, orig))
        source_chars_in_trans = len(re.findall(source_pattern, trans))
        target_chars_in_trans = len(re.findall(target_pattern, trans))

        total_original = len(orig) if orig else 1
        total_trans = len(trans) if trans else 1

        source_ratio_original = source_chars_in_original / total_original
        target_ratio_trans = target_chars_in_trans / total_trans
        source_ratio_trans = source_chars_in_trans / total_trans
        target_ratio_original = target_chars_in_original / total_original

        if source_ratio_original < 0.2 and source_ratio_trans > 0.4:
            issues.append(f"行{idx + 1}: 语言错位（原文疑似{target_lang}，译文疑似{source_lang}）")
        elif target_ratio_original > 0.4 and target_ratio_trans < 0.2:
            issues.append(f"行{idx + 1}: 语言错位（原文疑似{target_lang}，译文疑似{source_lang}）")

        lo, lt = len(orig), len(trans)
        if lo > 0 and lt > 0:
            ratio = max(lo, lt) / min(lo, lt)
            if ratio > 5:
                issues.append(f"行{idx + 1}: 长度比 {ratio:.1f}:1")
    return issues


# ── Excel 双文件对齐 ──────────────────────────────────────
def _needs_table_cell_split(text, source_lang):
    if not text or not text.strip():
        return False
    import re

    source_is_cjk = source_lang in ["中文", "日语", "韩语"]
    has_newline = '\n' in text or '\r' in text
    has_consecutive_spaces = '  ' in text

    if source_is_cjk:
        period_count = text.count('。')
        exclamation_count = text.count('！')
        question_count = text.count('？')
    else:
        period_count = text.count('.')
        exclamation_count = text.count('!')
        question_count = text.count('?')

    total_punct_count = period_count + exclamation_count + question_count

    has_numbered_list = bool(re.search(
        r'（[一二三四五六七八九十\d]+）|'
        r'[（\(]\d+[）\)]|'
        r'[①②③④⑤⑥⑦⑧⑨⑩]|'
        r'\d+\.(?!\d)\s*\S|'
        r'\d+、\s*\S|'
        r'\d+[）\)]\s*\S|'
        r'^[一二三四五六七八九十]+[、．.]',
        text
    ))

    return (
        total_punct_count > 1 or
        exclamation_count > 0 or
        question_count > 0 or
        has_newline or
        has_consecutive_spaces or
        has_numbered_list
    )


def _split_table_cell_with_ai(orig, trans, model_id, source_lang):
    prompt = f"""请根据原文的句末标点（。！？）和换行符进行分句，将译文对应拆分对齐：

【原文】
{orig}

【译文】
{trans}

输出分句对齐结果："""
    resp = _call_llm(
        _get_table_cell_split_prompt(source_lang),
        prompt,
        model_id,
        max_output=_get_model_max_output(model_id),
    )
    if resp:
        results = _parse_alignment_response(resp)
        if results:
            return results
    return None


def _process_excel_dual(orig_path, trans_path, output_path, model_id, source_lang, target_lang, task_id):
    orig_sheets = _read_excel_all_sheets(orig_path)
    trans_sheets = _read_excel_all_sheets(trans_path)
    if not orig_sheets or not trans_sheets:
        return False

    common = set(orig_sheets.keys()) & set(trans_sheets.keys())
    if common:
        pairs = [(n, n) for n in sorted(common)]
    else:
        pairs = list(zip(orig_sheets.keys(), trans_sheets.keys()))

    all_results = []
    total_pairs = len(pairs)

    for si, (orig_name, trans_name) in enumerate(pairs):
        _update_progress(task_id, 20 + int(60 * si / total_pairs),
                         f"处理工作簿 {si + 1}/{total_pairs}: {orig_name}")
        df_o = orig_sheets[orig_name]
        df_t = trans_sheets[trans_name]
        max_rows = max(df_o.shape[0], df_t.shape[0])
        max_cols = max(df_o.shape[1], df_t.shape[1])

        for r in range(max_rows):
            for c in range(max_cols):
                orig_text = ""
                if r < df_o.shape[0] and c < df_o.shape[1]:
                    v = df_o.iloc[r, c]
                    if pd.notna(v):
                        orig_text = str(v).strip()
                trans_text = ""
                if r < df_t.shape[0] and c < df_t.shape[1]:
                    v = df_t.iloc[r, c]
                    if pd.notna(v):
                        trans_text = str(v).strip()
                if not orig_text and not trans_text:
                    continue
                if orig_text and trans_text and _needs_table_cell_split(orig_text, source_lang):
                    split = _split_table_cell_with_ai(orig_text, trans_text, model_id, source_lang)
                    if split:
                        all_results.extend(split)
                    else:
                        all_results.append({"原文": orig_text, "译文": trans_text})
                else:
                    if orig_text or trans_text:
                        all_results.append({"原文": orig_text, "译文": trans_text})

    if not all_results:
        return False
    pd.DataFrame(all_results).to_excel(output_path, index=False)
    return True


# ── 合并去重 ─────────────────────────────────────────────
def _merge_excels(paths, final_path, source_lang, target_lang):
    dfs = []
    for p in paths:
        if os.path.exists(p):
            try:
                df = pd.read_excel(p)
                if not df.empty:
                    dfs.append(df)
            except Exception:
                pass
    if not dfs:
        return None
    combined = pd.concat(dfs, ignore_index=True)
    combined.drop_duplicates(subset=['原文', '译文'], keep='first', inplace=True)
    combined.to_excel(final_path, index=False)

    try:
        from openpyxl import load_workbook
        from openpyxl.styles import PatternFill
        wb = load_workbook(final_path)
        ws = wb.active
        yellow = PatternFill(start_color='FFFF00', end_color='FFFF00', fill_type='solid')
        orig_dup = combined.duplicated(subset=['原文'], keep=False)
        trans_dup = combined.duplicated(subset=['译文'], keep=False)
        for idx in combined[orig_dup | trans_dup].index:
            for col in range(1, 3):
                ws.cell(row=idx + 2, column=col).fill = yellow
        wb.save(final_path)
    except Exception:
        pass

    return final_path


# ── 单文件对齐 ───────────────────────────────────────────
def _run_single_alignment(orig_path, trans_path, output_path, model_id,
                          source_lang, target_lang, enable_post_split,
                          system_prompt_override=None,
                          anchor_orig=None, anchor_trans=None,
                          max_output=65536):
    print(f"[alignment] _run_single_alignment: orig={orig_path}")
    print(f"[alignment]   orig exists={os.path.exists(orig_path)}, trans exists={os.path.exists(trans_path)}")

    text_original = _read_file_content(orig_path)
    text_trans = _read_file_content(trans_path)

    print(f"[alignment]   原文长度={len(text_original)}, 译文长度={len(text_trans)}")

    if not text_original or not text_trans:
        print(f"[alignment]   内容为空，跳过")
        return False

    if system_prompt_override:
        sys_prompt = system_prompt_override
    else:
        sys_prompt = _get_docx_alignment_prompt(source_lang, target_lang)

    anchor_hint = ""
    if anchor_orig and anchor_trans:
        anchor_hint = f"""
## 上下文提示
- 原文开头: "{anchor_orig.get('first_anchor', '')[:100]}"
- 译文开头: "{anchor_trans.get('first_anchor', '')[:100]}"
"""

    user_prompt = f"""{anchor_hint}

<Stream_A_Original>
{text_original}
</Stream_A_Original>

<Stream_B_Translation>
{text_trans}
</Stream_B_Translation>

请严格按规则输出对齐结果：
"""

    print(f"[alignment]   调用 LLM (model={model_id})...")
    try:
        response = _call_llm(sys_prompt, user_prompt, model_id, max_output=max_output)
    except Exception as e:
        print(f"[alignment]   LLM 调用异常: {e}")
        traceback.print_exc()
        return False
    print(f"[alignment]   LLM 返回长度={len(response) if response else 0}")

    if not response:
        print(f"[alignment]   LLM 返回为空 (None 或空字符串)")
        return False

    print(f"[alignment]   LLM 前300字: {response[:300]}")
    last_500 = response[-500:] if len(response) > 500 else response
    if '|||' not in last_500:
        print("[alignment]   警告: 输出可能被截断（最后500字符无分隔符）")

    data = _parse_alignment_response(response)
    print(f"[alignment]   解析得到 {len(data)} 行")

    if not data:
        print(f"[alignment]   解析结果为空，返回 False")
        print(f"[alignment]   LLM 原始响应前500字: {response[:500]}")
        return False

    if enable_post_split and source_lang not in ["中文", "日语", "韩语"]:
        data = _post_process_split(data, model_id, source_lang)

    df = pd.DataFrame(data)
    df.to_excel(output_path, index=False)
    print(f"[alignment]   已保存: {output_path} ({len(df)} 行)")
    return True


# ── 列出中间文件 ──
def _list_intermediate_files(temp_dir: str) -> list:
    """扫描中间文件目录，返回 [{name, path, type}] 列表"""
    files = []
    if not os.path.isdir(temp_dir):
        return files
    for f in sorted(os.listdir(temp_dir)):
        full = os.path.join(temp_dir, f)
        if os.path.isfile(full):
            rel = os.path.relpath(full, ".").replace("\\", "/")
            ext = os.path.splitext(f)[1].lower()
            ftype = "excel" if ext in (".xlsx", ".xls") else "word" if ext in (".docx", ".doc") else "other"
            files.append({"name": f, "path": rel, "type": ftype})
    return files


# ── 同步主处理（1:1 复刻 memory.py GUI 的 run_processing）─────
def _run_alignment_sync(
    original_path: str,
    translated_path: str,
    task_id: str,
    display_no: Optional[str],
    source_lang: str,
    target_lang: str,
    model_name: str,
    gemini_route: str,
    enable_post_split: bool,
    threshold_2: int = 25000,
    threshold_3: int = 50000,
    threshold_4: int = 75000,
    threshold_5: int = 100000,
    threshold_6: int = 125000,
    threshold_7: int = 150000,
    threshold_8: int = 175000,
    buffer_chars: int = 2000,
):
    """同步执行对齐任务 - 线程安全版，支持多任务并发"""
    try:
        _update_progress(task_id, 5, "正在加载处理引擎...", stream_log="")
        gemini_route = ensure_gemini_route_configured(gemini_route)
        _set_current_gemini_route(gemini_route)

        memory_module = _get_memory_module()
        _install_log_patches()

        # 设置当前线程的 task_id，永久日志补丁通过 threading.local 自动路由
        _stream_task_id.task_id = task_id
        _log = memory_module.log_manager.log

        # 绝对路径
        original_path = os.path.abspath(original_path)
        translated_path = os.path.abspath(translated_path)

        model_info = AVAILABLE_MODELS.get(model_name, AVAILABLE_MODELS[DEFAULT_MODEL])
        model_id = model_info['id']

        file_type = memory_module.get_file_type(original_path)
        base_name = os.path.splitext(os.path.basename(original_path))[0]
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        task_dir = os.path.join(OUTPUT_DIR, "alignment", display_no or f"{base_name}_{timestamp}")
        temp_dir = os.path.join(task_dir, "中间文件")
        os.makedirs(task_dir, exist_ok=True)
        os.makedirs(temp_dir, exist_ok=True)

        print(f"[alignment] task_id={task_id}")
        print(f"[alignment] original={original_path}, exists={os.path.exists(original_path)}, size={os.path.getsize(original_path) if os.path.exists(original_path) else 'N/A'}")
        print(f"[alignment] translated={translated_path}, exists={os.path.exists(translated_path)}, size={os.path.getsize(translated_path) if os.path.exists(translated_path) else 'N/A'}")
        print(f"[alignment] file_type={file_type}, model={model_name} ({model_id})")
        print(f"[alignment] lang: {source_lang} → {target_lang}")
        _log(f"语言对: {source_lang} → {target_lang}")
        _log(f"后处理分句: {'启用' if enable_post_split else '禁用'}")
        _log(f"Gemini ???: {gemini_route}")
        _log(f"模型: {model_name} ({model_id})")
        _log(f"文件类型: {file_type.upper()}")

        _update_progress(task_id, 8, "正在分析文档...")

        # .doc → .docx 转换
        if file_type == 'doc':
            _update_progress(task_id, 8, "检测到 .doc 文件，正在转换...")
            converted_orig = memory_module.convert_doc_to_docx(original_path, temp_dir)
            converted_trans = memory_module.convert_doc_to_docx(translated_path, temp_dir)
            if converted_orig is None or converted_trans is None:
                _complete_task(task_id, error="无法转换 .doc 文件，请安装 pywin32")
                return
            original_path = converted_orig
            translated_path = converted_trans
            file_type = 'docx'

        # ── Excel 双文件模式 ──
        if file_type == 'excel':
            _update_progress(task_id, 10, "Excel 双文件对齐模式...")
            out_path = os.path.join(task_dir, f"{base_name}_对齐结果.xlsx")
            memory_module.set_gemini_route(gemini_route)
            success = memory_module.process_excel_dual_file_alignment(
                original_path, translated_path, out_path, model_id,
                source_lang=source_lang, target_lang=target_lang,
            )
            if success and os.path.exists(out_path):
                rel = os.path.relpath(out_path, ".").replace("\\", "/")
                _complete_task(task_id, result={
                    "output_excel": rel,
                    "row_count": len(pd.read_excel(out_path)),
                    "file_type": "excel",
                "gemini_route": gemini_route,
                    "intermediate_files": _list_intermediate_files(temp_dir),
                })
            else:
                _complete_task(task_id, error="Excel 对齐处理失败")
            return

        # ── Word / PPT ──
        _update_progress(task_id, 10, "分析文档结构...")
        count_a, _ = memory_module.analyze_document_structure(original_path, source_lang)
        count_b, _ = memory_module.analyze_document_structure(translated_path, target_lang)

        s_info = SUPPORTED_LANGUAGES.get(source_lang, {})
        t_info = SUPPORTED_LANGUAGES.get(target_lang, {})
        s_unit = "字" if not s_info.get('word_based', True) else "词"
        t_unit = "字" if not t_info.get('word_based', True) else "词"

        print(f"[alignment] 原文: {count_a:,} {s_unit}, 译文: {count_b:,} {t_unit}")
        _log(f"原文: {count_a:,} {s_unit}")
        _log(f"译文: {count_b:,} {t_unit}")
        _update_progress(task_id, 15, f"原文 {count_a:,} {s_unit}，译文 {count_b:,} {t_unit}")

        # 分割策略（与 GUI run_processing 1:1）
        if file_type == 'pptx':
            split_parts = 1
            _log("PPT文件：不进行分割")
        else:
            max_count = max(count_a, count_b)
            split_parts = 1
            if max_count > threshold_8:
                split_parts = 8
            elif max_count > threshold_7:
                split_parts = 7
            elif max_count > threshold_6:
                split_parts = 6
            elif max_count > threshold_5:
                split_parts = 5
            elif max_count > threshold_4:
                split_parts = 4
            elif max_count > threshold_3:
                split_parts = 3
            elif max_count > threshold_2:
                split_parts = 2
            _log(f"分割策略: {split_parts} 份")

        _update_progress(task_id, 20, f"分割策略: {split_parts} 份")

        tasks_queue = []
        generated_excel_paths = []

        if split_parts > 1 and file_type == 'docx':
            _update_progress(task_id, 25, f"正在分割文档（{split_parts} 份，缓冲区 {buffer_chars} 字）...")
            _log(f"分割原文（主文档，自主计算分割点）...")
            files_a, part_info_a, split_ratios = memory_module.smart_split_with_buffer(
                original_path, split_parts, temp_dir, source_lang, buffer_chars)
            _log(f"分割译文（从文档，使用原文的分割比例）...")
            files_b, part_info_b, _ = memory_module.smart_split_with_buffer(
                translated_path, split_parts, temp_dir, target_lang, buffer_chars,
                split_element_ratios=split_ratios)

            for i in range(len(files_a)):
                out = os.path.join(temp_dir, f"Part{i + 1}_对齐结果.xlsx")
                tasks_queue.append({
                    'original': files_a[i], 'trans': files_b[i], 'output': out,
                    'anchor_orig': part_info_a[i] if part_info_a else None,
                    'anchor_trans': part_info_b[i] if part_info_b else None,
                })
                generated_excel_paths.append(out)
        else:
            out = os.path.join(task_dir, f"{base_name}_对齐结果.xlsx")
            ppt_prompt = memory_module.get_ppt_alignment_prompt(source_lang, target_lang) if file_type == 'pptx' else None
            tasks_queue.append({
                'original': original_path, 'trans': translated_path, 'output': out,
                'anchor_orig': None, 'anchor_trans': None,
                'system_prompt_override': ppt_prompt,
            })
            generated_excel_paths.append(out)

        # ── AI 对齐 ──
        _update_progress(task_id, 30, f"AI 对齐中（共 {len(tasks_queue)} 个任务）...")
        _log(f"待处理任务数: {len(tasks_queue)}")
        progress_per_task = 50 / len(tasks_queue) if tasks_queue else 50

        for idx, task in enumerate(tasks_queue):
            progress = 30 + int((idx + 1) * progress_per_task)
            _update_progress(task_id, progress,
                             f"AI 对齐中 ({idx + 1}/{len(tasks_queue)})...")

            _log(f"处理任务 {idx + 1}/{len(tasks_queue)}: {os.path.basename(task['output'])}")
            print(f"[alignment] 处理任务 {idx + 1}/{len(tasks_queue)}")
            print(f"[alignment]   original: {task['original']}")
            print(f"[alignment]   trans: {task['trans']}")
            print(f"[alignment]   output: {task['output']}")

            # 诊断：直接用 memory 的 read_file_content 测试读取
            try:
                test_orig = memory_module.read_file_content(task['original'])
                test_trans = memory_module.read_file_content(task['trans'])
                print(f"[alignment]   read_file_content 原文: {len(test_orig)} 字符")
                print(f"[alignment]   read_file_content 译文: {len(test_trans)} 字符")
                if test_orig:
                    print(f"[alignment]   原文前200字: {test_orig[:200]}")
                if test_trans:
                    print(f"[alignment]   译文前200字: {test_trans[:200]}")
            except Exception as diag_e:
                print(f"[alignment]   read_file_content 诊断异常: {diag_e}")

            memory_module.set_gemini_route(gemini_route)
            success = memory_module.run_llm_alignment(
                task['original'],
                task['trans'],
                task['output'],
                model_id,
                anchor_info_orig=task.get('anchor_orig'),
                anchor_info_trans=task.get('anchor_trans'),
                system_prompt_override=task.get('system_prompt_override'),
                source_lang=source_lang,
                target_lang=target_lang,
                enable_post_split=enable_post_split,
            )

            print(f"[alignment] 任务 {idx + 1} 结果: success={success}, output_exists={os.path.exists(task['output'])}")
            if success and os.path.exists(task['output']):
                _log(f"任务 {idx + 1} 成功: {os.path.basename(task['output'])}")
            else:
                _log(f"任务 {idx + 1} 失败: {os.path.basename(task['output'])}")
                if task['output'] in generated_excel_paths:
                    generated_excel_paths.remove(task['output'])

        # ── 合并 ──
        _update_progress(task_id, 85, "合并与去重...")
        final_path = None
        if split_parts > 1 and len(generated_excel_paths) > 0:
            final_path = os.path.join(task_dir, f"「最终结果」{base_name}_对齐.xlsx")
            memory_module.merge_and_deduplicate_excels(
                generated_excel_paths, final_path,
                source_lang=source_lang, target_lang=target_lang,
            )
        else:
            final_path = generated_excel_paths[0] if generated_excel_paths else None

        print(f"[alignment] final_path={final_path}, exists={os.path.exists(final_path) if final_path else 'N/A'}")
        print(f"[alignment] generated_excel_paths={generated_excel_paths}")
        if final_path and os.path.exists(final_path):
            _update_progress(task_id, 90, "正在整理最终结果...")
            rel = os.path.relpath(final_path, ".").replace("\\", "/")
            final_df = pd.read_excel(final_path)
            row_count = len(final_df)
            if split_parts > 1:
                issues = []
            else:
                _update_progress(task_id, 95, "正在执行最终质量检查...")
                issues = _quality_check(final_df, source_lang, target_lang)
            print(f"[alignment] 准备完成任务: row_count={row_count}, split_parts={split_parts}, path={rel}")
            _log(f"准备写入任务结果，{row_count} 行, 路径: {rel}")
            _complete_task(task_id, result={
                "output_excel": rel,
                "row_count": row_count,
                "file_type": file_type,
                "gemini_route": gemini_route,
                "split_parts": split_parts,
                "issues": issues[:10] if issues else [],
                "intermediate_files": _list_intermediate_files(temp_dir),
            })
            print(f"[alignment] task completed: {task_id}")
        else:
            print(f"[alignment] 失败: 无有效输出文件")
            _complete_task(task_id, error="对齐处理失败，未生成结果文件。请查看实时输出了解详情")

    except Exception as e:
        tb = traceback.format_exc()
        print(f"[alignment] 异常:\n{tb}")
        _complete_task(task_id, error=str(e))
    finally:
        # 仅清理线程局部变量，永久补丁无需恢复
        if hasattr(_stream_task_id, "task_id"):
            del _stream_task_id.task_id


# ── 异步入口（供 BackgroundTasks 调用）────────────────────
async def run_alignment_task(
    original_path: str,
    translated_path: str,
    task_id: str,
    display_no: Optional[str] = None,
    source_lang: str = "中文",
    target_lang: str = "英语",
    model_name: str = DEFAULT_MODEL,
    gemini_route: str = settings.GEMINI_DEFAULT_ROUTE,
    enable_post_split: bool = True,
    threshold_2: int = 25000,
    threshold_3: int = 50000,
    threshold_4: int = 75000,
    threshold_5: int = 100000,
    threshold_6: int = 125000,
    threshold_7: int = 150000,
    threshold_8: int = 175000,
    buffer_chars: int = 2000,
    executor: Optional[Executor] = None,
):
    """在后台线程池中执行对齐任务"""
    import functools
    loop = asyncio.get_running_loop()
    await loop.run_in_executor(
        executor,
        functools.partial(
            _run_alignment_sync,
            original_path, translated_path, task_id, display_no,
            source_lang, target_lang, model_name, gemini_route, enable_post_split,
            threshold_2=threshold_2, threshold_3=threshold_3,
            threshold_4=threshold_4, threshold_5=threshold_5,
            threshold_6=threshold_6, threshold_7=threshold_7,
            threshold_8=threshold_8, buffer_chars=buffer_chars,
        ),
    )
