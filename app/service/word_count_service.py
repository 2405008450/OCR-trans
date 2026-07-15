# -*- coding: utf-8 -*-
from __future__ import annotations

import json
import os
import re
import unicodedata
from collections import Counter, defaultdict
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Awaitable, Callable, Iterable, Optional
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter

from app.core.config import settings
from app.service.cad_text_service import (
    CadConverterUnavailableError,
    extract_cad_text,
    get_cad_support_info,
)
from app.service.libreoffice_service import (
    convert_doc_to_docx_via_libreoffice,
    convert_presentation_to_pptx_via_libreoffice,
    convert_spreadsheet_to_xlsx_via_libreoffice,
)
from app.service.ocr_text_service import extract_ocr_plain_text
from app.service.pdf2docx_service import (
    PDF2DOCX_DEFAULT_GEMINI_ROUTE,
    PDF2DOCX_DEFAULT_MODEL,
    get_pdf2docx_models,
    normalize_pdf2docx_model,
)

ProgressCallback = Callable[[int, str], Awaitable[None]]

WORD_COUNT_TASK_TYPE = "word_count"
COUNTABLE_EXTENSIONS = {".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".pdf", ".txt"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff"}
CAD_EXTENSIONS = {".dwg", ".dxf", ".dwt", ".dws"}
CANDIDATE_EXTENSIONS = COUNTABLE_EXTENSIONS | IMAGE_EXTENSIONS | CAD_EXTENSIONS
DEFAULT_SCAN_EXTENSIONS = sorted(CANDIDATE_EXTENSIONS)

STATUS_COUNTED = "counted"
STATUS_NEEDS_OCR = "needs_ocr"
STATUS_NEEDS_CAD_PARSER = "needs_cad_parser"
STATUS_SKIPPED = "skipped"
STATUS_FAILED = "failed"

OCR_MODE_AUTO = "auto"
OCR_MODE_ON = "on"
OCR_MODE_OFF = "off"
OCR_MODES = {
    OCR_MODE_AUTO: {"label": "智能", "description": "单文件自动识别，目录任务不自动识别"},
    OCR_MODE_ON: {"label": "开启", "description": "识别扫描 PDF 和独立图片"},
    OCR_MODE_OFF: {"label": "关闭", "description": "仅标记需要 OCR 的文件"},
}
PDF_OCR_SPARSE_WORD_THRESHOLD = 20
PDF_OCR_IMAGE_COVERAGE_THRESHOLD = 0.5

EXTRA_SOURCE_TYPES = {"header", "footer", "footnote", "endnote", "chart", "comment"}

SCRIPT_COUNT_FIELDS = (
    "han_count",
    "kana_count",
    "hangul_count",
    "latin_word_count",
    "number_token_count",
    "mixed_latin_number_count",
    "cyrillic_word_count",
    "arabic_word_count",
    "greek_word_count",
    "hebrew_word_count",
    "thai_word_count",
    "cjk_punct_count",
    "other_count",
)

QUOTE_COUNT_FIELDS = (
    "billable_chinese_count",
    "billable_japanese_count",
    "billable_korean_count",
    "billable_latin_count",
    "billable_cyrillic_count",
    "billable_arabic_count",
    "billable_greek_count",
    "billable_hebrew_count",
    "billable_thai_count",
)

SCRIPT_COUNT_LABELS = {
    "han_count": "汉字",
    "kana_count": "假名",
    "hangul_count": "韩文",
    "latin_word_count": "拉丁词",
    "number_token_count": "纯数字",
    "mixed_latin_number_count": "拉丁数字混合词",
    "cyrillic_word_count": "西里尔词",
    "arabic_word_count": "阿拉伯字母词",
    "greek_word_count": "希腊字母词",
    "hebrew_word_count": "希伯来字母词",
    "thai_word_count": "泰文词",
    "cjk_punct_count": "CJK 标点",
    "other_count": "其他",
}

QUOTE_COUNT_LABELS = {
    "billable_chinese_count": "中文候选",
    "billable_japanese_count": "日语候选",
    "billable_korean_count": "韩语候选",
    "billable_latin_count": "拉丁系候选",
    "billable_cyrillic_count": "西里尔候选",
    "billable_arabic_count": "阿语候选",
    "billable_greek_count": "希腊语候选",
    "billable_hebrew_count": "希伯来语候选",
    "billable_thai_count": "泰语候选",
}

NS = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "wps": "http://schemas.microsoft.com/office/word/2010/wordprocessingShape",
    "v": "urn:schemas-microsoft-com:vml",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
REL_CHART = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart"


@dataclass(frozen=True)
class TextMetrics:
    word_count: int
    non_space_chars: int
    raw_chars: int
    han_count: int = 0
    kana_count: int = 0
    hangul_count: int = 0
    latin_word_count: int = 0
    number_token_count: int = 0
    mixed_latin_number_count: int = 0
    cyrillic_word_count: int = 0
    arabic_word_count: int = 0
    greek_word_count: int = 0
    hebrew_word_count: int = 0
    thai_word_count: int = 0
    cjk_punct_count: int = 0
    other_count: int = 0

    @property
    def script_count_total(self) -> int:
        return sum(int(getattr(self, field, 0)) for field in SCRIPT_COUNT_FIELDS)

    @property
    def billable_chinese_count(self) -> int:
        return self.han_count + self.cjk_punct_count

    @property
    def billable_japanese_count(self) -> int:
        return self.han_count + self.kana_count + self.cjk_punct_count

    @property
    def billable_korean_count(self) -> int:
        return self.hangul_count + self.cjk_punct_count

    @property
    def billable_latin_count(self) -> int:
        return self.latin_word_count + self.mixed_latin_number_count

    @property
    def billable_cyrillic_count(self) -> int:
        return self.cyrillic_word_count

    @property
    def billable_arabic_count(self) -> int:
        return self.arabic_word_count

    @property
    def billable_greek_count(self) -> int:
        return self.greek_word_count

    @property
    def billable_hebrew_count(self) -> int:
        return self.hebrew_word_count

    @property
    def billable_thai_count(self) -> int:
        return self.thai_word_count

    def script_counts(self) -> dict[str, int]:
        return {field: int(getattr(self, field, 0)) for field in SCRIPT_COUNT_FIELDS}

    def quote_counts(self) -> dict[str, int]:
        return {field: int(getattr(self, field, 0)) for field in QUOTE_COUNT_FIELDS}


@dataclass(frozen=True)
class TextItem:
    source_type: str
    source_label: str
    text: str
    is_extra: bool = False
    paragraph_count: int = 1
    line_count: int = 1
    image_count: int = 0


@dataclass(frozen=True)
class ExtractedContent:
    items: list[TextItem]
    page_count: int = 0
    paragraph_count: int = 0
    line_count: int = 0
    image_count: int = 0
    stat_method: str = ""
    file_type: str = ""
    warning: str = ""


def normalize_scan_extensions(extensions: Optional[Iterable[str]]) -> list[str]:
    if not extensions:
        return DEFAULT_SCAN_EXTENSIONS
    normalized: set[str] = set()
    for item in extensions:
        value = str(item or "").strip().lower()
        if not value:
            continue
        if not value.startswith("."):
            value = f".{value}"
        if value in CANDIDATE_EXTENSIONS:
            normalized.add(value)
    return sorted(normalized or CANDIDATE_EXTENSIONS)


def _normalize_ocr_mode(value: Optional[str]) -> str:
    mode = str(value or OCR_MODE_AUTO).strip().lower()
    if mode not in OCR_MODES:
        raise ValueError(f"不支持的 OCR 模式: {mode}")
    return mode


def _ocr_enabled_for_input(ocr_mode: str, input_kind: str) -> bool:
    return ocr_mode == OCR_MODE_ON or (ocr_mode == OCR_MODE_AUTO and input_kind == "file")


def _safe_path_exists(path: Path) -> bool:
    """检查路径是否存在；网络路径无权限或暂时不可达时返回 False。"""
    try:
        return path.exists()
    except OSError:
        return False


def get_word_count_config() -> dict[str, Any]:
    cad_support = get_cad_support_info(settings.ODA_FILE_CONVERTER_PATH)
    runtime_countable_extensions = COUNTABLE_EXTENSIONS | set(cad_support["supported_extensions"])
    allowed_roots = []
    for raw_root in settings.WORD_COUNT_ALLOWED_ROOTS:
        root = Path(raw_root).expanduser()
        mapped_root = _map_unc_path_to_local(str(raw_root))
        mapped_path = mapped_root[0] if mapped_root is not None else None
        root_exists = _safe_path_exists(root)
        scope_only = _is_unc_server_root(root) and not root_exists
        exists = _safe_path_exists(mapped_path) if mapped_path is not None else root_exists or scope_only
        allowed_roots.append(
            {
                "path": str(root),
                "exists": exists,
                "scope_only": scope_only,
                "mount_path": str(mapped_path) if mapped_path is not None else "",
            }
        )
    unc_mount_mappings = [
        {
            "unc": _normalize_unc_text(unc_path) or unc_path,
            "mount": mount_path,
        }
        for unc_path, mount_path in settings.WORD_COUNT_UNC_MOUNT_MAP.items()
    ]
    if settings.WORD_COUNT_ALLOW_LOCAL_PATHS_ENABLED:
        allowed_roots.append(
            {
                "path": "本地路径（测试放开：所有本机盘符和挂载路径）",
                "exists": True,
            }
        )
    return {
        "allowed_roots": allowed_roots,
        "countable_extensions": sorted(runtime_countable_extensions),
        "image_extensions": sorted(IMAGE_EXTENSIONS),
        "cad_extensions": sorted(CAD_EXTENSIONS),
        "cad_support": cad_support,
        "default_extensions": DEFAULT_SCAN_EXTENSIONS,
        "max_files": settings.WORD_COUNT_MAX_FILES,
        "max_file_mb": settings.WORD_COUNT_MAX_FILE_MB,
        "upload_max_file_mb": settings.WORD_COUNT_UPLOAD_MAX_MB,
        "upload_extensions": sorted(CANDIDATE_EXTENSIONS),
        "unc_mount_mappings": unc_mount_mappings,
        "unc_auto_mount_roots": settings.WORD_COUNT_UNC_AUTO_MOUNT_ROOTS,
        "follow_symlinks": settings.WORD_COUNT_FOLLOW_SYMLINKS_ENABLED,
        "allow_local_paths": settings.WORD_COUNT_ALLOW_LOCAL_PATHS_ENABLED,
        "count_policy": (
            "Word 近似口径：汉字/假名/韩文逐字计数，拉丁/数字等按连续词计数，"
            "部分 CJK 标点计数，空白不计；新增脚本分桶为互斥口径，脚本桶合计等于总字数。"
        ),
        "script_count_labels": SCRIPT_COUNT_LABELS,
        "quote_count_labels": QUOTE_COUNT_LABELS,
        "ocr_modes": OCR_MODES,
        "default_ocr_mode": OCR_MODE_AUTO,
        "ocr_models": get_pdf2docx_models(),
        "default_ocr_model": PDF2DOCX_DEFAULT_MODEL,
    }


def prepare_word_count_request(
    *,
    directory_path: str,
    recursive: bool = True,
    include_hidden: bool = False,
    extensions: Optional[Iterable[str]] = None,
    ocr_mode: str = OCR_MODE_AUTO,
    ocr_model: Optional[str] = None,
) -> dict[str, Any]:
    input_path, matched_root, input_kind = _resolve_allowed_input_path(directory_path)
    normalized_extensions = normalize_scan_extensions(extensions)
    normalized_ocr_mode = _normalize_ocr_mode(ocr_mode)
    normalized_ocr_model = normalize_pdf2docx_model(ocr_model)
    if normalized_ocr_model not in get_pdf2docx_models():
        raise ValueError(f"不支持的 OCR 模型: {normalized_ocr_model}")
    ocr_enabled = _ocr_enabled_for_input(normalized_ocr_mode, input_kind)
    params = {
        "directory_path": str(input_path),
        "input_source": "path",
        "input_kind": input_kind,
        "recursive": bool(recursive),
        "include_hidden": bool(include_hidden),
        "extensions": normalized_extensions,
        "max_files": settings.WORD_COUNT_MAX_FILES,
        "max_file_mb": settings.WORD_COUNT_MAX_FILE_MB,
        "follow_symlinks": settings.WORD_COUNT_FOLLOW_SYMLINKS_ENABLED,
        "ocr_mode": normalized_ocr_mode,
        "ocr_enabled": ocr_enabled,
        "ocr_model": normalized_ocr_model,
        "ocr_route": PDF2DOCX_DEFAULT_GEMINI_ROUTE,
    }
    return {
        "filename": str(input_path),
        "params": params,
        "input_files": {
            "directory_path": str(input_path),
            "input_kind": input_kind,
            "allowed_root": str(matched_root),
        },
    }


def prepare_word_count_upload_request(
    *,
    filename: str,
    ocr_mode: str = OCR_MODE_AUTO,
    ocr_model: Optional[str] = None,
) -> dict[str, Any]:
    original_filename = str(filename or "").replace("\\", "/").rsplit("/", 1)[-1].strip()
    if not original_filename:
        raise ValueError("请选择需要统计的文件")
    extension = Path(original_filename).suffix.lower()
    if extension not in CANDIDATE_EXTENSIONS:
        supported = "、".join(sorted(CANDIDATE_EXTENSIONS))
        raise ValueError(f"不支持的文件格式: {extension or '无扩展名'}。支持: {supported}")

    normalized_ocr_mode = _normalize_ocr_mode(ocr_mode)
    normalized_ocr_model = normalize_pdf2docx_model(ocr_model)
    if normalized_ocr_model not in get_pdf2docx_models():
        raise ValueError(f"不支持的 OCR 模型: {normalized_ocr_model}")

    return {
        "filename": original_filename,
        "extension": extension,
        "params": {
            "input_source": "upload",
            "input_kind": "file",
            "recursive": False,
            "include_hidden": False,
            "extensions": [extension],
            "max_files": 1,
            "max_file_mb": settings.WORD_COUNT_MAX_FILE_MB,
            "upload_max_file_mb": settings.WORD_COUNT_UPLOAD_MAX_MB,
            "follow_symlinks": False,
            "ocr_mode": normalized_ocr_mode,
            "ocr_enabled": _ocr_enabled_for_input(normalized_ocr_mode, "file"),
            "ocr_model": normalized_ocr_model,
            "ocr_route": PDF2DOCX_DEFAULT_GEMINI_ROUTE,
        },
    }


async def execute_word_count_task(
    *,
    task_id: str,
    display_no: Optional[str],
    directory_path: str,
    recursive: bool,
    include_hidden: bool,
    extensions: Iterable[str],
    ocr_mode: str = OCR_MODE_AUTO,
    ocr_model: str = PDF2DOCX_DEFAULT_MODEL,
    ocr_route: str = PDF2DOCX_DEFAULT_GEMINI_ROUTE,
    input_source: str = "path",
    original_filename: Optional[str] = None,
    progress_callback: Optional[ProgressCallback] = None,
    executor=None,
) -> dict[str, Any]:
    import asyncio

    loop = asyncio.get_running_loop()

    def sync_report(progress: int, message: str) -> None:
        if not progress_callback:
            return
        future = asyncio.run_coroutine_threadsafe(progress_callback(progress, message), loop)
        future.result(timeout=30)

    return await loop.run_in_executor(
        executor,
        lambda: run_word_count_task_sync(
            task_id=task_id,
            display_no=display_no,
            directory_path=directory_path,
            recursive=recursive,
            include_hidden=include_hidden,
            extensions=extensions,
            ocr_mode=ocr_mode,
            ocr_model=ocr_model,
            ocr_route=ocr_route,
            input_source=input_source,
            original_filename=original_filename,
            progress_callback=sync_report,
        ),
    )


def run_word_count_task_sync(
    *,
    task_id: str,
    display_no: Optional[str],
    directory_path: str,
    recursive: bool,
    include_hidden: bool,
    extensions: Iterable[str],
    ocr_mode: str = OCR_MODE_AUTO,
    ocr_model: str = PDF2DOCX_DEFAULT_MODEL,
    ocr_route: str = PDF2DOCX_DEFAULT_GEMINI_ROUTE,
    input_source: str = "path",
    original_filename: Optional[str] = None,
    progress_callback: Optional[Callable[[int, str], None]] = None,
) -> dict[str, Any]:
    started_at = datetime.now()
    normalized_input_source = str(input_source or "path").strip().lower()
    if normalized_input_source == "upload":
        input_path, matched_root, input_kind = _resolve_uploaded_word_count_path(directory_path)
    elif normalized_input_source == "path":
        input_path, matched_root, input_kind = _resolve_allowed_input_path(directory_path)
    else:
        raise ValueError(f"不支持的字数统计输入方式: {normalized_input_source}")
    normalized_ocr_mode = _normalize_ocr_mode(ocr_mode)
    normalized_ocr_model = normalize_pdf2docx_model(ocr_model)
    if normalized_ocr_model not in get_pdf2docx_models():
        raise ValueError(f"不支持的 OCR 模型: {normalized_ocr_model}")
    ocr_enabled = _ocr_enabled_for_input(normalized_ocr_mode, input_kind)
    scan_extensions = set(normalize_scan_extensions(extensions))
    output_dir = Path(settings.OUTPUT_DIR) / "word_count" / (display_no or task_id)
    output_dir.mkdir(parents=True, exist_ok=True)
    converted_dir = output_dir / "converted_inputs"
    converted_dir.mkdir(parents=True, exist_ok=True)
    ocr_text_dir = output_dir / "OCR识别文本"

    _report(progress_callback, 5, "正在扫描文件..." if input_kind == "file" else "正在扫描目录...")
    if input_kind == "file":
        candidates = [input_path] if input_path.suffix.lower() in scan_extensions else []
        relative_root = input_path.parent
    else:
        candidates = list(_iter_candidate_files(input_path, recursive, include_hidden, scan_extensions))
        relative_root = input_path
    max_files = max(1, int(settings.WORD_COUNT_MAX_FILES or 5000))
    truncated = len(candidates) > max_files
    if truncated:
        candidates = candidates[:max_files]

    file_results: list[dict[str, Any]] = []
    source_rows: list[dict[str, Any]] = []
    ocr_text_files: list[Path] = []
    total = len(candidates)
    max_bytes = max(1, int(settings.WORD_COUNT_MAX_FILE_MB or 200)) * 1024 * 1024

    if total == 0:
        _report(progress_callback, 80, "未发现可统计的候选文件，正在生成空报告...")

    for index, file_path in enumerate(candidates, start=1):
        progress = 8 + int(index / max(total, 1) * 72)
        _report(progress_callback, min(progress, 80), f"正在统计 {index}/{total}: {file_path.name}")
        result, rows, ocr_text_path = _count_single_file(
            file_path=file_path,
            root=relative_root,
            converted_dir=converted_dir,
            max_bytes=max_bytes,
            ocr_enabled=ocr_enabled,
            ocr_model=normalized_ocr_model,
            ocr_route=ocr_route,
            ocr_text_dir=ocr_text_dir,
            ocr_status_callback=lambda message, name=file_path.name, pct=min(progress, 80): _report(
                progress_callback, pct, f"{name}: {message}"
            ),
            display_relative_path=(
                str(original_filename or "").replace("\\", "/").rsplit("/", 1)[-1]
                if normalized_input_source == "upload" and input_kind == "file"
                else None
            ),
        )
        file_results.append(result)
        source_rows.extend(rows)
        if ocr_text_path is not None:
            ocr_text_files.append(ocr_text_path)

    ocr_archive_path = _write_ocr_text_archive(output_dir, ocr_text_dir, ocr_text_files)

    summary = _build_summary(file_results, truncated=truncated, started_at=started_at)
    report_payload = {
        "task_id": task_id,
        "directory_path": str(input_path),
        "input_path": str(input_path),
        "input_kind": input_kind,
        "input_source": normalized_input_source,
        "original_filename": original_filename or "",
        "allowed_root": str(matched_root),
        "recursive": bool(recursive),
        "include_hidden": bool(include_hidden),
        "extensions": sorted(scan_extensions),
        "ocr_mode": normalized_ocr_mode,
        "ocr_enabled": ocr_enabled,
        "ocr_model": normalized_ocr_model,
        "ocr_route": ocr_route,
        "ocr_text_archive": _output_web_path(ocr_archive_path) if ocr_archive_path else "",
        "summary": summary,
        "files": file_results,
        "source_details": source_rows,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "rules": [
            "主统计采用 Word 近似口径，不保证与 Word COM 组件 100% 一致。",
            "脚本分桶采用互斥口径：汉字、假名、韩文、拉丁词、纯数字等分桶合计等于主字数。",
            "报价候选字段用于按源语取列；中文候选与日语候选都可能包含同一批汉字，不能把所有候选列相加。",
            "纯数字 token 默认单独列出，不并入拉丁系候选；如报价规则要求数字计费，可与源语候选相加。",
            "Word 主数计入正文、正文表格和正文文本框；页眉页脚、脚注尾注、图表文字、批注列为额外内容。",
            "页数和行数采用快速近似口径：Word 优先读取文档属性，PDF/PPT 使用实际页数或幻灯片数，Excel 页数按工作表数近似。",
            "图片数量统计嵌入图片出现次数，用于后续 OCR 流程线索；含图片的可编辑 Office 文件仍保持已统计。",
            "PPT 主数统计幻灯片可见文本框、表格、组合形状和可提取图表文字；备注文本列为额外内容。",
            "启用 OCR 时，扫描 PDF、混合 PDF 的扫描页和独立图片使用现有 PDF2DOCX 视觉识别能力；Office 内嵌图片不识别。",
            "PDF 页面无可统计文字，或字数少于 20 且最大图片覆盖率达到 50% 时，以 OCR 结果替换该页文本层，避免重复统计。",
            "OCR 任一必要页面失败时，整个文件不计入总字数；已识别文本仍保留用于复核。",
            "CAD 按模型空间、图纸空间和实际插入块中的可提取文字实例统计；外部参照不自动追踪。",
            "DWG/DWS/DWT 依赖 ODA File Converter 转为 DXF；工具缺失时标记为需要 CAD 解析。",
        ],
    }

    _report(progress_callback, 86, "正在生成 JSON 报告...")
    json_path = output_dir / "字数统计结果.json"
    json_path.write_text(json.dumps(report_payload, ensure_ascii=False, indent=2), encoding="utf-8")

    _report(progress_callback, 92, "正在生成 Excel 报告...")
    excel_path = output_dir / "字数统计报告.xlsx"
    _write_excel_report(excel_path, report_payload)

    _report(progress_callback, 98, "正在整理输出结果...")
    report_payload["report_json"] = _output_web_path(json_path)
    report_payload["report_excel"] = _output_web_path(excel_path)
    report_payload["summary_text"] = (
        f"已统计 {summary['counted_files']} 个文件，主字数 {summary['total_main_word_count']}，"
        f"额外内容字数 {summary['total_extra_word_count']}；"
        f"中文候选 {summary.get('total_billable_chinese_count', 0)}，"
        f"拉丁系候选 {summary.get('total_billable_latin_count', 0)}。"
    )
    return report_payload


def count_words_word_like(text: str) -> int:
    return _count_text(text).word_count


def _report(callback: Optional[Callable[[int, str], None]], progress: int, message: str) -> None:
    if callback:
        callback(progress, message)


def _resolve_allowed_directory(raw_path: str) -> tuple[Path, Path]:
    candidate, allowed_root, input_kind = _resolve_allowed_input_path(raw_path)
    if input_kind != "directory":
        raise ValueError(f"路径不是目录: {candidate}")
    return candidate, allowed_root


def resolve_allowed_shared_input_path(raw_path: str) -> tuple[Path, Path, str]:
    """解析共享路径并执行与字数统计相同的白名单和 UNC 映射校验。"""
    return _resolve_allowed_input_path(raw_path)


def _resolve_allowed_input_path(raw_path: str) -> tuple[Path, Path, str]:
    raw_text = str(raw_path or "").strip()
    if not raw_text:
        raise ValueError("路径不能为空")

    mapped = _map_unc_path_to_local(raw_text)
    if mapped is not None:
        candidate, matched_root, matched_unc = mapped
        candidate = candidate.expanduser().resolve(strict=False)
        matched_root = matched_root.expanduser().resolve(strict=False)
        if not candidate.exists():
            raise FileNotFoundError(_build_mapped_path_missing_message(candidate, matched_root, matched_unc))
        input_kind = _input_kind(candidate)
        if not input_kind:
            raise ValueError(f"路径不是文件或目录: {candidate}（由 {matched_unc} 映射而来）")
        allowed_root = _match_allowed_root_for_mapped_unc(raw_text, candidate, matched_root)
        if allowed_root is None:
            allowed_text = "；".join(settings.WORD_COUNT_ALLOWED_ROOTS)
            raise PermissionError(f"路径不在允许扫描范围内。允许根目录: {allowed_text}")
        return candidate, allowed_root, input_kind

    if _is_unc_text(raw_text) and os.name != "nt":
        candidates = _format_unc_auto_candidates(raw_text)
        suffix = f" 自动尝试过: {candidates}" if candidates else ""
        raise FileNotFoundError(
            "Docker/Linux 容器不能直接读取 Windows UNC 路径；/mnt 是容器内挂载点，不是 Windows Server 路径。"
            "请在 docker-compose.yml 中把 Windows Server 共享目录挂载到 /mnt/win-server/服务器资料7，"
            "或配置 WORD_COUNT_UNC_MOUNT_MAP_JSON 将 UNC 映射到实际容器内路径。"
            f"{suffix}"
        )

    candidate = Path(raw_text).expanduser().resolve(strict=False)
    if not candidate.exists():
        raise FileNotFoundError(f"路径不存在: {candidate}")
    input_kind = _input_kind(candidate)
    if not input_kind:
        raise ValueError(f"路径不是文件或目录: {candidate}")

    if settings.WORD_COUNT_ALLOW_LOCAL_PATHS_ENABLED and _is_local_absolute_path(candidate):
        return candidate, _local_path_root(candidate), input_kind

    for raw_root in settings.WORD_COUNT_ALLOWED_ROOTS:
        allowed_root = Path(raw_root).expanduser().resolve(strict=False)
        if _is_relative_to_path(candidate, allowed_root):
            return candidate, allowed_root, input_kind
    allowed_text = "；".join(settings.WORD_COUNT_ALLOWED_ROOTS)
    raise PermissionError(f"路径不在允许扫描范围内。允许根目录: {allowed_text}")


def _resolve_uploaded_word_count_path(raw_path: str) -> tuple[Path, Path, str]:
    raw_text = str(raw_path or "").strip()
    if not raw_text:
        raise ValueError("上传文件路径不能为空")

    upload_root = (Path(settings.UPLOAD_DIR) / "word_count").expanduser().resolve(strict=False)
    candidate = Path(raw_text).expanduser().resolve(strict=False)
    if not _is_relative_to_path(candidate, upload_root):
        raise PermissionError("上传文件不在字数统计任务目录内")
    if not candidate.exists():
        raise FileNotFoundError(f"上传文件不存在: {candidate}")
    if not candidate.is_file():
        raise ValueError(f"上传目标不是文件: {candidate}")
    return candidate, upload_root, "file"


def _input_kind(path: Path) -> str:
    if path.is_dir():
        return "directory"
    if path.is_file():
        return "file"
    return ""


def _map_unc_path_to_local(raw_path: str) -> Optional[tuple[Path, Path, str]]:
    raw_unc = _normalize_unc_text(raw_path)
    if not raw_unc:
        return None
    mappings = []
    for unc_path, mount_path in settings.WORD_COUNT_UNC_MOUNT_MAP.items():
        normalized = _normalize_unc_text(unc_path)
        if normalized:
            mappings.append((normalized, str(mount_path)))
    mappings.sort(key=lambda item: len(_unc_parts_from_text(item[0])), reverse=True)
    for unc_prefix, mount_path in mappings:
        remainder = _unc_relative_parts(raw_unc, unc_prefix)
        if remainder is None:
            continue
        mount_root = Path(mount_path)
        return mount_root.joinpath(*remainder), mount_root, unc_prefix
    for unc_prefix, mount_root, candidate in _iter_unc_auto_mapping_candidates(raw_unc):
        if mount_root.exists() or candidate.exists():
            return candidate, mount_root, unc_prefix
    return None


def _build_mapped_path_missing_message(candidate: Path, matched_root: Path, matched_unc: str) -> str:
    base = f"目录不存在: {candidate}（由 {matched_unc} 映射而来）。"
    if not matched_root.exists():
        return (
            f"{base} 容器内共享根目录也不存在: {matched_root}。"
            "这表示 Docker 还没有把 Windows Server 共享目录挂载进容器；请检查 docker-compose 的 volume 或宿主机 CIFS 挂载。"
        )
    if not matched_root.is_dir():
        return (
            f"{base} 容器内共享根路径存在但不是目录: {matched_root}。"
            "请检查 volume 目标路径是否被文件占用。"
        )

    nearest = _nearest_existing_parent(candidate, stop_at=matched_root)
    entries = _sample_directory_entries(nearest)
    entry_text = "、".join(entries) if entries else "未看到任何条目"
    return (
        f"{base} 容器内共享根目录存在: {matched_root}；最近存在的父目录: {nearest}；"
        f"该目录下可见条目: {entry_text}。"
        "如果这里是空的或不是 Windows Server 上的内容，说明 Docker 当前挂载到了本地空目录，"
        "而不是 \\\\win-server\\服务器资料7。请检查 WORD_COUNT_SHARE7_HOST_PATH 或宿主机 CIFS 挂载。"
    )


def _nearest_existing_parent(path: Path, *, stop_at: Path) -> Path:
    current = path
    stop_text = os.path.normcase(str(stop_at.resolve(strict=False)))
    while not current.exists() and current != current.parent:
        current = current.parent
        current_text = os.path.normcase(str(current.resolve(strict=False)))
        if current_text == stop_text:
            break
    return current


def _sample_directory_entries(path: Path, limit: int = 8) -> list[str]:
    try:
        if not path.is_dir():
            return []
        return sorted(item.name for item in path.iterdir())[:limit]
    except OSError:
        return []


def _match_allowed_root_for_mapped_unc(raw_unc_path: str, candidate: Path, fallback_mapped_root: Optional[Path] = None) -> Optional[Path]:
    raw_unc = _normalize_unc_text(raw_unc_path)
    for raw_root in settings.WORD_COUNT_ALLOWED_ROOTS:
        root_unc = _normalize_unc_text(raw_root)
        if root_unc and raw_unc and _unc_relative_parts(raw_unc, root_unc) is not None:
            mapped_root = _map_unc_path_to_local(root_unc)
            if mapped_root is not None:
                return mapped_root[0].expanduser().resolve(strict=False)
            if fallback_mapped_root is not None:
                return fallback_mapped_root.expanduser().resolve(strict=False)
            return Path(raw_root)
        if not root_unc:
            allowed_root = Path(raw_root).expanduser().resolve(strict=False)
            if _is_relative_to_path(candidate, allowed_root):
                return allowed_root
    return None


def _iter_unc_auto_mapping_candidates(raw_unc: str) -> Iterable[tuple[str, Path, Path]]:
    parts = _unc_parts_from_text(raw_unc)
    if len(parts) < 2:
        return
    server, share = parts[0], parts[1]
    remainder = parts[2:]
    share_unc = _normalize_unc_text(f"\\\\{server}\\{share}\\")
    server_unc = _normalize_unc_text(f"\\\\{server}\\")
    for root_text in settings.WORD_COUNT_UNC_AUTO_MOUNT_ROOTS:
        base = Path(root_text).expanduser()
        share_mounts = [
            base / server / share,
            base / share,
        ]
        seen: set[str] = set()
        for mount_root in share_mounts:
            key = os.path.normcase(str(mount_root))
            if key in seen:
                continue
            seen.add(key)
            yield share_unc, mount_root, mount_root.joinpath(*remainder)
        server_mount = base / server
        server_key = os.path.normcase(str(server_mount))
        if server_key not in seen:
            yield server_unc, server_mount, server_mount.joinpath(share, *remainder)


def _format_unc_auto_candidates(raw_path: str, limit: int = 8) -> str:
    raw_unc = _normalize_unc_text(raw_path)
    if not raw_unc:
        return ""
    candidates: list[str] = []
    seen: set[str] = set()
    for _, _, candidate in _iter_unc_auto_mapping_candidates(raw_unc):
        text = str(candidate)
        key = os.path.normcase(text)
        if key in seen:
            continue
        seen.add(key)
        candidates.append(text)
        if len(candidates) >= limit:
            break
    return "；".join(candidates)


def _is_relative_to_path(candidate: Path, root: Path) -> bool:
    if _is_unc_server_root(root):
        return _unc_server_name(candidate) == _unc_server_name(root)

    try:
        candidate.resolve(strict=False).relative_to(root.resolve(strict=False))
        return True
    except ValueError:
        pass

    try:
        candidate_text = os.path.normcase(str(candidate.resolve(strict=False)))
        root_text = os.path.normcase(str(root.resolve(strict=False)))
        return os.path.commonpath([candidate_text, root_text]) == root_text
    except (ValueError, OSError):
        return False


def _is_local_absolute_path(path: Path) -> bool:
    return path.is_absolute() and not _is_unc_path(path)


def _local_path_root(path: Path) -> Path:
    anchor = path.anchor
    return Path(anchor) if anchor else path


def _is_unc_path(path: Path) -> bool:
    return _is_unc_text(str(path))


def _is_unc_server_root(path: Path) -> bool:
    return len(_unc_parts(path)) == 1


def _unc_server_name(path: Path) -> Optional[str]:
    parts = _unc_parts(path)
    return parts[0].lower() if parts else None


def _unc_parts(path: Path) -> list[str]:
    return _unc_parts_from_text(str(path))


def _is_unc_text(value: str) -> bool:
    return str(value or "").strip().replace("/", "\\").startswith("\\\\")


def _normalize_unc_text(value: str) -> str:
    text = str(value or "").strip().replace("/", "\\")
    if not text.startswith("\\\\"):
        return ""
    parts = _unc_parts_from_text(text)
    if not parts:
        return ""
    return "\\\\" + "\\".join(parts) + "\\"


def _unc_parts_from_text(value: str) -> list[str]:
    text = str(value or "").strip().replace("/", "\\")
    if not text.startswith("\\\\"):
        return []
    return [part for part in text.strip("\\").split("\\") if part]


def _unc_relative_parts(candidate_unc: str, root_unc: str) -> Optional[list[str]]:
    candidate_parts = _unc_parts_from_text(candidate_unc)
    root_parts = _unc_parts_from_text(root_unc)
    if not candidate_parts or not root_parts or len(candidate_parts) < len(root_parts):
        return None
    candidate_head = [part.lower() for part in candidate_parts[: len(root_parts)]]
    root_head = [part.lower() for part in root_parts]
    if candidate_head != root_head:
        return None
    return candidate_parts[len(root_parts):]


def _iter_candidate_files(
    root: Path,
    recursive: bool,
    include_hidden: bool,
    extensions: set[str],
) -> Iterable[Path]:
    followlinks = settings.WORD_COUNT_FOLLOW_SYMLINKS_ENABLED

    if not recursive:
        for child in sorted(root.iterdir(), key=lambda item: item.name.lower()):
            if not include_hidden and _is_hidden_name(child.name):
                continue
            if child.is_file() and child.suffix.lower() in extensions:
                if followlinks or not child.is_symlink():
                    yield child
        return

    for current_dir, dir_names, file_names in os.walk(root, followlinks=followlinks):
        if not include_hidden:
            dir_names[:] = [name for name in dir_names if not _is_hidden_name(name)]
            file_names = [name for name in file_names if not _is_hidden_name(name)]
        dir_names.sort(key=str.lower)
        for file_name in sorted(file_names, key=str.lower):
            file_path = Path(current_dir) / file_name
            if file_path.suffix.lower() not in extensions:
                continue
            if not followlinks and file_path.is_symlink():
                continue
            yield file_path


def _is_hidden_name(name: str) -> bool:
    return name.startswith(".") or name.startswith("~$")


def _count_single_file(
    *,
    file_path: Path,
    root: Path,
    converted_dir: Path,
    max_bytes: int,
    ocr_enabled: bool = False,
    ocr_model: str = PDF2DOCX_DEFAULT_MODEL,
    ocr_route: str = PDF2DOCX_DEFAULT_GEMINI_ROUTE,
    ocr_text_dir: Optional[Path] = None,
    ocr_status_callback: Optional[Callable[[str], None]] = None,
    display_relative_path: Optional[str] = None,
) -> tuple[dict[str, Any], list[dict[str, Any]], Optional[Path]]:
    now_status = STATUS_COUNTED
    ext = file_path.suffix.lower()
    relative_path = str(display_relative_path or _relative_path(file_path, root))
    error = ""
    warning = ""
    source_rows: list[dict[str, Any]] = []

    try:
        stat = file_path.stat()
        size_bytes = int(stat.st_size)
        modified_at = datetime.fromtimestamp(stat.st_mtime).isoformat(timespec="seconds")
    except OSError as exc:
        return (
            _base_file_result(file_path, relative_path, ext, STATUS_FAILED, error=f"无法读取文件信息: {exc}"),
            [],
            None,
        )

    base = _base_file_result(
        file_path,
        relative_path,
        ext,
        now_status,
        size_bytes=size_bytes,
        modified_at=modified_at,
    )

    if size_bytes > max_bytes:
        base.update(
            status=STATUS_SKIPPED,
            message=f"文件超过大小限制 {settings.WORD_COUNT_MAX_FILE_MB} MB",
            counted_at=_now_iso(),
        )
        return base, [], None

    if ext in IMAGE_EXTENSIONS and not ocr_enabled:
        base.update(
            status=STATUS_NEEDS_OCR,
            image_count=1,
            message="图片文件需要 OCR 后再统计",
            stat_method=_stat_method_for_extension(ext),
            counted_at=_now_iso(),
        )
        return base, [], None
    ocr_info = {
        "ocr_used": False,
        "ocr_page_count": 0,
        "ocr_model": "",
        "ocr_failed_pages": [],
    }
    try:
        if ext in IMAGE_EXTENSIONS:
            extracted, ocr_info = _extract_image_ocr_content(
                file_path,
                model=ocr_model,
                gemini_route=ocr_route,
                status_callback=ocr_status_callback,
            )
        elif ext == ".pdf" and ocr_enabled:
            extracted, ocr_info = _extract_pdf_content_with_ocr(
                file_path,
                model=ocr_model,
                gemini_route=ocr_route,
                status_callback=ocr_status_callback,
            )
        else:
            extracted = _extract_content(file_path, converted_dir)
    except CadConverterUnavailableError as exc:
        base.update(
            status=STATUS_NEEDS_CAD_PARSER,
            message="CAD 解析工具未就绪",
            error=str(exc),
            stat_method=_stat_method_for_extension(ext),
            counted_at=_now_iso(),
        )
        return base, [], None
    except Exception as exc:
        if ext in IMAGE_EXTENSIONS or (ext == ".pdf" and ocr_enabled):
            failed_message = "OCR 失败"
        elif ext in CAD_EXTENSIONS:
            failed_message = "CAD 解析失败"
        else:
            failed_message = "解析失败"
        base.update(
            status=STATUS_FAILED,
            error=str(exc),
            message=failed_message,
            ocr_used=bool(ext in IMAGE_EXTENSIONS or (ext == ".pdf" and ocr_enabled)),
            ocr_model=ocr_model if ext in IMAGE_EXTENSIONS or ext == ".pdf" else "",
            counted_at=_now_iso(),
        )
        return base, [], None

    items = extracted.items
    warning = extracted.warning
    if ext == ".pdf" and not ocr_info.get("ocr_used") and not any(item.text.strip() for item in items):
        base.update(
            status=STATUS_NEEDS_OCR,
            page_count=extracted.page_count,
            image_count=extracted.image_count,
            stat_method=extracted.stat_method,
            message="PDF 未提取到可编辑文本，可能是扫描件",
            counted_at=_now_iso(),
        )
        return base, [], None

    ocr_text_path: Optional[Path] = None
    if ocr_info.get("ocr_used") and ocr_text_dir is not None:
        ocr_text_path = _write_ocr_plain_text(
            ocr_text_dir=ocr_text_dir,
            relative_path=relative_path,
            items=items,
        )
        base["ocr_text_path"] = _output_web_path(ocr_text_path)

    failed_pages = [int(page) for page in ocr_info.get("ocr_failed_pages") or []]
    if failed_pages:
        base.update(
            status=STATUS_FAILED,
            page_count=extracted.page_count,
            paragraph_count=extracted.paragraph_count,
            line_count=extracted.line_count,
            image_count=extracted.image_count,
            stat_method=extracted.stat_method,
            file_type=extracted.file_type or base.get("file_type", ""),
            message="OCR 页面识别不完整，文件未计入总字数",
            error=f"OCR 失败页: {', '.join(str(page) for page in failed_pages)}",
            counted_at=_now_iso(),
            **ocr_info,
        )
        return base, [], ocr_text_path

    totals = {
        "main_word_count": 0,
        "extra_word_count": 0,
        "main_non_space_chars": 0,
        "main_raw_chars": 0,
        "extra_non_space_chars": 0,
        "extra_raw_chars": 0,
    }
    for field in SCRIPT_COUNT_FIELDS:
        totals[f"main_{field}"] = 0
        totals[f"extra_{field}"] = 0
    source_counter: Counter[str] = Counter()
    for item in items:
        metrics = _count_text(item.text)
        if item.is_extra:
            totals["extra_word_count"] += metrics.word_count
            totals["extra_non_space_chars"] += metrics.non_space_chars
            totals["extra_raw_chars"] += metrics.raw_chars
            for field in SCRIPT_COUNT_FIELDS:
                totals[f"extra_{field}"] += int(getattr(metrics, field, 0))
        else:
            totals["main_word_count"] += metrics.word_count
            totals["main_non_space_chars"] += metrics.non_space_chars
            totals["main_raw_chars"] += metrics.raw_chars
            for field in SCRIPT_COUNT_FIELDS:
                totals[f"main_{field}"] += int(getattr(metrics, field, 0))
        source_counter[item.source_type] += 1
        item_script_counts = metrics.script_counts()
        item_quote_counts = metrics.quote_counts()
        source_rows.append(
            {
                "file_path": str(file_path),
                "relative_path": relative_path,
                "extension": ext,
                "source_type": item.source_type,
                "source_label": item.source_label,
                "is_extra": item.is_extra,
                "word_count": metrics.word_count,
                "non_space_chars": metrics.non_space_chars,
                "raw_chars": metrics.raw_chars,
                "char_count_no_spaces": metrics.non_space_chars,
                "char_count_with_spaces": metrics.raw_chars,
                "script_counts": item_script_counts,
                "quote_counts": item_quote_counts,
                "script_count_total": metrics.script_count_total,
                "paragraph_count": item.paragraph_count,
                "line_count": item.line_count,
                "image_count": item.image_count,
                "text_preview": _preview_text(item.text),
                **item_script_counts,
                **item_quote_counts,
            }
        )

    if not items:
        warning = "；".join(part for part in (warning, "未提取到文本") if part)

    main_script_counts = {field: int(totals.get(f"main_{field}") or 0) for field in SCRIPT_COUNT_FIELDS}
    extra_script_counts = {field: int(totals.get(f"extra_{field}") or 0) for field in SCRIPT_COUNT_FIELDS}
    main_quote_counts = _quote_counts_from_script_counts(main_script_counts)
    extra_quote_counts = _quote_counts_from_script_counts(extra_script_counts)

    base.update(
        status=STATUS_COUNTED,
        main_word_count=totals["main_word_count"],
        word_count=totals["main_word_count"],
        extra_word_count=totals["extra_word_count"],
        char_count_no_spaces=totals["main_non_space_chars"],
        char_count_with_spaces=totals["main_raw_chars"],
        non_space_chars=totals["main_non_space_chars"],
        raw_chars=totals["main_raw_chars"],
        extra_non_space_chars=totals["extra_non_space_chars"],
        extra_raw_chars=totals["extra_raw_chars"],
        script_counts=main_script_counts,
        quote_counts=main_quote_counts,
        extra_script_counts=extra_script_counts,
        extra_quote_counts=extra_quote_counts,
        script_count_total=sum(main_script_counts.values()),
        extra_script_count_total=sum(extra_script_counts.values()),
        page_count=extracted.page_count,
        paragraph_count=extracted.paragraph_count,
        line_count=extracted.line_count,
        image_count=extracted.image_count,
        stat_method=extracted.stat_method,
        file_type=extracted.file_type or base.get("file_type", ""),
        source_counts=dict(source_counter),
        warning=warning,
        message="统计完成",
        counted_at=_now_iso(),
        **ocr_info,
        **main_script_counts,
        **main_quote_counts,
        **{f"extra_{field}": value for field, value in extra_script_counts.items()},
        **{f"extra_{field}": value for field, value in extra_quote_counts.items()},
    )
    if ocr_info.get("ocr_used"):
        base["message"] = f"OCR 统计完成（{int(ocr_info.get('ocr_page_count') or 0)} 页）"
    return base, source_rows, ocr_text_path


def _base_file_result(
    file_path: Path,
    relative_path: str,
    extension: str,
    status: str,
    *,
    size_bytes: int = 0,
    modified_at: str = "",
    error: str = "",
) -> dict[str, Any]:
    script_counts = {field: 0 for field in SCRIPT_COUNT_FIELDS}
    quote_counts = _quote_counts_from_script_counts(script_counts)
    return {
        "file_path": str(file_path),
        "relative_path": relative_path,
        "filename": Path(relative_path).name or file_path.name,
        "extension": extension,
        "file_type": _file_type_for_extension(extension),
        "status": status,
        "message": "",
        "main_word_count": 0,
        "word_count": 0,
        "extra_word_count": 0,
        "char_count_no_spaces": 0,
        "char_count_with_spaces": 0,
        "non_space_chars": 0,
        "raw_chars": 0,
        "extra_non_space_chars": 0,
        "extra_raw_chars": 0,
        "script_counts": script_counts,
        "quote_counts": quote_counts,
        "extra_script_counts": dict(script_counts),
        "extra_quote_counts": dict(quote_counts),
        "script_count_total": 0,
        "extra_script_count_total": 0,
        "page_count": 0,
        "paragraph_count": 0,
        "line_count": 0,
        "image_count": 0,
        "stat_method": _stat_method_for_extension(extension),
        "counted_at": _now_iso(),
        "size_bytes": size_bytes,
        "modified_at": modified_at,
        "source_counts": {},
        "ocr_used": False,
        "ocr_page_count": 0,
        "ocr_model": "",
        "ocr_failed_pages": [],
        "ocr_text_path": "",
        "warning": "",
        "error": error,
        **script_counts,
        **quote_counts,
        **{f"extra_{field}": 0 for field in SCRIPT_COUNT_FIELDS},
        **{f"extra_{field}": 0 for field in QUOTE_COUNT_FIELDS},
    }


def _extract_text_items(file_path: Path, converted_dir: Path) -> list[TextItem]:
    return _extract_content(file_path, converted_dir).items


def _extract_content(file_path: Path, converted_dir: Path) -> ExtractedContent:
    ext = file_path.suffix.lower()
    converted_from = ""
    if ext in CAD_EXTENSIONS:
        cad_content = extract_cad_text(
            file_path,
            workspace_dir=converted_dir,
            oda_path=settings.ODA_FILE_CONVERTER_PATH,
            timeout_seconds=settings.WORD_COUNT_CAD_CONVERT_TIMEOUT_SECONDS,
        )
        return ExtractedContent(
            items=[
                TextItem(
                    source_type=item.source_type,
                    source_label=item.source_label,
                    text=item.text,
                    paragraph_count=item.paragraph_count,
                    line_count=item.line_count,
                )
                for item in cad_content.items
            ],
            page_count=cad_content.page_count,
            paragraph_count=cad_content.paragraph_count,
            line_count=cad_content.line_count,
            stat_method=cad_content.stat_method,
            file_type="CAD",
            warning=cad_content.warning,
        )
    if ext == ".doc":
        target = converted_dir / f"{_safe_stem(file_path)}.docx"
        file_path = Path(convert_doc_to_docx_via_libreoffice(file_path, target))
        converted_from = "DOC经LibreOffice转换"
        ext = ".docx"
    elif ext == ".xls":
        target = converted_dir / f"{_safe_stem(file_path)}.xlsx"
        file_path = Path(convert_spreadsheet_to_xlsx_via_libreoffice(file_path, target))
        converted_from = "XLS经LibreOffice转换"
        ext = ".xlsx"
    elif ext == ".ppt":
        target = converted_dir / f"{_safe_stem(file_path)}.pptx"
        file_path = Path(convert_presentation_to_pptx_via_libreoffice(file_path, target))
        converted_from = "PPT经LibreOffice转换"
        ext = ".pptx"

    if ext == ".docx":
        content = _extract_docx_content(file_path)
    if ext == ".xlsx":
        content = _extract_xlsx_content(file_path)
    if ext == ".pptx":
        content = _extract_pptx_content(file_path)
    if ext == ".pdf":
        content = _extract_pdf_content(file_path)
    if ext == ".txt":
        content = _extract_txt_content(file_path)
    if ext not in {".docx", ".xlsx", ".pptx", ".pdf", ".txt"}:
        raise ValueError(f"不支持的文件格式: {ext}")
    if converted_from:
        content = ExtractedContent(
            items=content.items,
            page_count=content.page_count,
            paragraph_count=content.paragraph_count,
            line_count=content.line_count,
            image_count=content.image_count,
            stat_method=f"{converted_from}+{content.stat_method}",
            file_type=content.file_type,
            warning=content.warning,
        )
    return content


def _safe_stem(path: Path) -> str:
    normalized = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff._-]+", "_", path.stem).strip("._")
    return normalized or "converted"


def _extract_txt_content(path: Path) -> ExtractedContent:
    items = _extract_txt_text_items(path)
    text = "\n".join(item.text for item in items)
    return ExtractedContent(
        items=items,
        page_count=1 if text.strip() else 0,
        paragraph_count=_count_text_paragraphs(text),
        line_count=_count_text_lines(text),
        image_count=0,
        stat_method="TXT文本解析+Word近似计数",
        file_type="TXT",
    )


def _extract_txt_text_items(path: Path) -> list[TextItem]:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "big5"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    else:
        try:
            import chardet

            detected = chardet.detect(raw).get("encoding") or "utf-8"
            text = raw.decode(detected, errors="replace")
        except Exception:
            text = raw.decode("utf-8", errors="replace")
    return [
        TextItem(
            "txt",
            "文本文件",
            text,
            is_extra=False,
            paragraph_count=_count_text_paragraphs(text),
            line_count=_count_text_lines(text),
        )
    ] if text else []


def _extract_pdf_content(path: Path) -> ExtractedContent:
    import fitz

    items: list[TextItem] = []
    with fitz.open(str(path)) as doc:
        page_count = doc.page_count
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                items.append(
                    TextItem(
                        "pdf_page",
                        f"第 {page_index} 页",
                        text,
                        is_extra=False,
                        paragraph_count=_count_text_lines(text),
                        line_count=_count_text_lines(text),
                    )
                )
    return ExtractedContent(
        items=items,
        page_count=page_count,
        paragraph_count=sum(item.paragraph_count for item in items),
        line_count=sum(item.line_count for item in items),
        image_count=_count_pdf_images(path),
        stat_method="PDF文本层解析+Word近似计数",
        file_type="PDF",
    )


def _extract_image_ocr_content(
    path: Path,
    *,
    model: str,
    gemini_route: str,
    status_callback: Optional[Callable[[str], None]] = None,
) -> tuple[ExtractedContent, dict[str, Any]]:
    payload = extract_ocr_plain_text(
        file_path=str(path),
        model=model,
        gemini_route=gemini_route,
        page_progress_callback=(
            (lambda current, total: status_callback(f"正在 OCR 第 {current}/{total} 页"))
            if status_callback
            else None
        ),
        status_callback=status_callback,
        continue_on_error=True,
    )
    page_results = payload.get("page_results") or []
    items = _ocr_page_items(page_results, source_type="image_ocr_page")
    total_pages = max(int(payload.get("total_pages") or 0), 1)
    content = ExtractedContent(
        items=items,
        page_count=total_pages,
        paragraph_count=sum(item.paragraph_count for item in items),
        line_count=sum(item.line_count for item in items),
        image_count=total_pages,
        stat_method=f"图片视觉OCR（{model}）+Word近似计数",
        file_type="图片",
    )
    return content, _ocr_info_from_payload(payload, model)


def _extract_pdf_content_with_ocr(
    path: Path,
    *,
    model: str,
    gemini_route: str,
    status_callback: Optional[Callable[[str], None]] = None,
) -> tuple[ExtractedContent, dict[str, Any]]:
    import fitz

    page_data: list[dict[str, Any]] = []
    image_count = 0
    with fitz.open(str(path)) as doc:
        page_count = doc.page_count
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            word_count = _count_text(text).word_count
            image_count += len(page.get_images(full=True))
            image_coverage = _max_pdf_page_image_coverage(page)
            needs_ocr = word_count == 0 or (
                word_count < PDF_OCR_SPARSE_WORD_THRESHOLD
                and image_coverage >= PDF_OCR_IMAGE_COVERAGE_THRESHOLD
            )
            page_data.append(
                {
                    "page_number": page_index,
                    "text": text,
                    "needs_ocr": needs_ocr,
                }
            )

    ocr_pages = [int(item["page_number"]) for item in page_data if item["needs_ocr"]]
    if not ocr_pages:
        items = _pdf_text_layer_items(page_data)
        return (
            ExtractedContent(
                items=items,
                page_count=page_count,
                paragraph_count=sum(item.paragraph_count for item in items),
                line_count=sum(item.line_count for item in items),
                image_count=image_count,
                stat_method="PDF文本层解析+Word近似计数",
                file_type="PDF",
            ),
            {"ocr_used": False, "ocr_page_count": 0, "ocr_model": "", "ocr_failed_pages": []},
        )

    payload = extract_ocr_plain_text(
        file_path=str(path),
        model=model,
        gemini_route=gemini_route,
        page_numbers=ocr_pages,
        page_progress_callback=(
            (lambda current, total: status_callback(f"正在 OCR 第 {current}/{total} 页"))
            if status_callback
            else None
        ),
        status_callback=status_callback,
        continue_on_error=True,
    )
    ocr_results = {
        int(item.get("page_number") or 0): item
        for item in (payload.get("page_results") or [])
        if int(item.get("page_number") or 0) > 0
    }
    items: list[TextItem] = []
    for page in page_data:
        page_number = int(page["page_number"])
        if page["needs_ocr"]:
            ocr_result = ocr_results.get(page_number) or {}
            text = "" if ocr_result.get("error") else str(ocr_result.get("text") or "")
            source_type = "pdf_ocr_page"
        else:
            text = str(page.get("text") or "")
            source_type = "pdf_page"
        if text.strip():
            items.append(
                TextItem(
                    source_type,
                    f"第 {page_number} 页",
                    text,
                    is_extra=False,
                    paragraph_count=_count_text_paragraphs(text),
                    line_count=_count_text_lines(text),
                )
            )

    method = "PDF视觉OCR+Word近似计数" if len(ocr_pages) == page_count else "PDF文本层+按页视觉OCR+Word近似计数"
    content = ExtractedContent(
        items=items,
        page_count=page_count,
        paragraph_count=sum(item.paragraph_count for item in items),
        line_count=sum(item.line_count for item in items),
        image_count=image_count,
        stat_method=f"{method}（{model}）",
        file_type="PDF",
    )
    return content, _ocr_info_from_payload(payload, model)


def _max_pdf_page_image_coverage(page: Any) -> float:
    page_area = max(float(page.rect.width) * float(page.rect.height), 1.0)
    max_area = 0.0
    try:
        blocks = page.get_text("blocks") or []
    except Exception:
        return 0.0
    for block in blocks:
        if len(block) < 7 or int(block[6]) != 1:
            continue
        width = max(float(block[2]) - float(block[0]), 0.0)
        height = max(float(block[3]) - float(block[1]), 0.0)
        max_area = max(max_area, width * height)
    return min(max_area / page_area, 1.0)


def _pdf_text_layer_items(page_data: list[dict[str, Any]]) -> list[TextItem]:
    items: list[TextItem] = []
    for page in page_data:
        text = str(page.get("text") or "")
        if not text.strip():
            continue
        page_number = int(page.get("page_number") or 0)
        items.append(
            TextItem(
                "pdf_page",
                f"第 {page_number} 页",
                text,
                is_extra=False,
                paragraph_count=_count_text_paragraphs(text),
                line_count=_count_text_lines(text),
            )
        )
    return items


def _ocr_page_items(page_results: list[dict[str, Any]], *, source_type: str) -> list[TextItem]:
    items: list[TextItem] = []
    for page in page_results:
        if page.get("error"):
            continue
        text = str(page.get("text") or "")
        if not text.strip():
            continue
        page_number = int(page.get("page_number") or 0)
        items.append(
            TextItem(
                source_type,
                f"第 {page_number} 页",
                text,
                is_extra=False,
                paragraph_count=_count_text_paragraphs(text),
                line_count=_count_text_lines(text),
            )
        )
    return items


def _ocr_info_from_payload(payload: dict[str, Any], model: str) -> dict[str, Any]:
    processed_pages = [int(page) for page in payload.get("processed_pages") or []]
    failed_pages = [int(page) for page in payload.get("failed_pages") or []]
    return {
        "ocr_used": True,
        "ocr_page_count": len(processed_pages),
        "ocr_model": model,
        "ocr_failed_pages": failed_pages,
    }


def _extract_pdf_text_items(path: Path) -> list[TextItem]:
    return _extract_pdf_content(path).items


def _extract_xlsx_content(path: Path) -> ExtractedContent:
    workbook = load_workbook(path, read_only=True, data_only=True)
    items: list[TextItem] = []
    sheet_count = 0
    total_non_empty_rows = 0
    total_non_empty_cells = 0
    try:
        for sheet in workbook.worksheets:
            sheet_count += 1
            values: list[str] = []
            sheet_non_empty_rows = 0
            sheet_non_empty_cells = 0
            for row in sheet.iter_rows():
                row_values: list[str] = []
                for cell in row:
                    if cell.value is None:
                        continue
                    text = str(cell.value).strip()
                    if text:
                        row_values.append(text)
                if row_values:
                    sheet_non_empty_rows += 1
                    sheet_non_empty_cells += len(row_values)
                    values.extend(row_values)
            total_non_empty_rows += sheet_non_empty_rows
            total_non_empty_cells += sheet_non_empty_cells
            if values:
                items.append(
                    TextItem(
                        "excel_cell",
                        f"工作表: {sheet.title}",
                        "\n".join(values),
                        is_extra=False,
                        paragraph_count=sheet_non_empty_cells,
                        line_count=sheet_non_empty_rows,
                    )
                )
    finally:
        workbook.close()
    return ExtractedContent(
        items=items,
        page_count=sheet_count,
        paragraph_count=total_non_empty_cells,
        line_count=total_non_empty_rows,
        image_count=_count_ooxml_image_references(path, ("xl/drawings/",)),
        stat_method="XLSX单元格解析+Word近似计数（页数按工作表数近似）",
        file_type="Excel",
    )


def _extract_xlsx_text_items(path: Path) -> list[TextItem]:
    return _extract_xlsx_content(path).items


def _extract_pptx_content(path: Path) -> ExtractedContent:
    from pptx import Presentation

    prs = Presentation(str(path))
    items: list[TextItem] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        paragraph_count = 0
        line_count = 0
        for shape in _iter_ppt_shapes(slide.shapes):
            shape_texts, shape_paragraphs, shape_lines = _extract_ppt_shape_texts(shape)
            texts.extend(shape_texts)
            paragraph_count += shape_paragraphs
            line_count += shape_lines
        if texts:
            items.append(
                TextItem(
                    "ppt_slide",
                    f"第 {slide_index} 张幻灯片",
                    "\n".join(texts),
                    is_extra=False,
                    paragraph_count=paragraph_count or _count_text_paragraphs("\n".join(texts)),
                    line_count=line_count or _count_text_lines("\n".join(texts)),
                )
            )
        notes_text = _extract_ppt_notes_text(slide)
        if notes_text:
            items.append(
                TextItem(
                    "ppt_notes",
                    f"第 {slide_index} 张幻灯片备注",
                    notes_text,
                    is_extra=True,
                    paragraph_count=_count_text_paragraphs(notes_text),
                    line_count=_count_text_lines(notes_text),
                )
            )
    return ExtractedContent(
        items=items,
        page_count=len(prs.slides),
        paragraph_count=sum(item.paragraph_count for item in items),
        line_count=sum(item.line_count for item in items),
        image_count=_count_ooxml_image_references(path, ("ppt/slides/",)),
        stat_method="PPTX形状文本解析+Word近似计数",
        file_type="PPT",
    )


def _extract_pptx_text_items(path: Path) -> list[TextItem]:
    return _extract_pptx_content(path).items


def _iter_ppt_shapes(shapes) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            yield from _iter_ppt_shapes(nested)


def _extract_ppt_shape_texts(shape) -> tuple[list[str], int, int]:
    texts: list[str] = []
    paragraph_count = 0
    line_count = 0
    if getattr(shape, "has_text_frame", False):
        text, paragraphs, lines = _extract_text_frame_text(shape.text_frame)
        if text:
            texts.append(text)
            paragraph_count += paragraphs
            line_count += lines
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text, paragraphs, lines = _extract_text_frame_text(cell.text_frame)
                if text:
                    texts.append(text)
                    paragraph_count += paragraphs
                    line_count += lines
    if getattr(shape, "has_chart", False):
        chart_text = _extract_ppt_chart_text(shape)
        if chart_text:
            texts.append(chart_text)
            paragraph_count += _count_text_paragraphs(chart_text)
            line_count += _count_text_lines(chart_text)
    return texts, paragraph_count, line_count


def _extract_text_frame_text(text_frame) -> tuple[str, int, int]:
    texts: list[str] = []
    paragraph_count = 0
    line_count = 0
    for paragraph in getattr(text_frame, "paragraphs", []) or []:
        text = (paragraph.text or "").strip()
        if not text:
            continue
        texts.append(text)
        paragraph_count += 1
        line_count += _count_text_lines(text)
    if not texts:
        text = (getattr(text_frame, "text", "") or "").strip()
        if text:
            texts.append(text)
            paragraph_count = _count_text_paragraphs(text)
            line_count = _count_text_lines(text)
    return "\n".join(texts), paragraph_count, line_count


def _extract_ppt_chart_text(shape) -> str:
    try:
        chart_space = getattr(shape.chart, "_chartSpace", None)
        if chart_space is None:
            return ""
        return "\n".join(_extract_chart_texts(etree.tostring(chart_space)))
    except Exception:
        return ""


def _extract_ppt_notes_text(slide) -> str:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
        if notes_frame is None:
            return ""
        text, _, _ = _extract_text_frame_text(notes_frame)
        return text
    except Exception:
        return ""


def _qn(tag: str) -> str:
    prefix, local = tag.split(":")
    return f"{{{NS[prefix]}}}{local}"


def _extract_docx_content(path: Path) -> ExtractedContent:
    items = _extract_docx_text_items(path)
    app_props = _read_ooxml_app_properties(path)
    page_count = _to_positive_int(app_props.get("Pages"))
    app_line_count = _to_positive_int(app_props.get("Lines"))
    fallback_line_count = sum(item.line_count for item in items)
    return ExtractedContent(
        items=items,
        page_count=page_count,
        paragraph_count=sum(item.paragraph_count for item in items),
        line_count=app_line_count or fallback_line_count,
        image_count=_count_ooxml_image_references(path, ("word/",)),
        stat_method="DOCX XML解析+Word近似计数",
        file_type="Word",
    )


def _extract_docx_text_items(path: Path) -> list[TextItem]:
    items: list[TextItem] = []
    with ZipFile(path) as zf:
        names = set(zf.namelist())
        if "word/document.xml" not in names:
            raise ValueError("DOCX 结构异常，缺少 word/document.xml")

        document_root = etree.fromstring(zf.read("word/document.xml"))
        body = document_root.find(_qn("w:body"))
        if body is not None:
            for child in body:
                if child.tag == _qn("w:p"):
                    _append_docx_paragraph_items(items, child, "body", "正文")
                elif child.tag == _qn("w:tbl"):
                    _append_docx_table_items(items, child)

        for xml_name in sorted(name for name in names if name.startswith("word/header") and name.endswith(".xml")):
            text = _extract_xml_text(zf.read(xml_name))
            if text.strip():
                items.append(
                    TextItem(
                        "header",
                        "页眉",
                        text,
                        is_extra=True,
                        paragraph_count=_count_text_paragraphs(text),
                        line_count=_count_text_lines(text),
                    )
                )

        for xml_name in sorted(name for name in names if name.startswith("word/footer") and name.endswith(".xml")):
            text = _extract_xml_text(zf.read(xml_name))
            if text.strip():
                items.append(
                    TextItem(
                        "footer",
                        "页脚",
                        text,
                        is_extra=True,
                        paragraph_count=_count_text_paragraphs(text),
                        line_count=_count_text_lines(text),
                    )
                )

        for xml_name, source_type, label in (
            ("word/footnotes.xml", "footnote", "脚注"),
            ("word/endnotes.xml", "endnote", "尾注"),
            ("word/comments.xml", "comment", "批注"),
        ):
            if xml_name in names:
                for text in _extract_note_like_texts(zf.read(xml_name), source_type):
                    items.append(
                        TextItem(
                            source_type,
                            label,
                            text,
                            is_extra=True,
                            paragraph_count=_count_text_paragraphs(text),
                            line_count=_count_text_lines(text),
                        )
                    )

        for xml_name in sorted(name for name in names if name.startswith("word/charts/") and name.endswith(".xml")):
            chart_text = "\n".join(_extract_chart_texts(zf.read(xml_name)))
            if chart_text.strip():
                items.append(
                    TextItem(
                        "chart",
                        "图表",
                        chart_text,
                        is_extra=True,
                        paragraph_count=_count_text_paragraphs(chart_text),
                        line_count=_count_text_lines(chart_text),
                    )
                )

    return items


def _append_docx_paragraph_items(items: list[TextItem], paragraph, source_type: str, label: str) -> None:
    normal_text = _extract_text_from_element(paragraph, skip_textboxes=True)
    if normal_text.strip():
        items.append(
            TextItem(
                source_type,
                label,
                normal_text,
                is_extra=False,
                paragraph_count=1,
                line_count=_count_text_lines(normal_text),
            )
        )
    for textbox_text in _extract_textbox_texts(paragraph):
        if textbox_text.strip():
            items.append(
                TextItem(
                    "textbox",
                    "正文文本框",
                    textbox_text,
                    is_extra=False,
                    paragraph_count=_count_text_paragraphs(textbox_text),
                    line_count=_count_text_lines(textbox_text),
                )
            )


def _append_docx_table_items(items: list[TextItem], table) -> None:
    for row in table.findall("w:tr", NS):
        row_texts: list[str] = []
        textbox_texts: list[str] = []
        for cell in row.findall("w:tc", NS):
            text = _extract_text_from_element(cell, skip_textboxes=True)
            if text.strip():
                row_texts.append(text)
            textbox_texts.extend(_extract_textbox_texts(cell))
        if row_texts:
            row_text = "\t".join(row_texts)
            items.append(
                TextItem(
                    "table",
                    "正文表格",
                    row_text,
                    is_extra=False,
                    paragraph_count=len(row_texts),
                    line_count=max(1, len(row_texts)),
                )
            )
        for textbox_text in textbox_texts:
            if textbox_text.strip():
                items.append(
                    TextItem(
                        "textbox",
                        "正文文本框",
                        textbox_text,
                        is_extra=False,
                        paragraph_count=_count_text_paragraphs(textbox_text),
                        line_count=_count_text_lines(textbox_text),
                    )
                )


def _extract_xml_text(xml_bytes: bytes) -> str:
    root = etree.fromstring(xml_bytes)
    return _extract_text_from_element(root, skip_textboxes=False)


def _extract_note_like_texts(xml_bytes: bytes, source_type: str) -> list[str]:
    root = etree.fromstring(xml_bytes)
    if source_type == "footnote":
        item_tag = _qn("w:footnote")
    elif source_type == "endnote":
        item_tag = _qn("w:endnote")
    else:
        item_tag = _qn("w:comment")

    texts: list[str] = []
    for elem in root.iter(item_tag):
        note_id = elem.get(_qn("w:id"), "")
        if note_id in {"-1", "0", "1"} and source_type in {"footnote", "endnote"}:
            continue
        text = _extract_text_from_element(elem, skip_textboxes=False)
        if text.strip():
            texts.append(text)
    return texts


def _extract_text_from_element(element, *, skip_textboxes: bool) -> str:
    textbox_text_ids: set[int] = set()
    if skip_textboxes:
        for textbox in _iter_textbox_elements(element):
            for text_node in textbox.iter(_qn("w:t")):
                textbox_text_ids.add(id(text_node))

    deleted_tag = _qn("w:del")
    instr_tag = _qn("w:instrText")
    parts: list[str] = []
    for text_node in element.iter():
        if text_node.tag not in {_qn("w:t"), instr_tag}:
            continue
        if text_node.tag == instr_tag:
            continue
        if id(text_node) in textbox_text_ids:
            continue
        if _has_ancestor(text_node, deleted_tag):
            continue
        if text_node.text:
            parts.append(text_node.text)
    return "".join(parts)


def _iter_textbox_elements(element) -> Iterable[Any]:
    for tag in (_qn("wps:txbx"), _qn("w:txbxContent"), _qn("v:textbox")):
        yield from element.iter(tag)


def _extract_textbox_texts(element) -> list[str]:
    texts: list[str] = []
    seen: set[int] = set()
    for textbox in _iter_textbox_elements(element):
        if id(textbox) in seen:
            continue
        seen.add(id(textbox))
        text = _extract_text_from_element(textbox, skip_textboxes=False)
        if text.strip():
            texts.append(text)
    return texts


def _has_ancestor(node, ancestor_tag: str) -> bool:
    parent = node.getparent()
    while parent is not None:
        if parent.tag == ancestor_tag:
            return True
        parent = parent.getparent()
    return False


def _extract_chart_texts(xml_bytes: bytes) -> list[str]:
    root = etree.fromstring(xml_bytes)
    seen: set[str] = set()
    texts: list[str] = []

    def add(value: Optional[str]) -> None:
        text = (value or "").strip()
        if text and text not in seen:
            seen.add(text)
            texts.append(text)

    for node in root.iter(_qn("a:t")):
        add(node.text)
    for node in root.iter(_qn("c:v")):
        value = (node.text or "").strip()
        if not value:
            continue
        try:
            float(value)
            continue
        except ValueError:
            add(value)
    return texts


def _read_ooxml_app_properties(path: Path) -> dict[str, str]:
    try:
        with ZipFile(path) as zf:
            if "docProps/app.xml" not in zf.namelist():
                return {}
            root = etree.fromstring(zf.read("docProps/app.xml"))
    except Exception:
        return {}
    props: dict[str, str] = {}
    for child in root:
        tag = etree.QName(child).localname
        props[tag] = (child.text or "").strip()
    return props


def _count_ooxml_image_references(path: Path, prefixes: tuple[str, ...]) -> int:
    count = 0
    try:
        with ZipFile(path) as zf:
            for name in zf.namelist():
                if not name.endswith(".xml") or "/_rels/" in name:
                    continue
                if prefixes and not any(name.startswith(prefix) for prefix in prefixes):
                    continue
                try:
                    root = etree.fromstring(zf.read(name))
                except Exception:
                    continue
                for node in root.iter():
                    local_name = etree.QName(node).localname
                    if local_name in {"blip", "imagedata"}:
                        count += 1
    except Exception:
        return 0
    return count


def _count_pdf_images(path: Path) -> int:
    try:
        import fitz

        with fitz.open(str(path)) as doc:
            return sum(len(page.get_images(full=True)) for page in doc)
    except Exception:
        return 0


def _count_text_lines(text: str) -> int:
    return sum(1 for line in re.split(r"\r\n|\r|\n", text or "") if line.strip())


def _count_text_paragraphs(text: str) -> int:
    content = (text or "").strip()
    if not content:
        return 0
    blocks = [block for block in re.split(r"(?:\r?\n\s*){2,}", content) if block.strip()]
    return len(blocks) if blocks else _count_text_lines(content)


def _to_positive_int(value: Any) -> int:
    try:
        number = int(str(value or "").strip())
    except (TypeError, ValueError):
        return 0
    return number if number > 0 else 0


def _now_iso() -> str:
    return datetime.now().isoformat(timespec="seconds")


def _file_type_for_extension(extension: str) -> str:
    ext = (extension or "").lower()
    if ext in {".doc", ".docx"}:
        return "Word"
    if ext in {".xls", ".xlsx"}:
        return "Excel"
    if ext in {".ppt", ".pptx"}:
        return "PPT"
    if ext == ".pdf":
        return "PDF"
    if ext == ".txt":
        return "TXT"
    if ext in IMAGE_EXTENSIONS:
        return "图片"
    if ext in CAD_EXTENSIONS:
        return "CAD"
    return "未知"


def _stat_method_for_extension(extension: str) -> str:
    ext = (extension or "").lower()
    if ext == ".docx":
        return "DOCX XML解析+Word近似计数"
    if ext == ".doc":
        return "DOC经LibreOffice转换+DOCX XML解析+Word近似计数"
    if ext == ".xlsx":
        return "XLSX单元格解析+Word近似计数（页数按工作表数近似）"
    if ext == ".xls":
        return "XLS经LibreOffice转换+XLSX单元格解析+Word近似计数（页数按工作表数近似）"
    if ext == ".pptx":
        return "PPTX形状文本解析+Word近似计数"
    if ext == ".ppt":
        return "PPT经LibreOffice转换+PPTX形状文本解析+Word近似计数"
    if ext == ".pdf":
        return "PDF文本层解析+Word近似计数"
    if ext == ".txt":
        return "TXT文本解析+Word近似计数"
    if ext in IMAGE_EXTENSIONS:
        return "图片候选文件，需OCR后统计"
    if ext in CAD_EXTENSIONS:
        if ext == ".dxf":
            return "DXF实体解析+Word近似计数"
        return "ODA File Converter转DXF+ezdxf实体解析+Word近似计数"
    return ""


def _count_text(text: str) -> TextMetrics:
    content = text or ""
    counts = {field: 0 for field in SCRIPT_COUNT_FIELDS}
    index = 0
    while index < len(content):
        char = content[index]
        cjk_field = _cjk_script_field(char)
        if cjk_field:
            counts[cjk_field] += 1
            index += 1
            continue
        if _is_token_char(char):
            token, index = _consume_word_like_token(content, index)
            counts[_token_count_field(token)] += 1
            continue
        previous_char = content[index - 1] if index > 0 else ""
        next_char = content[index + 1] if index + 1 < len(content) else ""
        if _is_word_count_cjk_punctuation(char, previous_char, next_char):
            counts["cjk_punct_count"] += 1
            index += 1
            continue
        index += 1
    word_count = sum(counts.values())
    return TextMetrics(
        word_count=word_count,
        non_space_chars=sum(1 for char in text or "" if not char.isspace()),
        raw_chars=len(text or ""),
        **counts,
    )


def _quote_counts_from_script_counts(script_counts: dict[str, int]) -> dict[str, int]:
    counts = {field: int(script_counts.get(field) or 0) for field in SCRIPT_COUNT_FIELDS}
    return {
        "billable_chinese_count": counts["han_count"] + counts["cjk_punct_count"],
        "billable_japanese_count": counts["han_count"] + counts["kana_count"] + counts["cjk_punct_count"],
        "billable_korean_count": counts["hangul_count"] + counts["cjk_punct_count"],
        "billable_latin_count": counts["latin_word_count"] + counts["mixed_latin_number_count"],
        "billable_cyrillic_count": counts["cyrillic_word_count"],
        "billable_arabic_count": counts["arabic_word_count"],
        "billable_greek_count": counts["greek_word_count"],
        "billable_hebrew_count": counts["hebrew_word_count"],
        "billable_thai_count": counts["thai_word_count"],
    }


def _consume_word_like_token(text: str, start: int) -> tuple[str, int]:
    index = start
    while index < len(text):
        char = text[index]
        if _is_token_char(char):
            index += 1
            continue
        if char in {"'", "’", "-", "_", ".", "/"}:
            index += 1
            continue
        break
    return text[start:index], index


def _cjk_script_field(char: str) -> str:
    if _is_han(char):
        return "han_count"
    if _is_kana(char):
        return "kana_count"
    if _is_hangul(char):
        return "hangul_count"
    return ""


def _is_cjk(char: str) -> bool:
    return bool(_cjk_script_field(char))


def _is_han(char: str) -> bool:
    if not char or len(char) != 1:
        return False
    code = ord(char)
    return (
        0x3400 <= code <= 0x4DBF
        or 0x4E00 <= code <= 0x9FFF
        or 0xF900 <= code <= 0xFAFF
        or 0x20000 <= code <= 0x2A6DF
        or 0x2A700 <= code <= 0x2B73F
        or 0x2B740 <= code <= 0x2B81F
        or 0x2B820 <= code <= 0x2CEAF
        or 0x2CEB0 <= code <= 0x2EBEF
        or 0x2F800 <= code <= 0x2FA1F
        or 0x30000 <= code <= 0x3134F
    )


def _is_kana(char: str) -> bool:
    if not char or len(char) != 1:
        return False
    code = ord(char)
    return (
        0x3040 <= code <= 0x30FF
        or 0x31F0 <= code <= 0x31FF
        or 0xFF66 <= code <= 0xFF9D
    )


def _is_hangul(char: str) -> bool:
    if not char or len(char) != 1:
        return False
    code = ord(char)
    return (
        0x1100 <= code <= 0x11FF
        or 0x3130 <= code <= 0x318F
        or 0xA960 <= code <= 0xA97F
        or 0xAC00 <= code <= 0xD7AF
        or 0xD7B0 <= code <= 0xD7FF
    )


def _is_token_char(char: str) -> bool:
    if not char or _is_cjk(char):
        return False
    category = unicodedata.category(char)
    return category[0] in {"L", "N"}


def _token_count_field(token: str) -> str:
    letter_scripts: set[str] = set()
    has_number = False
    for char in token:
        if not _is_token_char(char):
            continue
        category = unicodedata.category(char)
        if category[0] == "N":
            has_number = True
            continue
        letter_scripts.add(_letter_script(char))

    if not letter_scripts and has_number:
        return "number_token_count"
    if len(letter_scripts) == 1:
        script = next(iter(letter_scripts))
        if script == "latin":
            return "mixed_latin_number_count" if has_number else "latin_word_count"
        script_field = {
            "cyrillic": "cyrillic_word_count",
            "arabic": "arabic_word_count",
            "greek": "greek_word_count",
            "hebrew": "hebrew_word_count",
            "thai": "thai_word_count",
        }.get(script)
        if script_field:
            return script_field
    return "other_count"


def _letter_script(char: str) -> str:
    code = ord(char)
    name = unicodedata.name(char, "")
    if "LATIN" in name:
        return "latin"
    if "CYRILLIC" in name:
        return "cyrillic"
    if "ARABIC" in name or 0x0600 <= code <= 0x06FF or 0x0750 <= code <= 0x077F or 0x08A0 <= code <= 0x08FF:
        return "arabic"
    if "GREEK" in name:
        return "greek"
    if "HEBREW" in name:
        return "hebrew"
    if 0x0E00 <= code <= 0x0E7F:
        return "thai"
    return "other"


def _is_word_count_cjk_punctuation(char: str, previous_char: str = "", next_char: str = "") -> bool:
    if not char or char.isspace() or _is_token_char(char):
        return False
    code = ord(char)
    category = unicodedata.category(char)
    if category[0] not in {"P", "S"}:
        return False

    # Word 的中文统计会把中文/全角标点算入“中文字符”，但不会把 © 这类普通符号算入字数。
    if (
        0x3000 <= code <= 0x303F
        or 0xFE10 <= code <= 0xFE1F
        or 0xFE30 <= code <= 0xFE4F
        or 0xFF00 <= code <= 0xFFEF
    ):
        return True

    chinese_context_punctuation = {"“", "”", "‘", "’", "—", "–", "…", "·"}
    if char in chinese_context_punctuation and (_is_cjk(previous_char) or _is_cjk(next_char)):
        return True

    return False


def _build_summary(file_results: list[dict[str, Any]], *, truncated: bool, started_at: datetime) -> dict[str, Any]:
    status_counts = Counter(item.get("status") for item in file_results)
    summary = {
        "total_files": len(file_results),
        "counted_files": status_counts.get(STATUS_COUNTED, 0),
        "failed_files": status_counts.get(STATUS_FAILED, 0),
        "skipped_files": status_counts.get(STATUS_SKIPPED, 0),
        "needs_ocr_files": status_counts.get(STATUS_NEEDS_OCR, 0),
        "needs_cad_parser_files": status_counts.get(STATUS_NEEDS_CAD_PARSER, 0),
        "total_main_word_count": sum(int(item.get("main_word_count") or 0) for item in file_results),
        "total_word_count": sum(int(item.get("word_count") or item.get("main_word_count") or 0) for item in file_results),
        "total_extra_word_count": sum(int(item.get("extra_word_count") or 0) for item in file_results),
        "total_char_count_no_spaces": sum(int(item.get("char_count_no_spaces") or 0) for item in file_results),
        "total_char_count_with_spaces": sum(int(item.get("char_count_with_spaces") or 0) for item in file_results),
        "total_non_space_chars": sum(int(item.get("char_count_no_spaces") or item.get("non_space_chars") or 0) for item in file_results),
        "total_raw_chars": sum(int(item.get("char_count_with_spaces") or item.get("raw_chars") or 0) for item in file_results),
        "total_page_count": sum(int(item.get("page_count") or 0) for item in file_results),
        "total_paragraph_count": sum(int(item.get("paragraph_count") or 0) for item in file_results),
        "total_line_count": sum(int(item.get("line_count") or 0) for item in file_results),
        "total_image_count": sum(int(item.get("image_count") or 0) for item in file_results),
        "ocr_files": sum(1 for item in file_results if item.get("ocr_used")),
        "ocr_pages": sum(int(item.get("ocr_page_count") or 0) for item in file_results),
        "ocr_failed_files": sum(
            1 for item in file_results if item.get("ocr_used") and item.get("status") == STATUS_FAILED
        ),
        "truncated": truncated,
        "started_at": started_at.isoformat(timespec="seconds"),
        "finished_at": datetime.now().isoformat(timespec="seconds"),
    }
    main_script_counts = {
        field: sum(int(item.get(field) or 0) for item in file_results)
        for field in SCRIPT_COUNT_FIELDS
    }
    extra_script_counts = {
        field: sum(int(item.get(f"extra_{field}") or 0) for item in file_results)
        for field in SCRIPT_COUNT_FIELDS
    }
    main_quote_counts = _quote_counts_from_script_counts(main_script_counts)
    extra_quote_counts = _quote_counts_from_script_counts(extra_script_counts)
    summary.update(
        script_counts=main_script_counts,
        quote_counts=main_quote_counts,
        extra_script_counts=extra_script_counts,
        extra_quote_counts=extra_quote_counts,
        script_count_total=sum(main_script_counts.values()),
        extra_script_count_total=sum(extra_script_counts.values()),
    )
    summary.update({f"total_{field}": value for field, value in main_script_counts.items()})
    summary.update({f"total_{field}": value for field, value in main_quote_counts.items()})
    summary.update({f"total_extra_{field}": value for field, value in extra_script_counts.items()})
    summary.update({f"total_extra_{field}": value for field, value in extra_quote_counts.items()})
    return summary


def _write_excel_report(path: Path, payload: dict[str, Any]) -> None:
    workbook = Workbook()
    summary_sheet = workbook.active
    summary_sheet.title = "汇总"
    _write_summary_sheet(summary_sheet, payload)
    _write_file_sheet(workbook.create_sheet("文件明细"), payload.get("files") or [])
    _write_source_sheet(workbook.create_sheet("来源明细"), payload.get("source_details") or [])
    _write_fail_sheet(workbook.create_sheet("跳过与失败"), payload.get("files") or [])
    _write_rules_sheet(workbook.create_sheet("规则说明"), payload.get("rules") or [])
    workbook.save(path)


def _style_header(sheet, row: int = 1) -> None:
    fill = PatternFill("solid", fgColor="1F4E78")
    font = Font(color="FFFFFF", bold=True)
    for cell in sheet[row]:
        cell.fill = fill
        cell.font = font
        cell.alignment = Alignment(horizontal="center", vertical="center")


def _autosize(sheet) -> None:
    for column in sheet.columns:
        max_length = 10
        column_letter = get_column_letter(column[0].column)
        for cell in column:
            value = "" if cell.value is None else str(cell.value)
            max_length = max(max_length, min(len(value) + 2, 60))
        sheet.column_dimensions[column_letter].width = max_length


def _write_summary_sheet(sheet, payload: dict[str, Any]) -> None:
    summary = payload.get("summary") or {}
    input_kind = "单文件" if payload.get("input_kind") == "file" else "目录"
    input_source = "直接上传" if payload.get("input_source") == "upload" else "共享路径"
    rows = [
        ("统计类型", input_kind),
        ("输入方式", input_source),
        ("原始文件名", payload.get("original_filename", "")),
        ("统计路径", payload.get("input_path") or payload.get("directory_path", "")),
        ("允许根目录", payload.get("allowed_root", "")),
        ("主字数合计", summary.get("total_main_word_count", 0)),
        ("额外内容字数合计", summary.get("total_extra_word_count", 0)),
        ("中文候选字数", summary.get("total_billable_chinese_count", 0)),
        ("日语候选字数", summary.get("total_billable_japanese_count", 0)),
        ("韩语候选字数", summary.get("total_billable_korean_count", 0)),
        ("拉丁系候选词数", summary.get("total_billable_latin_count", 0)),
        ("纯数字 token", summary.get("total_number_token_count", 0)),
        ("脚本桶合计", summary.get("script_count_total", 0)),
        ("字符数(不计空格)合计", summary.get("total_char_count_no_spaces", summary.get("total_non_space_chars", 0))),
        ("字符数(计空格)合计", summary.get("total_char_count_with_spaces", summary.get("total_raw_chars", 0))),
        ("页数合计", summary.get("total_page_count", 0)),
        ("段落数合计", summary.get("total_paragraph_count", 0)),
        ("行数合计", summary.get("total_line_count", 0)),
        ("图片数量合计", summary.get("total_image_count", 0)),
        ("OCR 文件数", summary.get("ocr_files", 0)),
        ("OCR 页数", summary.get("ocr_pages", 0)),
        ("OCR 失败文件数", summary.get("ocr_failed_files", 0)),
        ("OCR 模式", payload.get("ocr_mode", OCR_MODE_AUTO)),
        ("OCR 模型", payload.get("ocr_model", "")),
        ("文件总数", summary.get("total_files", 0)),
        ("已统计文件数", summary.get("counted_files", 0)),
        ("失败文件数", summary.get("failed_files", 0)),
        ("跳过文件数", summary.get("skipped_files", 0)),
        ("需 OCR 文件数", summary.get("needs_ocr_files", 0)),
        ("需 CAD 解析文件数", summary.get("needs_cad_parser_files", 0)),
        ("是否截断", "是" if summary.get("truncated") else "否"),
        ("生成时间", payload.get("generated_at", "")),
    ]
    rows.extend(
        (f"脚本明细：{SCRIPT_COUNT_LABELS.get(field, field)}", summary.get(f"total_{field}", 0))
        for field in SCRIPT_COUNT_FIELDS
    )
    sheet.append(["项目", "值"])
    for row in rows:
        sheet.append(list(row))
    _style_header(sheet)
    _autosize(sheet)


def _write_file_sheet(sheet, files: list[dict[str, Any]]) -> None:
    headers = [
        "相对路径",
        "页数",
        "字数",
        "额外字数",
        "中文候选",
        "日语候选",
        "韩语候选",
        "拉丁系候选",
        "纯数字",
        "脚本桶合计",
        "字符数(不计空格)",
        "字符数(计空格)",
        "段落数",
        "行数",
        "图片数量",
        "是否使用 OCR",
        "OCR 页数",
        "OCR 模型",
        "OCR 失败页",
        "OCR 文本路径",
        "统计方法",
        "文件类型",
        "状态",
        "错误信息",
        "统计时间",
        "扩展名",
        "文件大小",
        "修改时间",
        "消息",
        "警告",
        "错误",
    ]
    headers.extend(SCRIPT_COUNT_LABELS[field] for field in SCRIPT_COUNT_FIELDS)
    sheet.append(headers)
    for item in files:
        sheet.append(
            [
                item.get("relative_path", ""),
                item.get("page_count", 0),
                item.get("word_count", item.get("main_word_count", 0)),
                item.get("extra_word_count", 0),
                item.get("billable_chinese_count", 0),
                item.get("billable_japanese_count", 0),
                item.get("billable_korean_count", 0),
                item.get("billable_latin_count", 0),
                item.get("number_token_count", 0),
                item.get("script_count_total", 0),
                item.get("char_count_no_spaces", item.get("non_space_chars", 0)),
                item.get("char_count_with_spaces", item.get("raw_chars", 0)),
                item.get("paragraph_count", 0),
                item.get("line_count", 0),
                item.get("image_count", 0),
                "是" if item.get("ocr_used") else "否",
                item.get("ocr_page_count", 0),
                item.get("ocr_model", ""),
                ", ".join(str(page) for page in (item.get("ocr_failed_pages") or [])),
                item.get("ocr_text_path", ""),
                item.get("stat_method", ""),
                item.get("file_type", ""),
                item.get("status", ""),
                item.get("error", ""),
                item.get("counted_at", ""),
                item.get("extension", ""),
                item.get("size_bytes", 0),
                item.get("modified_at", ""),
                item.get("message", ""),
                item.get("warning", ""),
                item.get("error", ""),
            ]
            + [item.get(field, 0) for field in SCRIPT_COUNT_FIELDS]
        )
    _style_header(sheet)
    _autosize(sheet)


def _write_source_sheet(sheet, rows: list[dict[str, Any]]) -> None:
    headers = [
        "相对路径",
        "扩展名",
        "来源类型",
        "来源标签",
        "是否额外内容",
        "字数",
        "中文候选",
        "日语候选",
        "韩语候选",
        "拉丁系候选",
        "纯数字",
        "脚本桶合计",
        "字符数(不计空格)",
        "字符数(计空格)",
        "段落数",
        "行数",
        "图片数量",
        "文本预览",
    ]
    headers.extend(SCRIPT_COUNT_LABELS[field] for field in SCRIPT_COUNT_FIELDS)
    sheet.append(headers)
    for row in rows:
        sheet.append(
            [
                row.get("relative_path", ""),
                row.get("extension", ""),
                row.get("source_type", ""),
                row.get("source_label", ""),
                "是" if row.get("is_extra") else "否",
                row.get("word_count", 0),
                row.get("billable_chinese_count", 0),
                row.get("billable_japanese_count", 0),
                row.get("billable_korean_count", 0),
                row.get("billable_latin_count", 0),
                row.get("number_token_count", 0),
                row.get("script_count_total", 0),
                row.get("char_count_no_spaces", row.get("non_space_chars", 0)),
                row.get("char_count_with_spaces", row.get("raw_chars", 0)),
                row.get("paragraph_count", 0),
                row.get("line_count", 0),
                row.get("image_count", 0),
                row.get("text_preview", ""),
            ]
            + [row.get(field, 0) for field in SCRIPT_COUNT_FIELDS]
        )
    _style_header(sheet)
    _autosize(sheet)


def _write_fail_sheet(sheet, files: list[dict[str, Any]]) -> None:
    sheet.append(["相对路径", "扩展名", "文件类型", "状态", "统计方法", "图片数量", "OCR 页数", "OCR 失败页", "消息", "警告", "错误", "统计时间"])
    for item in files:
        if item.get("status") == STATUS_COUNTED:
            continue
        sheet.append(
            [
                item.get("relative_path", ""),
                item.get("extension", ""),
                item.get("file_type", ""),
                item.get("status", ""),
                item.get("stat_method", ""),
                item.get("image_count", 0),
                item.get("ocr_page_count", 0),
                ", ".join(str(page) for page in (item.get("ocr_failed_pages") or [])),
                item.get("message", ""),
                item.get("warning", ""),
                item.get("error", ""),
                item.get("counted_at", ""),
            ]
        )
    _style_header(sheet)
    _autosize(sheet)


def _write_rules_sheet(sheet, rules: list[str]) -> None:
    sheet.append(["规则说明"])
    for rule in rules:
        sheet.append([rule])
    _style_header(sheet)
    _autosize(sheet)


def _relative_path(path: Path, root: Path) -> str:
    try:
        return path.resolve(strict=False).relative_to(root.resolve(strict=False)).as_posix()
    except ValueError:
        return path.as_posix()


def _preview_text(text: str, limit: int = 120) -> str:
    normalized = re.sub(r"\s+", " ", text or "").strip()
    return normalized[:limit]


def _write_ocr_plain_text(*, ocr_text_dir: Path, relative_path: str, items: list[TextItem]) -> Path:
    relative = Path(relative_path)
    safe_parts = [part for part in relative.parts if part not in {"", ".", "..", relative.anchor}]
    safe_relative = Path(*safe_parts) if safe_parts else Path("ocr_result")
    target = ocr_text_dir / safe_relative.parent / f"{safe_relative.name}.txt"
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(
        "\n\n".join(item.text.strip() for item in items if item.text.strip()),
        encoding="utf-8",
    )
    return target


def _write_ocr_text_archive(output_dir: Path, ocr_text_dir: Path, files: list[Path]) -> Optional[Path]:
    existing_files = [path for path in files if path.exists() and path.is_file()]
    if not existing_files:
        return None
    archive_path = output_dir / "OCR识别文本.zip"
    with ZipFile(archive_path, "w", compression=ZIP_DEFLATED) as archive:
        for file_path in existing_files:
            try:
                archive_name = file_path.relative_to(ocr_text_dir).as_posix()
            except ValueError:
                archive_name = file_path.name
            archive.write(file_path, arcname=archive_name)
    return archive_path


def _output_web_path(path: Path) -> str:
    output_root = Path(settings.OUTPUT_DIR).resolve()
    resolved = path.resolve()
    try:
        return f"outputs/{resolved.relative_to(output_root).as_posix()}"
    except ValueError:
        return str(path).replace("\\", "/")
