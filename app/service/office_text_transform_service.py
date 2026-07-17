# -*- coding: utf-8 -*-
"""复用项目现有 Office 遍历器，对文档中的文本做保格式批量转换。"""

from __future__ import annotations

import asyncio
import zipfile
from collections import Counter
from concurrent.futures import Executor
from pathlib import Path
from typing import Any, Awaitable, Callable

from docx import Document
from docx.text.run import Run
from lxml import etree
from openpyxl.cell.cell import MergedCell
from pptx import Presentation

from app.core.config import settings
from app.core.file_naming import ensure_unique_path
from app.service.english_variant_service import EnglishVariantConverter, get_converter, normalize_target_style
from app.service.libreoffice_service import (
    convert_doc_to_docx_via_libreoffice,
    convert_presentation_to_pptx_via_libreoffice,
    convert_spreadsheet_to_xlsx_via_libreoffice,
)


ProgressCallback = Callable[[int, str], Awaitable[None]]
W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
W_TEXT = f"{{{W_NS}}}t"
LEGACY_OUTPUT_SUFFIX = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}


class ConversionSummary:
    def __init__(self, converter: EnglishVariantConverter, target_style: str) -> None:
        self.converter = converter
        self.target_style = normalize_target_style(target_style)
        self.replacements: Counter[tuple[str, str]] = Counter()
        self.ambiguous: Counter[tuple[str, tuple[str, ...]]] = Counter()
        self.processed_text_units = 0

    def add(self, result: dict[str, Any]) -> None:
        self.processed_text_units += 1
        for item in result.get("replacements") or []:
            self.replacements[(str(item["source"]), str(item["target"]))] += int(item["count"])
        for item in result.get("ambiguous_hits") or []:
            key = (str(item["term"]), tuple(str(value) for value in item.get("candidates") or []))
            self.ambiguous[key] += int(item["count"])

    def to_dict(self) -> dict[str, Any]:
        replacements = [
            {"source": source, "target": target, "count": count}
            for (source, target), count in sorted(
                self.replacements.items(), key=lambda item: (-item[1], item[0][0].casefold())
            )
        ]
        ambiguous_hits = [
            {"term": term, "candidates": list(candidates), "count": count}
            for (term, candidates), count in sorted(
                self.ambiguous.items(), key=lambda item: (-item[1], item[0][0].casefold())
            )
        ]
        return {
            "target_style": self.target_style,
            "replacement_count": sum(self.replacements.values()),
            "distinct_rule_count": len(self.replacements),
            "replacements": replacements,
            "ambiguous_hit_count": sum(self.ambiguous.values()),
            "ambiguous_hits": ambiguous_hits,
            "processed_text_units": self.processed_text_units,
            "dictionary_version": self.converter.dictionary_version,
            "dictionary_sha256": self.converter.source_sha256,
        }


def _load_specialist_office_helpers():
    """通过现有中翻译服务的命名空间装载器复用 Office 遍历/回写类。"""
    from app.service import zhongfanyi_service

    with zhongfanyi_service._specialist_import_lock:
        zhongfanyi_service._prepare_zhongfanyi_import_path()
        from replace.excel.excel_replacer import ExcelReplacer
        from replace.pptx.pptx_replacer import PPTXReplacer
        from replace.word.replace_clean import (
            iter_body_paragraphs,
            iter_footer_paragraphs,
            iter_header_paragraphs,
        )

    return ExcelReplacer, PPTXReplacer, iter_body_paragraphs, iter_header_paragraphs, iter_footer_paragraphs


def _find_segment_index(lengths: list[int], position: int) -> tuple[int, int]:
    consumed = 0
    for index, length in enumerate(lengths):
        if position < consumed + length:
            return index, position - consumed
        consumed += length
    if not lengths:
        return 0, 0
    return len(lengths) - 1, lengths[-1]


def _apply_edits_to_segments(segments: list[Any], edits: list[dict[str, Any]]) -> None:
    original_texts = [str(getattr(segment, "text", "") or "") for segment in segments]
    lengths = [len(text) for text in original_texts]
    for edit in reversed(edits):
        start = int(edit["start"])
        end = int(edit["end"])
        if start >= end or not segments:
            continue
        start_index, start_offset = _find_segment_index(lengths, start)
        end_index, end_offset_last = _find_segment_index(lengths, end - 1)
        end_offset = end_offset_last + 1
        replacement = str(edit["after"])

        start_text = str(getattr(segments[start_index], "text", "") or "")
        if start_index == end_index:
            segments[start_index].text = start_text[:start_offset] + replacement + start_text[end_offset:]
            continue

        end_text = str(getattr(segments[end_index], "text", "") or "")
        segments[start_index].text = start_text[:start_offset] + replacement
        for index in range(start_index + 1, end_index):
            segments[index].text = ""
        segments[end_index].text = end_text[end_offset:]


def _paragraph_runs(paragraph) -> list[Run]:
    try:
        elements = paragraph._p.xpath("./w:r | ./w:hyperlink/w:r")
        return [Run(element, paragraph) for element in elements]
    except Exception:
        return list(paragraph.runs)


def _transform_segments(
    segments: list[Any],
    converter: EnglishVariantConverter,
    summary: ConversionSummary,
) -> None:
    text = "".join(str(getattr(segment, "text", "") or "") for segment in segments)
    if not text:
        return
    result = converter.convert(text, summary.target_style, include_edits=True)
    summary.add(result)
    edits = result.get("_edits") or []
    if edits:
        _apply_edits_to_segments(segments, edits)


def _transform_docx(input_path: Path, output_path: Path, summary: ConversionSummary) -> None:
    _, _, iter_body, iter_headers, iter_footers = _load_specialist_office_helpers()
    document = Document(str(input_path))
    # 直接保存 XML 元素对象，避免临时包装对象被回收后 Python 重用 id，
    # 导致页脚等后续容器被误判为已经处理。
    seen: set[Any] = set()
    for iterator in (iter_body(document), iter_headers(document), iter_footers(document)):
        for paragraph in iterator:
            ancestors = {ancestor.tag for ancestor in paragraph._element.iterancestors()}
            if f"{{{W_NS}}}footnote" in ancestors or f"{{{W_NS}}}endnote" in ancestors:
                continue
            element = paragraph._element
            if element in seen:
                continue
            seen.add(element)
            _transform_segments(_paragraph_runs(paragraph), summary.converter, summary)
    document.save(str(output_path))
    _transform_docx_note_parts(output_path, summary)


def _transform_docx_note_parts(output_path: Path, summary: ConversionSummary) -> None:
    part_names = {"word/footnotes.xml", "word/endnotes.xml"}
    temp_path = output_path.with_suffix(".notes.tmp.docx")
    changed = False
    with zipfile.ZipFile(output_path, "r") as source_zip, zipfile.ZipFile(temp_path, "w") as target_zip:
        for item in source_zip.infolist():
            data = source_zip.read(item.filename)
            if item.filename in part_names:
                root = etree.fromstring(data)
                for paragraph in root.findall(f".//{{{W_NS}}}p"):
                    nodes = paragraph.findall(f".//{W_TEXT}")
                    if nodes:
                        before = "".join(node.text or "" for node in nodes)
                        _transform_segments(nodes, summary.converter, summary)
                        if "".join(node.text or "" for node in nodes) != before:
                            changed = True
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            target_zip.writestr(item, data)
    if changed:
        temp_path.replace(output_path)
    else:
        temp_path.unlink(missing_ok=True)


def _transform_xlsx(input_path: Path, output_path: Path, summary: ConversionSummary) -> None:
    ExcelReplacer, _, _, _, _ = _load_specialist_office_helpers()
    replacer = ExcelReplacer(str(input_path))
    for worksheet in replacer.wb.worksheets:
        for row in worksheet.iter_rows():
            for cell in row:
                if isinstance(cell, MergedCell) or cell.data_type == "f" or not isinstance(cell.value, str):
                    continue
                result = summary.converter.convert(cell.value, summary.target_style)
                summary.add(result)
                if result["converted_text"] != cell.value:
                    cell.value = result["converted_text"]
    replacer.save(str(output_path))


def _iter_pptx_text_frames(replacer, presentation: Presentation):
    for slide_index, slide in enumerate(presentation.slides):
        yield from replacer._iter_text_frames(slide.shapes, slide_index)
        yield from replacer._iter_table_text_frames(slide.shapes, slide_index)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            yield {"text_frame": slide.notes_slide.notes_text_frame, "slide_idx": slide_index}


def _transform_pptx(input_path: Path, output_path: Path, summary: ConversionSummary) -> None:
    _, PPTXReplacer, _, _, _ = _load_specialist_office_helpers()
    presentation = Presentation(str(input_path))
    replacer = PPTXReplacer.__new__(PPTXReplacer)
    replacer.prs = presentation
    seen: set[Any] = set()
    for item in _iter_pptx_text_frames(replacer, presentation):
        text_frame = item["text_frame"]
        tx_body = text_frame._txBody
        if tx_body in seen:
            continue
        seen.add(tx_body)
        for paragraph in text_frame.paragraphs:
            _transform_segments(list(paragraph.runs), summary.converter, summary)
    presentation.save(str(output_path))


def _normalize_office_input(input_path: Path, work_dir: Path) -> Path:
    suffix = input_path.suffix.lower()
    if suffix == ".doc":
        return Path(convert_doc_to_docx_via_libreoffice(input_path, work_dir / f"{input_path.stem}.docx"))
    if suffix == ".xls":
        return Path(convert_spreadsheet_to_xlsx_via_libreoffice(input_path, work_dir / f"{input_path.stem}.xlsx"))
    if suffix == ".ppt":
        return Path(convert_presentation_to_pptx_via_libreoffice(input_path, work_dir / f"{input_path.stem}.pptx"))
    return input_path


def transform_office_file(
    input_path: str | Path,
    output_path: str | Path,
    target_style: str,
) -> dict[str, Any]:
    source = Path(input_path)
    output = Path(output_path)
    converter = get_converter()
    summary = ConversionSummary(converter, target_style)
    suffix = source.suffix.lower()
    if suffix == ".docx":
        _transform_docx(source, output, summary)
    elif suffix == ".xlsx":
        _transform_xlsx(source, output, summary)
    elif suffix == ".pptx":
        _transform_pptx(source, output, summary)
    else:
        raise ValueError(f"不支持的 Office 格式: {source.name}")
    return summary.to_dict()


def _web_output_path(file_path: Path) -> str:
    relative = file_path.resolve().relative_to(Path(settings.OUTPUT_DIR).resolve())
    return f"outputs/{relative.as_posix()}"


async def execute_english_variant_task(
    *,
    task_id: str,
    display_no: str | None,
    input_path: str,
    original_filename: str,
    target_style: str,
    progress_callback: ProgressCallback | None = None,
    executor: Executor | None = None,
) -> dict[str, Any]:
    style = normalize_target_style(target_style)
    source = Path(input_path)
    original_suffix = Path(original_filename).suffix.lower() or source.suffix.lower()
    output_suffix = LEGACY_OUTPUT_SUFFIX.get(original_suffix, original_suffix)
    if output_suffix not in {".docx", ".xlsx", ".pptx"}:
        raise ValueError(f"不支持的 Office 格式: {original_filename}")

    output_dir = Path(settings.OUTPUT_DIR) / "english_variant" / (display_no or task_id)
    output_dir.mkdir(parents=True, exist_ok=True)
    work_dir = output_dir / "converted_inputs"
    work_dir.mkdir(parents=True, exist_ok=True)
    loop = asyncio.get_running_loop()

    if progress_callback:
        await progress_callback(10, "正在准备 Office 文档...")
    prepared = await loop.run_in_executor(executor, lambda: _normalize_office_input(source, work_dir))

    stem = Path(original_filename).stem or source.stem
    desired = output_dir / f"{stem}_{style}{output_suffix}"
    output_path = ensure_unique_path(desired)
    if progress_callback:
        await progress_callback(30, "正在转换英美式英语词汇...")
    summary = await loop.run_in_executor(
        executor,
        lambda: transform_office_file(prepared, output_path, style),
    )
    if progress_callback:
        await progress_callback(95, "正在整理转换结果...")

    return {
        "task_id": task_id,
        "input_format": original_suffix,
        "output_format": output_suffix,
        "output_file": _web_output_path(output_path),
        "output_filename": output_path.name,
        **summary,
    }


__all__ = [
    "ConversionSummary",
    "execute_english_variant_task",
    "transform_office_file",
]
